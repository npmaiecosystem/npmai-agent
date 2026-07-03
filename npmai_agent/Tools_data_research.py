"""
tools_data_research.py
NPM Agent — NPMAI ECOSYSTEM (by Sonu Kumar)
Data & Research vertical: DataAnalysis, Visualization, Web Scraping,
Search/Research, Financial, Social Media, Weather/Geo, Text Analytics,
Database, and Report Generation tools.
"""

import os
import sys
import json
import re
import shutil
import subprocess
import tempfile
import time
import threading
import hashlib
import base64
import platform
import glob
import zipfile
import sqlite3
from pathlib import Path
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any, Callable, Union, Tuple

# ── auto-installer ────────────────────────────────────────────────────────────

def _ensure(pkg: str, import_name: str = None):
    name = import_name or pkg
    try:
        __import__(name)
    except:
        try:
          subprocess.run([sys.executable, "-m", "pip", "install", pkg, "-q"], check=False)
        except:
          print(f"Some packages is not installed properly in your environment due to some reasons these are the packages {n}")

for _pkg, _imp in [
    ("pandas",                  "pandas"),
    ("numpy",                   "numpy"),
    ("scikit-learn",            "sklearn"),
    ("scipy",                   "scipy"),
    ("plotly",                  "plotly"),
    ("matplotlib",              "matplotlib"),
    ("seaborn",                 "seaborn"),
    ("folium",                  "folium"),
    ("requests",                "requests"),
    ("beautifulsoup4",          "bs4"),
    ("httpx",                   "httpx"),
    ("lxml",                    "lxml"),
    ("arxiv",                   "arxiv"),
    ("biopython",               "Bio"),
    ("newspaper3k",             "newspaper"),
    ("yfinance",                "yfinance"),
    ("tweepy",                  "tweepy"),
    ("praw",                    "praw"),
    ("google-api-python-client","googleapiclient"),
    ("textblob",                "textblob"),
    ("langdetect",              "langdetect"),
    ("psycopg2-binary",         "psycopg2"),
    ("pymongo",                 "pymongo"),
    ("redis",                   "redis"),
    ("reportlab",               "reportlab"),
    ("python-docx",             "docx"),
    ("openpyxl",                "openpyxl"),
    ("python-pptx",             "pptx"),
    ("googlemaps",              "googlemaps"),
    ("ccxt",                    "ccxt"),
    ("sentence-transformers",   "sentence_transformers"),
    ("polars",                  "polars"),
]:
    _ensure(_pkg, _imp)

# ── agent_core import ─────────────────────────────────────────────────────────
try:
    from core import ToolResult, CredStore
except ImportError:
    class ToolResult:
        def __init__(self, success: bool, output: str, data=None):
            self.success = success
            self.output  = output
            self.data    = data
        def __str__(self): return self.output

    class CredStore:
        @classmethod
        def load(cls, name: str) -> dict:
            p = Path.home() / ".npmai_agent" / "creds.json"
            if not p.exists(): return {}
            try:
                return json.loads(p.read_text()).get(name, {})
            except: return {}


# ─────────────────────────────────────────────────────────────────────────────
# 1. DataAnalysisTool
# ─────────────────────────────────────────────────────────────────────────────

class DataAnalysisTool:
    name = "data_analysis"
    description = (
        "Advanced pandas/polars data operations: load, clean, profile, transform, "
        "merge, pivot, time-series, correlation, clustering, ML feature importance, "
        "and natural-language queries via Ollama."
    )
    use = (
        """Name of Tool:- DataAnalysisTool

Purpose of Tool:-
The DataAnalysisTool is an advanced data intelligence and feature engineering wrapper built on top of high-performance ecosystems like pandas, scikit-learn, and Plotly. It serves as an automated pipeline engine designed to load multi-format tabular datasets (such as CSV, Parquet, Excel, and JSON), clean anomalies via robust median/mode imputation, and profile data structures dynamically. Beyond structural formatting, the tool provides native statistical and analytical modules—including K-Means clustering configurations, unsupervised outlier extraction algorithms, feature significance modeling via Random Forests, rolling time-series analyses, and interactive correlation engines. Additionally, it integrates a localized natural language parsing gateway that passes dataset schemas to an LLM context layer to execute zero-shot data extraction queries programmatically.

Methods:-
- load: Parses raw external data objects into standardized in-memory structures while extracting baseline dimensional logs.
- save: Dispatches operational dataset mutations back to disk storage configurations using specified tabular serialization protocols.
- profile: Automates exploratory data reporting pipelines, compiling structural indices, duplicate evaluations, and descriptive statistics into portable HTML files.
- clean: Programmatically strips record duplicates, automatically resolves messy column data types, and replaces null nodes with statistical markers.
- transform: Runs sequenced structural modifications (such as string filters, row sorting, text modifications, and computed evaluations) across active tables.
- filter_data: Isolates target cohorts from source arrays using specialized tabular filtering queries.
- merge_files: Combines split database fragments into centralized matrices over specified alignment keys.
- pivot: Reshapes transactional tables into strategic multi-dimensional summaries aggregated by key metrics.
- time_series_analysis: Extracts sequential trends by running rolling metrics over temporal features and rendering the output into interactive line components.
- correlation_matrix: Computes spatial dependency strengths across linear data properties using specified similarity tests.
- outlier_detection: Uncovers statistical data anomalies using IQR bounds or standard variations.
- feature_importance: Measures target variable correlation weights using machine learning tree splits to isolate top predictive properties.
- cluster_data: Separates highly dimensioned coordinate records into discrete relational cohorts using unsupervised spatial modeling.
- natural_language_query: Translates human language requests into executed syntax expressions to query structural data.
- auto_visualize: Scans target structures to generate distributed charts and frequency graphs automatically.

How to use Tool Methods:-

1. load:
   - Purpose: Translates local files into an operative data context.
   - Arguments:
     a) path: str - Target dataset storage path.
     b) sheet: str (default: "Sheet1") - Specific worksheet index identifier for spreadsheet files.
     c) sep: str (default: ",") - Delimiter character tracking separated value entries.
     d) encoding: str (default: "utf-8") - Standard file bit format descriptor.
     e) dtype: Optional[Dict[str, str]] (default: None) - Key-value map enforcing hard column data types.
   - Returns: ToolResult tracking parsed dimension attributes.
   - How to call: DataAnalysisTool.load(path="data/metrics.csv", sep=",", dtype={"user_id": "str"})

2. save:
   - Purpose: Commits active dataset changes permanently back to file systems.
   - Arguments:
     a) df: DataFrame - The internal source data matrix to output.
     b) path: str - Destination tracking path string.
     c) format: str (default: "csv") - Target layout storage type ("csv", "xlsx", "parquet", "json").
     d) index: bool (default: False) - Flag determining whether to include index rows.
   - Returns: ToolResult highlighting record serialization summaries.
   - How to call: DataAnalysisTool.save(df=clean_df, path="exports/final_report.parquet", format="parquet")

3. profile:
   - Purpose: Constructs automated auditing dossiers of specific dataset traits.
   - Arguments:
     a) df_or_path: DataFrame or str - Native data frame object or file system address.
     b) output_html: str (default: "profile_report.html") - File target endpoint identifier.
   - Returns: ToolResult embedding structured parameter matrix tallies.
   - How to call: DataAnalysisTool.profile(df_or_path="raw_data.json", output_html="reports/audit.html")

4. clean:
   - Purpose: Normalizes structural tabular sets by filling null fields and cleaning row items.
   - Arguments:
     a) df_or_path: DataFrame or str - Source object or path reference tracker.
     b) drop_duplicates: bool (default: True) - Flag stating whether matching rows should be removed.
     c) fill_nulls: Optional[Dict[str, Any]] (default: None) - Dedicated custom value dictionary mapping for target column null blocks.
     d) fix_dtypes: bool (default: True) - Activates automated object-to-datetime or numeric processing checks.
   - Returns: ToolResult logging mutation actions and remaining null values.
   - How to call: DataAnalysisTool.clean(df_or_path=raw_df, fill_nulls={"status": "unknown"})

5. transform:
   - Purpose: Executes sequential column, structural, or value transformations via a serialized operations array.
   - Arguments:
     a) df_or_path: DataFrame or str - Processing target tracking baseline.
     b) operations: List[Dict[str, Any]] - An array of command dictionaries containing actions like `rename`, `drop`, `add_column`, `sort`, `sample`, `filter`, or `astype`.
   - Returns: ToolResult containing the updated data array.
   - How to call: DataAnalysisTool.transform(df_or_path="users.csv", operations=[{"type": "drop", "columns": ["ssn"]}, {"type": "rename", "columns": {"id": "user_id"}}])

6. filter_data:
   - Purpose: Isolates a specific cohort of rows using custom validation rules.
   - Arguments:
     a) df_or_path: DataFrame or str - Operational dataset variable or address.
     b) conditions: List[str] - Query string expressions describing target evaluation checks.
   - Returns: ToolResult holding the subset records.
   - How to call: DataAnalysisTool.filter_data(df_or_path=my_df, conditions=["age >= 18", "country == 'IN'"])

7. merge_files:
   - Purpose: Joins separate files along identical tracking columns.
   - Arguments:
     a) paths: List[str] - Collection of file resource storage locations.
     b) on: Union[str, List[str]] - Shared key property or list of properties to join on.
     c) how: str (default: "inner") - Structural alignment rule mapping ("inner", "outer", "left", "right").
     d) output: str (default: "") - Optional destination tracking path.
   - Returns: ToolResult summarizing unified output parameters.
   - How to call: DataAnalysisTool.merge_files(paths=["orders.csv", "customers.csv"], on="customer_id", how="left")

8. pivot:
   - Purpose: Condenses large rows of logs into structured cross-tabulated dashboards.
   - Arguments:
     a) df_or_path: DataFrame or str - Target transactional record array.
     b) index: str - Property used as the row grouping keys.
     c) columns: str - Property used as the header column categories.
     b) values: str - Target numeric column to calculate across cells.
     e) aggfunc: str (default: "mean") - Evaluation calculator rule type ("sum", "mean", "count", "max", "min").
     f) output: str (default: "") - Optional export file address string.
   - Returns: ToolResult holding the summarized structure tracking dimensions.
   - How to call: DataAnalysisTool.pivot(df_or_path=sales_df, index="region", columns="product", values="revenue", aggfunc="sum")

9. time_series_analysis:
   - Purpose: Plots sequential variations and trend averages over temporal logs.
   - Arguments:
     a) df_or_path: DataFrame or str - Target time-stamped tracker logs.
     b) date_col: str - Column representing dates or times.
     c) value_col: str - Continuous measurement line value.
     d) output: str (default: "time_series.html") - Destination file tracking path.
   - Returns: ToolResult listing critical minimum, maximum, and deviation details.
   - How to call: DataAnalysisTool.time_series_analysis(df_or_path=stock_df, date_col="timestamp", value_col="price")

10. correlation_matrix:
    - Purpose: Visualizes linear association maps across numerical parameters.
    - Arguments:
      a) df_or_path: DataFrame or str - Source table matrix.
      b) method: str (default: "pearson") - Correlation algorithm metric ("pearson", "spearman", "kendall").
      b) output: str (default: "correlation.html") - Visualization export path tracking address.
    - Returns: ToolResult exposing computed correlation metrics.
    - How to call: DataAnalysisTool.correlation_matrix(df_or_path=survey_df, method="spearman")

11. outlier_detection:
    - Purpose: Highlights data anomalies using mathematical distribution checks.
    - Arguments:
      a) df_or_path: DataFrame or str - Targeted dataset structure.
      b) method: str (default: "iqr") - Variance identifier algorithm selection ("iqr", "zscore").
      c) threshold: float (default: 1.5) - Factor setting the outlier boundary sensitivity.
    - Returns: ToolResult mapping absolute anomaly counts found across columns.
    - How to call: DataAnalysisTool.outlier_detection(df_or_path=sensor_df, method="zscore", threshold=3.0)

12. feature_importance:
    - Purpose: Pinpoints which variable columns exert the strongest predictive force on a target outcome.
    - Arguments:
      a) df_or_path: DataFrame or str - Source dataset reference.
      b) target_col: str - The specific dependent variable column to classify or predict.
      c) output: str (default: "feature_importance.html") - Plotly export layout file target destination.
    - Returns: ToolResult listing sorted importance ratios per variable column.
    - How to call: DataAnalysisTool.feature_importance(df_or_path=churn_df, target_col="has_churned")

13. cluster_data:
    - Purpose: Segregates unstructured records into distinct multidimensional spatial groupings.
    - Arguments:
      a) df_or_path: DataFrame or str - Input dataset structure.
      b) n_clusters: int (default: 3) - Fixed count of target group clusters to model.
      c) output: str (default: "clusters.html") - PCA visual plot file target destination.
    - Returns: ToolResult containing population assignment values per cluster category.
    - How to call: DataAnalysisTool.cluster_data(df_or_path=user_features, n_clusters=5, output="charts/segments.html")

14. natural_language_query:
    - Purpose: Evaluates plain-language search questions into dynamic pandas command statements.
    - Arguments:
      a) df_or_path: DataFrame or str - Active workspace table structure.
      b) query: str - Plain-text data exploration request.
      c) model: str (default: "llama3.2:3b") - LLM engine selector tag.
    - Returns: ToolResult tracking parsed data metrics alongside the executed python expression.
    - How to call: DataAnalysisTool.natural_language_query(df_or_path="store.parquet", query="Find top 5 customers with highest spend")

15. auto_visualize:
    - Purpose: Automatically scans structural columns to generate distribution charts and frequency bar plots.
    - Arguments:
      a) df_or_path: DataFrame or str - Tabular data tracking repository.
      b) output_folder: str (default: "auto_charts") - Storage directory path tracking outcoming HTML assets.
    - Returns: ToolResult returning an array of generated image chart location strings.
    - How to call: DataAnalysisTool.auto_visualize(df_or_path=leads_df, output_folder="dashboards/lead_analysis")
""")
    # ── helpers ──────────────────────────────────────────────────────────────

    @staticmethod
    def _to_df(df_or_path):
        import pandas as pd
        if isinstance(df_or_path, pd.DataFrame):
            return df_or_path
        p = str(df_or_path)
        ext = Path(p).suffix.lower()
        if ext in (".xlsx", ".xls"):
            return pd.read_excel(p)
        if ext == ".parquet":
            return pd.read_parquet(p)
        if ext == ".json":
            return pd.read_json(p)
        return pd.read_csv(p)

    # ── methods ──────────────────────────────────────────────────────────────

    @staticmethod
    def load(
        path: str,
        sheet: str = "Sheet1",
        sep: str = ",",
        encoding: str = "utf-8",
        dtype: Optional[Dict[str, str]] = None,
    ) -> ToolResult:
        try:
            import pandas as pd
            ext = Path(path).suffix.lower()
            if ext in (".xlsx", ".xls"):
                df = pd.read_excel(path, sheet_name=sheet, dtype=dtype)
            elif ext == ".parquet":
                df = pd.read_parquet(path)
            elif ext == ".json":
                df = pd.read_json(path, encoding=encoding)
            elif ext in (".tsv",):
                df = pd.read_csv(path, sep="\t", encoding=encoding, dtype=dtype)
            else:
                df = pd.read_csv(path, sep=sep, encoding=encoding, dtype=dtype)
            info = {"rows": len(df), "cols": len(df.columns), "columns": list(df.columns), "dtypes": {c: str(t) for c, t in df.dtypes.items()}}
            return ToolResult(True, f"✓ Loaded {info['rows']} rows × {info['cols']} cols from '{path}'", {"df": df, "info": info})
        except Exception as e:
            return ToolResult(False, f"✗ load failed: {e}")

    @staticmethod
    def save(
        df,
        path: str,
        format: str = "csv",
        index: bool = False,
    ) -> ToolResult:
        try:
            import pandas as pd
            Path(path).parent.mkdir(parents=True, exist_ok=True)
            fmt = format.lower()
            if fmt == "csv":
                df.to_csv(path, index=index)
            elif fmt in ("xlsx", "excel"):
                df.to_excel(path, index=index)
            elif fmt == "parquet":
                df.to_parquet(path, index=index)
            elif fmt == "json":
                df.to_json(path, orient="records", indent=2)
            else:
                df.to_csv(path, index=index)
            return ToolResult(True, f"✓ Saved {len(df)} rows to '{path}' ({fmt})")
        except Exception as e:
            return ToolResult(False, f"✗ save failed: {e}")

    @staticmethod
    def profile(df_or_path, output_html: str = "profile_report.html") -> ToolResult:
        try:
            import pandas as pd
            df = DataAnalysisTool._to_df(df_or_path)
            # Try ydata-profiling first, fall back to manual profile
            try:
                from ydata_profiling import ProfileReport
                report = ProfileReport(df, title="Data Profile", minimal=True)
                report.to_file(output_html)
                return ToolResult(True, f"✓ Profile report saved to '{output_html}'", {"rows": len(df), "cols": len(df.columns)})
            except ImportError:
                pass
            # Manual profile
            profile_data = {
                "shape": {"rows": len(df), "cols": len(df.columns)},
                "dtypes": {c: str(t) for c, t in df.dtypes.items()},
                "missing": df.isnull().sum().to_dict(),
                "missing_pct": (df.isnull().sum() / len(df) * 100).round(2).to_dict(),
                "unique": {c: int(df[c].nunique()) for c in df.columns},
                "duplicates": int(df.duplicated().sum()),
            }
            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            if numeric_cols:
                profile_data["statistics"] = df[numeric_cols].describe().to_dict()
            # write simple HTML
            rows_html = "".join(
                f"<tr><td>{k}</td><td>{json.dumps(v, default=str)[:200]}</td></tr>"
                for k, v in profile_data.items()
            )
            html = f"<html><body><h1>Data Profile</h1><table border='1'>{rows_html}</table></body></html>"
            Path(output_html).write_text(html)
            return ToolResult(True, f"✓ Profile saved to '{output_html}'", profile_data)
        except Exception as e:
            return ToolResult(False, f"✗ profile failed: {e}")

    @staticmethod
    def clean(
        df_or_path,
        drop_duplicates: bool = True,
        fill_nulls: Optional[Dict[str, Any]] = None,
        fix_dtypes: bool = True,
    ) -> ToolResult:
        try:
            import pandas as pd
            df = DataAnalysisTool._to_df(df_or_path).copy()
            before = len(df)
            if drop_duplicates:
                df = df.drop_duplicates()
            if fill_nulls:
                for col, val in fill_nulls.items():
                    if col in df.columns:
                        df[col] = df[col].fillna(val)
            else:
                # default: numeric → median, string → mode
                for col in df.columns:
                    if df[col].isnull().any():
                        if pd.api.types.is_numeric_dtype(df[col]):
                            df[col] = df[col].fillna(df[col].median())
                        else:
                            mode = df[col].mode()
                            if len(mode):
                                df[col] = df[col].fillna(mode[0])
            if fix_dtypes:
                for col in df.select_dtypes(include="object").columns:
                    # try numeric
                    try:
                        df[col] = pd.to_numeric(df[col])
                        continue
                    except (ValueError, TypeError):
                        pass
                    # try datetime
                    try:
                        converted = pd.to_datetime(df[col], infer_datetime_format=True, errors="raise")
                        df[col] = converted
                    except Exception:
                        pass
            after = len(df)
            return ToolResult(True, f"✓ Cleaned: {before - after} rows removed, {df.isnull().sum().sum()} nulls remaining", {"df": df, "rows_before": before, "rows_after": after})
        except Exception as e:
            return ToolResult(False, f"✗ clean failed: {e}")

    @staticmethod
    def transform(df_or_path, operations: List[Dict[str, Any]]) -> ToolResult:
        """
        operations: list of dicts like:
          {"type": "rename",     "columns": {"old": "new"}}
          {"type": "drop",       "columns": ["col1"]}
          {"type": "add_column", "name": "col", "expr": "col1 + col2"}
          {"type": "sort",       "by": ["col"], "ascending": True}
          {"type": "sample",     "n": 100}
        """
        try:
            import pandas as pd
            df = DataAnalysisTool._to_df(df_or_path).copy()
            for op in operations:
                t = op.get("type", "")
                if t == "rename":
                    df = df.rename(columns=op.get("columns", {}))
                elif t == "drop":
                    df = df.drop(columns=[c for c in op.get("columns", []) if c in df.columns])
                elif t == "add_column":
                    df[op["name"]] = df.eval(op["expr"])
                elif t == "sort":
                    df = df.sort_values(by=op.get("by", []), ascending=op.get("ascending", True))
                elif t == "sample":
                    df = df.sample(n=min(op.get("n", 100), len(df)), random_state=42)
                elif t == "fillna":
                    df[op["column"]] = df[op["column"]].fillna(op.get("value"))
                elif t == "astype":
                    df[op["column"]] = df[op["column"]].astype(op.get("dtype", "str"))
                elif t == "filter":
                    df = df.query(op.get("expr", ""))
            return ToolResult(True, f"✓ Transformed: {len(df)} rows × {len(df.columns)} cols", {"df": df})
        except Exception as e:
            return ToolResult(False, f"✗ transform failed: {e}")

    @staticmethod
    def filter_data(df_or_path, conditions: List[str]) -> ToolResult:
        """conditions: list of pandas query strings like ['age > 30', 'country == "US"']"""
        try:
            df = DataAnalysisTool._to_df(df_or_path).copy()
            for cond in conditions:
                df = df.query(cond)
            return ToolResult(True, f"✓ Filtered to {len(df)} rows", {"df": df})
        except Exception as e:
            return ToolResult(False, f"✗ filter_data failed: {e}")

    @staticmethod
    def merge_files(
        paths: List[str],
        on: Union[str, List[str]],
        how: str = "inner",
        output: str = "",
    ) -> ToolResult:
        try:
            import pandas as pd
            dfs = [DataAnalysisTool._to_df(p) for p in paths]
            merged = dfs[0]
            for df in dfs[1:]:
                merged = merged.merge(df, on=on, how=how, suffixes=("", "_dup"))
            if output:
                DataAnalysisTool.save(merged, output)
            return ToolResult(True, f"✓ Merged {len(paths)} files → {len(merged)} rows", {"df": merged})
        except Exception as e:
            return ToolResult(False, f"✗ merge_files failed: {e}")

    @staticmethod
    def pivot(
        df_or_path,
        index: str,
        columns: str,
        values: str,
        aggfunc: str = "mean",
        output: str = "",
    ) -> ToolResult:
        try:
            import pandas as pd
            df = DataAnalysisTool._to_df(df_or_path)
            pivot_df = df.pivot_table(index=index, columns=columns, values=values, aggfunc=aggfunc)
            if output:
                DataAnalysisTool.save(pivot_df.reset_index(), output)
            return ToolResult(True, f"✓ Pivot table: {pivot_df.shape}", {"df": pivot_df})
        except Exception as e:
            return ToolResult(False, f"✗ pivot failed: {e}")

    @staticmethod
    def time_series_analysis(
        df_or_path,
        date_col: str,
        value_col: str,
        output: str = "time_series.html",
    ) -> ToolResult:
        try:
            import pandas as pd
            import plotly.graph_objects as go
            df = DataAnalysisTool._to_df(df_or_path).copy()
            df[date_col] = pd.to_datetime(df[date_col])
            df = df.sort_values(date_col)
            # rolling mean
            df["rolling_7"]  = df[value_col].rolling(7,  min_periods=1).mean()
            df["rolling_30"] = df[value_col].rolling(30, min_periods=1).mean()
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=df[date_col], y=df[value_col],       name="Actual",    mode="lines"))
            fig.add_trace(go.Scatter(x=df[date_col], y=df["rolling_7"],     name="7d MA",     line=dict(dash="dash")))
            fig.add_trace(go.Scatter(x=df[date_col], y=df["rolling_30"],    name="30d MA",    line=dict(dash="dot")))
            fig.update_layout(title=f"Time Series: {value_col}", xaxis_title=date_col, yaxis_title=value_col)
            fig.write_html(output)
            stats = {"min": float(df[value_col].min()), "max": float(df[value_col].max()), "mean": float(df[value_col].mean()), "std": float(df[value_col].std())}
            return ToolResult(True, f"✓ Time-series chart saved to '{output}'", stats)
        except Exception as e:
            return ToolResult(False, f"✗ time_series_analysis failed: {e}")

    @staticmethod
    def correlation_matrix(
        df_or_path,
        method: str = "pearson",
        output: str = "correlation.html",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = DataAnalysisTool._to_df(df_or_path)
            num_df = df.select_dtypes(include="number")
            corr = num_df.corr(method=method).round(3)
            fig = px.imshow(corr, text_auto=True, color_continuous_scale="RdBu", title=f"{method.capitalize()} Correlation Matrix")
            fig.write_html(output)
            return ToolResult(True, f"✓ Correlation matrix saved to '{output}'", {"correlation": corr.to_dict()})
        except Exception as e:
            return ToolResult(False, f"✗ correlation_matrix failed: {e}")

    @staticmethod
    def outlier_detection(
        df_or_path,
        method: str = "iqr",
        threshold: float = 1.5,
    ) -> ToolResult:
        try:
            import numpy as np
            df = DataAnalysisTool._to_df(df_or_path)
            num_df = df.select_dtypes(include="number")
            outlier_mask = {}
            for col in num_df.columns:
                series = num_df[col].dropna()
                if method == "iqr":
                    q1, q3 = series.quantile(0.25), series.quantile(0.75)
                    iqr = q3 - q1
                    mask = (series < q1 - threshold * iqr) | (series > q3 + threshold * iqr)
                elif method == "zscore":
                    z = np.abs((series - series.mean()) / series.std())
                    mask = z > threshold
                else:
                    mask = series < -threshold
                outlier_mask[col] = int(mask.sum())
            total_outliers = sum(outlier_mask.values())
            return ToolResult(True, f"✓ Found {total_outliers} total outliers across {len(outlier_mask)} numeric columns", outlier_mask)
        except Exception as e:
            return ToolResult(False, f"✗ outlier_detection failed: {e}")

    @staticmethod
    def feature_importance(
        df_or_path,
        target_col: str,
        output: str = "feature_importance.html",
    ) -> ToolResult:
        try:
            import pandas as pd
            from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
            from sklearn.preprocessing import LabelEncoder
            import plotly.express as px
            df = DataAnalysisTool._to_df(df_or_path).dropna()
            X = df.drop(columns=[target_col])
            y = df[target_col]
            # encode categoricals
            for col in X.select_dtypes(include="object").columns:
                X[col] = LabelEncoder().fit_transform(X[col].astype(str))
            is_classifier = y.dtype == "object" or y.nunique() < 20
            if is_classifier:
                y = LabelEncoder().fit_transform(y.astype(str))
                model = RandomForestClassifier(n_estimators=100, random_state=42)
            else:
                model = RandomForestRegressor(n_estimators=100, random_state=42)
            model.fit(X, y)
            importances = dict(sorted(zip(X.columns, model.feature_importances_), key=lambda x: x[1], reverse=True))
            fig = px.bar(x=list(importances.keys()), y=list(importances.values()), title="Feature Importance", labels={"x": "Feature", "y": "Importance"})
            fig.write_html(output)
            return ToolResult(True, f"✓ Feature importance chart saved to '{output}'", importances)
        except Exception as e:
            return ToolResult(False, f"✗ feature_importance failed: {e}")

    @staticmethod
    def cluster_data(
        df_or_path,
        n_clusters: int = 3,
        output: str = "clusters.html",
    ) -> ToolResult:
        try:
            import pandas as pd
            from sklearn.cluster import KMeans
            from sklearn.preprocessing import StandardScaler
            from sklearn.decomposition import PCA
            import plotly.express as px
            df = DataAnalysisTool._to_df(df_or_path)
            num_df = df.select_dtypes(include="number").dropna()
            scaler = StandardScaler()
            scaled = scaler.fit_transform(num_df)
            kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
            labels = kmeans.fit_predict(scaled)
            pca = PCA(n_components=2)
            components = pca.fit_transform(scaled)
            plot_df = pd.DataFrame({"PC1": components[:, 0], "PC2": components[:, 1], "Cluster": labels.astype(str)})
            fig = px.scatter(plot_df, x="PC1", y="PC2", color="Cluster", title=f"K-Means Clustering (k={n_clusters})")
            fig.write_html(output)
            cluster_counts = {f"cluster_{i}": int((labels == i).sum()) for i in range(n_clusters)}
            return ToolResult(True, f"✓ Clustered into {n_clusters} groups, chart saved to '{output}'", cluster_counts)
        except Exception as e:
            return ToolResult(False, f"✗ cluster_data failed: {e}")

    @staticmethod
    def natural_language_query(
        df_or_path,
        query: str,
        model: str = "llama3.2:3b",
    ) -> ToolResult:
        try:
            import pandas as pd
            df = DataAnalysisTool._to_df(df_or_path)
            sample = df.head(5).to_string()
            dtypes = str(df.dtypes.to_dict())
            prompt = (
                f"You are a pandas expert. Given the following DataFrame info:\n"
                f"Columns and dtypes: {dtypes}\n"
                f"Sample (first 5 rows):\n{sample}\n\n"
                f"Write ONLY Python code using 'df' as the variable to answer:\n"
                f"'{query}'\n"
                f"Return a single expression or assignment. No explanation."
            )
            try:
                from npmai import Ollama
                llm = Ollama(model=model, temperature=0.1, change=True, Models=["mistral:7b"])
                code = llm.invoke(prompt).strip().strip("```python").strip("```").strip()
            except ImportError:
                return ToolResult(False, "✗ npmai Ollama not available. Install npmai package.")
            local_ns = {"df": df.copy(), "pd": pd}
            exec(code, local_ns)
            result = local_ns.get("result", local_ns.get("df"))
            if isinstance(result, pd.DataFrame):
                return ToolResult(True, f"✓ Query result: {len(result)} rows", {"df": result, "code": code})
            return ToolResult(True, f"✓ Query result: {result}", {"result": result, "code": code})
        except Exception as e:
            return ToolResult(False, f"✗ natural_language_query failed: {e}")

    @staticmethod
    def auto_visualize(df_or_path, output_folder: str = "auto_charts") -> ToolResult:
        try:
            import pandas as pd
            import plotly.express as px
            df = DataAnalysisTool._to_df(df_or_path)
            Path(output_folder).mkdir(parents=True, exist_ok=True)
            charts = []
            num_cols = df.select_dtypes(include="number").columns.tolist()
            cat_cols = df.select_dtypes(include="object").columns.tolist()
            # histogram for each numeric
            for col in num_cols[:5]:
                fig = px.histogram(df, x=col, title=f"Distribution of {col}")
                out = str(Path(output_folder) / f"hist_{col}.html")
                fig.write_html(out); charts.append(out)
            # correlation
            if len(num_cols) > 1:
                corr = df[num_cols].corr()
                fig  = px.imshow(corr, text_auto=True, title="Correlation Matrix")
                out  = str(Path(output_folder) / "correlation.html")
                fig.write_html(out); charts.append(out)
            # bar chart for top categorical
            for col in cat_cols[:3]:
                vc = df[col].value_counts().head(15)
                fig = px.bar(x=vc.index, y=vc.values, title=f"Top values: {col}")
                out = str(Path(output_folder) / f"bar_{col}.html")
                fig.write_html(out); charts.append(out)
            return ToolResult(True, f"✓ {len(charts)} charts saved to '{output_folder}'", charts)
        except Exception as e:
            return ToolResult(False, f"✗ auto_visualize failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 2. VisualizationTool
# ─────────────────────────────────────────────────────────────────────────────

class VisualizationTool:
    name = "visualization"
    description = (
        "Charts and interactive dashboards: bar, line, scatter, pie, heatmap, "
        "histogram, box, violin, geo map, sankey, treemap, sunburst, waterfall, candlestick."
    )
    use = (
        """
Name of Tool:- VisualizationTool,

Purpose of Tool:- 
The VisualizationTool provides a powerful, unified interface for generating a wide variety of interactive and static charts and dashboards using Plotly. 
It supports bar charts, line charts, scatter plots, pie charts, heatmaps, histograms, box plots, violin plots, geographic maps, Sankey diagrams, treemaps, sunburst charts, waterfall charts, candlestick charts, and multi-chart dashboards. 
Data can be provided as pandas DataFrames, lists, dictionaries, CSV, or Excel files. 
All charts can be exported as interactive HTML or static images (PNG, JPG, PDF). 
This tool is essential for data exploration, business intelligence reporting, interactive dashboards, and agentic data storytelling.

Methods:-
- _to_df: Internal helper to convert various data sources to pandas DataFrame.
- _save: Internal helper to save Plotly figures as HTML or image files.
- bar_chart: Creates bar (column) charts with optional stacking and animation.
- line_chart: Creates line charts with multiple series, filling, and markers.
- scatter_plot: Creates scatter plots with size, color, trendlines, and animation.
- pie_chart: Creates pie or donut charts with optional exploding slices.
- heatmap: Creates annotated heatmaps from pivot data.
- histogram: Creates histograms with optional KDE and grouping.
- box_plot: Creates box-and-whisker plots.
- violin_plot: Creates violin plots with box and points.
- geographic_map: Creates interactive scatter maps using Mapbox.
- create_dashboard: Combines multiple charts into a single HTML dashboard.
- sankey_diagram: Creates Sankey flow diagrams.
- treemap: Creates hierarchical treemap visualizations.
- sunburst: Creates sunburst (radial hierarchical) charts.
- waterfall_chart: Creates waterfall charts for cumulative changes.
- candlestick_chart: Creates financial candlestick charts.

How to use Tool Methods:-

1. bar_chart:
   - Purpose: Generates vertical or horizontal bar charts with optional color grouping and animation.
   - Arguments:
     a) data - DataFrame, list of dicts, CSV path, etc.
     b) x: str - Column for x-axis.
     c) y: str - Column for y-axis values.
     d) title: str (default: "Bar Chart")
     e) output: str (default: "bar_chart.html") - Supports .html, .png, .jpg, .pdf
     f) color: str (optional) - Column for color grouping.
     g) orientation: str (default: "v") - "v" or "h".
     h) stacked: bool (default: False)
     i) animated: bool (default: False)
   - How to call: VisualizationTool.bar_chart(data=df, x="category", y="sales", color="region", output="sales.html")

2. line_chart:
   - Purpose: Creates multi-line time series or trend charts.
   - Arguments:
     a) data
     b) x: str
     c) y_cols: str or List[str] - One or multiple y columns.
     d) title, output
     e) fill: bool (default: False) - Area fill under lines.
     f) markers: bool (default: False)
     g) log_scale: bool (default: False)
   - How to call: VisualizationTool.line_chart(data=df, x="date", y_cols=["sales", "profit"], fill=True)

3. scatter_plot:
   - Purpose: Creates scatter plots with optional size, color, trendline, and animation.
   - Arguments:
     a) data
     b) x, y
     c) size: str (optional) - Column for bubble size.
     d) color: str (optional)
     e) trendline: bool (default: False) - Adds OLS regression line.
     f) animated: bool (default: False)
     g) title, output
   - How to call: VisualizationTool.scatter_plot(data=df, x="age", y="income", color="gender", trendline=True)

4. pie_chart:
   - Purpose: Creates pie or donut charts.
   - Arguments:
     a) data
     b) values: str - Value column.
     c) names: str - Label column.
     d) hole: float (default: 0.0) - Donut hole ratio.
     e) explode: List[str] (optional) - Slice names to pull out.
     f) title, output
   - How to call: VisualizationTool.pie_chart(data=df, values="sales", names="product")

5. heatmap:
   - Purpose: Creates color-coded heatmap from pivot data.
   - Arguments:
     a) data
     b) x, y, z - Columns for pivot.
     c) colorscale: str (default: "Viridis")
     d) annotate: bool (default: True)
     e) title, output
   - How to call: VisualizationTool.heatmap(data=df, x="month", y="category", z="value")

6. histogram, box_plot, violin_plot, geographic_map, sankey_diagram, treemap, sunburst, waterfall_chart, candlestick_chart:
   - Similar structured arguments tailored to each chart type. See individual method signatures for details.
   - All support flexible data input and output to HTML/image formats.

7. create_dashboard:
   - Purpose: Combines multiple pre-generated chart HTML files into a single dashboard page.
   - Arguments:
     a) charts: List[str] - List of existing chart HTML file paths.
     b) layout: List[Tuple[int,int]] (optional) - Grid layout control.
     c) output_html: str (default: "dashboard.html")
     d) title: str (default: "Dashboard")
   - How to call: VisualizationTool.create_dashboard(charts=["chart1.html", "chart2.html"], title="Monthly Report")

**Note**: All chart methods return a ToolResult with success status, message, and output path. Interactive HTML charts are highly recommended for full interactivity.
""")
    
    @staticmethod
    def _to_df(data):
        import pandas as pd
        if isinstance(data, pd.DataFrame): return data
        if isinstance(data, (str, Path)):  return pd.read_csv(str(data)) if str(data).endswith(".csv") else pd.read_excel(str(data))
        if isinstance(data, list):         return pd.DataFrame(data)
        if isinstance(data, dict):         return pd.DataFrame(data)
        return data

    @staticmethod
    def _save(fig, output: str) -> str:
        import plotly
        if output.endswith(".html"):
            fig.write_html(output)
        elif output.endswith(".png") or output.endswith(".jpg"):
            fig.write_image(output)
        elif output.endswith(".pdf"):
            fig.write_image(output)
        else:
            fig.write_html(output + ".html")
            output = output + ".html"
        return output

    @staticmethod
    def bar_chart(
        data,
        x: str,
        y: str,
        title: str = "Bar Chart",
        output: str = "bar_chart.html",
        color: str = "",
        orientation: str = "v",
        stacked: bool = False,
        animated: bool = False,
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            kwargs: dict = {"x": x, "y": y, "title": title, "orientation": orientation}
            if color: kwargs["color"] = color
            if animated: kwargs["animation_frame"] = color or x
            if stacked and color: kwargs["barmode"] = "stack"
            fig = px.bar(df, **kwargs)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Bar chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ bar_chart failed: {e}")

    @staticmethod
    def line_chart(
        data,
        x: str,
        y_cols: Union[str, List[str]],
        title: str = "Line Chart",
        output: str = "line_chart.html",
        fill: bool = False,
        markers: bool = False,
        log_scale: bool = False,
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            y_list = [y_cols] if isinstance(y_cols, str) else y_cols
            kwargs: dict = {"x": x, "y": y_list, "title": title}
            if markers: kwargs["markers"] = True
            fig = px.line(df, **kwargs)
            if fill:
                for trace in fig.data: trace.fill = "tozeroy"
            if log_scale:
                fig.update_yaxes(type="log")
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Line chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ line_chart failed: {e}")

    @staticmethod
    def scatter_plot(
        data,
        x: str,
        y: str,
        size: str = "",
        color: str = "",
        title: str = "Scatter Plot",
        output: str = "scatter.html",
        trendline: bool = False,
        animated: bool = False,
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            kwargs: dict = {"x": x, "y": y, "title": title}
            if size and size in df.columns:  kwargs["size"]           = size
            if color and color in df.columns: kwargs["color"]          = color
            if trendline:                      kwargs["trendline"]      = "ols"
            if animated:                       kwargs["animation_frame"] = color or x
            fig = px.scatter(df, **kwargs)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Scatter plot saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ scatter_plot failed: {e}")

    @staticmethod
    def pie_chart(
        data,
        values: str,
        names: str,
        title: str = "Pie Chart",
        output: str = "pie_chart.html",
        hole: float = 0.0,
        explode: Optional[List[str]] = None,
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            kwargs: dict = {"values": values, "names": names, "title": title, "hole": hole}
            if explode:
                pull_vals = [0.1 if n in explode else 0 for n in df[names]]
                kwargs["pull"] = pull_vals
            fig = px.pie(df, **kwargs)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Pie chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ pie_chart failed: {e}")

    @staticmethod
    def heatmap(
        data,
        x: str,
        y: str,
        z: str,
        title: str = "Heatmap",
        output: str = "heatmap.html",
        colorscale: str = "Viridis",
        annotate: bool = True,
    ) -> ToolResult:
        try:
            import plotly.express as px
            import pandas as pd
            df = VisualizationTool._to_df(data)
            pivot = df.pivot_table(index=y, columns=x, values=z, aggfunc="mean")
            fig = px.imshow(pivot, color_continuous_scale=colorscale, title=title, text_auto=annotate)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Heatmap saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ heatmap failed: {e}")

    @staticmethod
    def histogram(
        data,
        column: str,
        bins: int = 30,
        title: str = "Histogram",
        output: str = "histogram.html",
        kde: bool = False,
        by: str = "",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            kwargs: dict = {"x": column, "nbins": bins, "title": title}
            if by and by in df.columns:
                kwargs["color"] = by
            if kde:
                import plotly.figure_factory as ff
                import numpy as np
                if by and by in df.columns:
                    groups = [df[df[by] == g][column].dropna().tolist() for g in df[by].unique()]
                    labels = list(df[by].unique())
                    fig = ff.create_distplot(groups, labels, show_hist=True)
                else:
                    fig = ff.create_distplot([df[column].dropna().tolist()], [column], show_hist=True)
                fig.update_layout(title=title)
            else:
                fig = px.histogram(df, **kwargs)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Histogram saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ histogram failed: {e}")

    @staticmethod
    def box_plot(
        data,
        x: str,
        y: str,
        title: str = "Box Plot",
        output: str = "box_plot.html",
        notched: bool = False,
        points: str = "outliers",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            fig = px.box(df, x=x, y=y, title=title, notched=notched, points=points)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Box plot saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ box_plot failed: {e}")

    @staticmethod
    def violin_plot(
        data,
        x: str,
        y: str,
        title: str = "Violin Plot",
        output: str = "violin_plot.html",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            fig = px.violin(df, x=x, y=y, title=title, box=True, points="all")
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Violin plot saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ violin_plot failed: {e}")

    @staticmethod
    def geographic_map(
        data,
        lat: str,
        lon: str,
        value: str = "",
        title: str = "Geographic Map",
        output: str = "geo_map.html",
        map_style: str = "open-street-map",
        zoom: int = 4,
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            kwargs: dict = {"lat": lat, "lon": lon, "title": title, "zoom": zoom, "mapbox_style": map_style}
            if value and value in df.columns:
                kwargs["color"] = value
                kwargs["size"]  = value
            fig = px.scatter_mapbox(df, **kwargs)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Geo map saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ geographic_map failed: {e}")

    @staticmethod
    def create_dashboard(
        charts: List[str],
        layout: Optional[List[Tuple[int, int]]] = None,
        output_html: str = "dashboard.html",
        title: str = "Dashboard",
    ) -> ToolResult:
        try:
            from plotly.subplots import make_subplots
            import plotly.io as pio
            # embed each chart HTML into a combined page
            chart_divs = []
            for chart_path in charts:
                if Path(chart_path).exists():
                    content = Path(chart_path).read_text()
                    # extract just the body content
                    body_match = re.search(r"<body[^>]*>(.*?)</body>", content, re.DOTALL)
                    if body_match:
                        chart_divs.append(body_match.group(1))
                    else:
                        chart_divs.append(content)
            cols = 2
            rows_html = []
            for i in range(0, len(chart_divs), cols):
                row_items = chart_divs[i:i+cols]
                row_cells = "".join(f"<td style='width:{100//cols}%;vertical-align:top'>{div}</td>" for div in row_items)
                rows_html.append(f"<tr>{row_cells}</tr>")
            html = f"""<!DOCTYPE html><html><head><title>{title}</title>
<style>body{{font-family:sans-serif;margin:0;padding:10px}} h1{{text-align:center}} table{{width:100%;border-collapse:collapse}}</style>
</head><body><h1>{title}</h1><table>{"".join(rows_html)}</table></body></html>"""
            Path(output_html).write_text(html)
            return ToolResult(True, f"✓ Dashboard with {len(charts)} charts saved to '{output_html}'", {"output": output_html})
        except Exception as e:
            return ToolResult(False, f"✗ create_dashboard failed: {e}")

    @staticmethod
    def sankey_diagram(
        sources: List[str],
        targets: List[str],
        values: List[float],
        output: str = "sankey.html",
    ) -> ToolResult:
        try:
            import plotly.graph_objects as go
            all_nodes = list(set(sources + targets))
            node_idx  = {n: i for i, n in enumerate(all_nodes)}
            fig = go.Figure(go.Sankey(
                node=dict(label=all_nodes, pad=15, thickness=20),
                link=dict(
                    source=[node_idx[s] for s in sources],
                    target=[node_idx[t] for t in targets],
                    value=values,
                ),
            ))
            fig.update_layout(title="Sankey Diagram")
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Sankey diagram saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ sankey_diagram failed: {e}")

    @staticmethod
    def treemap(
        data,
        path: List[str],
        values: str,
        output: str = "treemap.html",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            fig = px.treemap(df, path=path, values=values, title="Treemap")
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Treemap saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ treemap failed: {e}")

    @staticmethod
    def sunburst(
        data,
        path: List[str],
        values: str,
        output: str = "sunburst.html",
    ) -> ToolResult:
        try:
            import plotly.express as px
            df = VisualizationTool._to_df(data)
            fig = px.sunburst(df, path=path, values=values, title="Sunburst Chart")
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Sunburst chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ sunburst failed: {e}")

    @staticmethod
    def waterfall_chart(
        categories: List[str],
        values: List[float],
        output: str = "waterfall.html",
    ) -> ToolResult:
        try:
            import plotly.graph_objects as go
            measures = ["relative"] * len(values)
            if len(values):
                measures[-1] = "total"
            fig = go.Figure(go.Waterfall(
                x=categories,
                y=values,
                measure=measures,
                connector={"line": {"color": "rgb(63, 63, 63)"}},
            ))
            fig.update_layout(title="Waterfall Chart", showlegend=False)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Waterfall chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ waterfall_chart failed: {e}")

    @staticmethod
    def candlestick_chart(
        data,
        date: str,
        open: str,
        high: str,
        low: str,
        close: str,
        output: str = "candlestick.html",
    ) -> ToolResult:
        try:
            import plotly.graph_objects as go
            df = VisualizationTool._to_df(data)
            fig = go.Figure(go.Candlestick(
                x=df[date],
                open=df[open], high=df[high],
                low=df[low],   close=df[close],
            ))
            fig.update_layout(title="Candlestick Chart", xaxis_rangeslider_visible=True)
            out = VisualizationTool._save(fig, output)
            return ToolResult(True, f"✓ Candlestick chart saved to '{out}'", {"output": out})
        except Exception as e:
            return ToolResult(False, f"✗ candlestick_chart failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 3. WebScrapingAdvancedTool
# ─────────────────────────────────────────────────────────────────────────────

class WebScrapingAdvancedTool:
    name = "web_scraping_advanced"
    description = (
        "Production-grade web scraping: JS rendering, pagination, login-protected pages, "
        "bulk scraping, email/phone extraction, site mapping, screenshots, form submission."
    )
    use = (
        """Name of Tool:- WebScrapingAdvancedTool

Purpose of Tool:-
The WebScrapingAdvancedTool is a production-grade, enterprise web intelligence and crawling wrapper engineered to navigate, extract, and monitor data across complex web environments. It unifies high-performance synchronous networking (requests, BeautifulSoup) with automated browser headless execution engines (Playwright). This combination allows the tool to bypass modern scraping obstacles such as client-side JavaScript rendering, deep infinite scrolling layouts, multi-page paginated catalogs, and session-locked login walls. Beyond standard content ingestion, it includes specialized routines for bulk asynchronous parsing, full-frame viewport screenshots, structural site mapping, regex-driven data extraction (emails/phone numbers), metadata schema auditing (JSON-LD, Open Graph, Microdata), and transactional form submissions.

Methods:-
- scrape_with_js: Emulates an automated browser instance to execute JavaScript, scroll layouts, capture viewports, and query UI components.
- scrape_paginated: Automates cyclical pagination indexing workflows across nested web lists by programmatically executing page-forward actions.
- scrape_login_protected: Simulates interactive authentication workflows to inject access credentials before routing into access-restricted directories.
- extract_structured_data: Audits schema graphs embedded inside landing layers, parsing structured markup families like JSON-LD and Open Graph.
- monitor_page_changes: Runs background threads to verify source consistency by comparing page checksum hashes at scheduled intervals.
- bulk_scrape: Orchestrates distributed parallel data queries across an array of independent URLs using multithreaded network pools.
- extract_emails_phones: Scans raw textual payloads or remote endpoints via matching patterns to harvest valid contact credentials.
- map_website_structure: Crawls structural architectures recursively to generate relational hyperlink network maps.
- take_full_screenshot: Fires automated browser engines to record structural visual screenshots of a targeted web address.
- extract_all_links: Collects and lists hyperlinked paths discovered within specific source pages.
- download_all_images: Compiles asset links from target nodes and stores assets that meet specified minimum file-size criteria.
- fill_and_submit_form: Inputs transactional keys into text arrays or drop-down elements before executing programmatic target submission triggers.

How to use Tool Methods:-

1. scrape_with_js:
   - Purpose: Extracts dynamic components that require real-time execution of client-side scripts.
   - Arguments:
     a) url: str - Target website web address.
     b) selectors: Dict[str, str] - Key-value map of tracking names and their matching CSS selectors.
     c) wait_for: str (default: "") - Specific target element string to resolve prior to parsing.
     d) scroll: bool (default: False) - Flag determining if infinite layouts should scroll down before extraction.
     e) screenshot: bool (default: False) - Flag defining whether to save a visual snapshot alongside data extraction.
     f) output: str (default: "scraped_data.json") - Target storage file path.
   - Returns: ToolResult holding targeted component string matches.
   - How to call: WebScrapingAdvancedTool.scrape_with_js(url="https://example.com/ajax-feed", selectors={"titles": "h2.item-title"}, scroll=True)

2. scrape_paginated:
   - Purpose: Indexes lists across websites that span multiple sequential pages.
   - Arguments:
     a) base_url: str - Initial start page tracking path.
     b) next_selector: str - The CSS identifier targeting the "Next Page" navigational element.
     c) max_pages: int (default: 10) - Maximum index depth to crawl before terminating the loop.
     d) data_selectors: Optional[Dict[str, str]] (default: None) - Extraction fields to collect per page.
   - Returns: ToolResult storing structured lists mapped by page numbers.
   - How to call: WebScrapingAdvancedTool.scrape_paginated(base_url="https://example.com/shop", next_selector="a.next-page-btn", max_pages=5)

3. scrape_login_protected:
   - Purpose: Accesses secure internal dashboards behind authentication walls.
   - Arguments:
     a) url: str - Destination tracking target resource path.
     b) login_url: str - Gateway login address page.
     c) credentials: Dict[str, str] - Dictionary containing access values and alternative CSS input selectors (e.g., `username`, `password`).
     d) selectors: Optional[Dict[str, str]] (default: None) - Destination properties to extract post-authorization.
   - Returns: ToolResult containing protected data metrics.
   - How to call: WebScrapingAdvancedTool.scrape_login_protected(url="https://example.com/dashboard", login_url="https://example.com/login", credentials={"username": "user1", "password": "pass123"})

4. extract_structured_data:
   - Purpose: Collects formal schema objects hidden inside page source headers.
   - Arguments:
     a) url: str - Target address tracking path.
     b) schema_type: str (default: "all") - Targeted framework style choice ("all", "json-ld", "meta", "opengraph", "microdata").
   - Returns: ToolResult containing categorized structural schema records.
   - How to call: WebScrapingAdvancedTool.extract_structured_data(url="https://example.com/product/102", schema_type="json-ld")

5. monitor_page_changes:
   - Purpose: Monitors web endpoints in the background to track content updates or modifications.
   - Arguments:
     a) url: str - Tracking page address.
     b) interval: int (default: 60) - Checking frequency cadence measured in seconds.
     b) alert_on_change: bool (default: True) - Activates automated terminal logs if changes are identified.
     d) cred_key: str (default: "alerts") - Unique monitoring configuration tag.
   - Returns: ToolResult confirming active background thread attachment.
   - How to call: WebScrapingAdvancedTool.monitor_page_changes(url="https://example.com/status", interval=300)

6. bulk_scrape:
   - Purpose: Scrapes data from a batch of target URLs concurrently to optimize network throughput.
   - Arguments:
     a) urls: List[str] - Collection of target resource endpoints.
     b) selectors: Dict[str, str] - Property selectors applied uniformly across pages.
     c) concurrent: int (default: 3) - Number of parallel threads to run concurrently.
     d) delay: float (default: 0.5) - Time buffer pause injected between network calls.
     e) output: str (default: "bulk_scrape.json") - Storage destination path.
   - Returns: ToolResult containing an array of page summary dictionaries.
   - How to call: WebScrapingAdvancedTool.bulk_scrape(urls=["https://site.com/p1", "https://site.com/p2"], selectors={"header": "h1"})

7. extract_emails_phones:
   - Purpose: Scans raw text strings or live web pages for contact information.
   - Arguments:
     a) url_or_text: str - Live web URL string or block of plain text.
   - Returns: ToolResult containing list components for verified emails and phone matches.
   - How to call: WebScrapingAdvancedTool.extract_emails_phones(url_or_text="https://example.com/contact-us")

8. map_website_structure:
   - Purpose: Visualizes internal routing pathways to identify hidden directories within a domain.
   - Arguments:
     a) url: str - Starting site root url.
     b) depth: int (default: 2) - Recursive link tracking depth boundary.
     c) output: str (default: "site_map.json") - Map output tracking path.
   - Returns: ToolResult capturing page counts and structural layouts.
   - How to call: WebScrapingAdvancedTool.map_website_structure(url="https://example.com", depth=3)

9. take_full_screenshot:
   - Purpose: Captures high-resolution visual layouts of complex modern websites.
   - Arguments:
     a) url: str - Target endpoint address.
     b) output: str (default: "screenshot.png") - File save path destination.
     c) width: int (default: 1280) - Target viewport browser base width.
     d) height: int (default: 900) - Target viewport browser base height.
   - Returns: ToolResult highlighting final file payload footprints.
   - How to call: WebScrapingAdvancedTool.take_full_screenshot(url="https://example.com", output="exports/landing.png")

10. extract_all_links:
    - Purpose: Extracts all anchor linkages from a page for structural inspection.
    - Arguments:
      a) url: str - Target site address.
      b) internal_only: bool (default: False) - Restricts results exclusively to matching root domain configurations.
    - Returns: ToolResult bundling reference link text and destination targets.
    - How to call: WebScrapingAdvancedTool.extract_all_links(url="https://example.com/resources", internal_only=True)

11. download_all_images:
    - Purpose: Harvests remote graphical media matching specific file size parameters from a web page.
    - Arguments:
      a) url: str - Target asset repository source.
      b) output_folder: str (default: "images") - Directory location to store image files.
      c) min_size: int (default: 5000) - Minimum byte size threshold required to save an image file.
    - Returns: ToolResult declaring total downloaded asset tallies.
    - How to call: WebScrapingAdvancedTool.download_all_images(url="https://example.com/gallery", min_size=20000)

12. fill_and_submit_form:
    - Purpose: Automates interactive input fields and clicks submission components.
    - Arguments:
      a) url: str - Target page host form location.
      b) fields: Dict[str, str] - Key-value pair configurations mapping selector pathways to text values.
      c) submit_selector: str (default: "button[type='submit']") - Targeted trigger element component.
    - Returns: ToolResult logging structural final URL redirection parameters.
    - How to call: WebScrapingAdvancedTool.fill_and_submit_form(url="https://example.com/search", fields={"input#query": "AI tools"})
    """)
    
    @staticmethod
    def _get_page(url: str) -> "BeautifulSoup":
        import requests
        from bs4 import BeautifulSoup
        headers = {"User-Agent": "Mozilla/5.0 (compatible; NPMAgent/1.0)"}
        resp = requests.get(url, headers=headers, timeout=20)
        resp.raise_for_status()
        return BeautifulSoup(resp.text, "lxml")

    @staticmethod
    def scrape_with_js(
        url: str,
        selectors: Dict[str, str],
        wait_for: str = "",
        scroll: bool = False,
        screenshot: bool = False,
        output: str = "scraped_data.json",
    ) -> ToolResult:
        try:
            from playwright.sync_api import sync_playwright
            results: dict = {}
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True)
                page    = browser.new_page()
                page.goto(url, wait_until="networkidle", timeout=30000)
                if wait_for:
                    page.wait_for_selector(wait_for, timeout=10000)
                if scroll:
                    for _ in range(5):
                        page.evaluate("window.scrollBy(0, document.body.scrollHeight/5)")
                        time.sleep(0.5)
                for key, selector in selectors.items():
                    elements = page.query_selector_all(selector)
                    results[key] = [el.inner_text().strip() for el in elements]
                if screenshot:
                    page.screenshot(path=output.replace(".json", ".png"), full_page=True)
                browser.close()
            Path(output).write_text(json.dumps(results, indent=2))
            total = sum(len(v) for v in results.values())
            return ToolResult(True, f"✓ Scraped {total} items, saved to '{output}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ scrape_with_js failed: {e}")

    @staticmethod
    def scrape_paginated(
        base_url: str,
        next_selector: str,
        max_pages: int = 10,
        data_selectors: Optional[Dict[str, str]] = None,
    ) -> ToolResult:
        try:
            from playwright.sync_api import sync_playwright
            all_data: list = []
            selectors = data_selectors or {"text": "p, h2, h3"}
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True)
                page    = browser.new_page()
                page.goto(base_url, wait_until="networkidle", timeout=30000)
                for page_num in range(max_pages):
                    page_data: dict = {"page": page_num + 1}
                    for key, selector in selectors.items():
                        elements = page.query_selector_all(selector)
                        page_data[key] = [el.inner_text().strip() for el in elements[:50]]
                    all_data.append(page_data)
                    next_btn = page.query_selector(next_selector)
                    if not next_btn:
                        break
                    next_btn.click()
                    page.wait_for_load_state("networkidle", timeout=10000)
                browser.close()
            return ToolResult(True, f"✓ Scraped {len(all_data)} pages", all_data)
        except Exception as e:
            return ToolResult(False, f"✗ scrape_paginated failed: {e}")

    @staticmethod
    def scrape_login_protected(
        url: str,
        login_url: str,
        credentials: Dict[str, str],
        selectors: Optional[Dict[str, str]] = None,
    ) -> ToolResult:
        try:
            from playwright.sync_api import sync_playwright
            results: dict = {}
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True)
                page    = browser.new_page()
                page.goto(login_url, wait_until="networkidle", timeout=20000)
                # fill username/email
                for field in ("username", "email", "user"):
                    sel = credentials.get(f"{field}_selector", f"input[name='{field}'], input[type='email']")
                    try:
                        page.fill(sel, credentials.get(field, credentials.get("username", "")))
                        break
                    except Exception:
                        continue
                # fill password
                pwd_sel = credentials.get("password_selector", "input[type='password']")
                page.fill(pwd_sel, credentials.get("password", ""))
                # submit
                submit_sel = credentials.get("submit_selector", "button[type='submit']")
                page.click(submit_sel)
                page.wait_for_load_state("networkidle", timeout=15000)
                # now scrape target URL
                page.goto(url, wait_until="networkidle", timeout=20000)
                for key, selector in (selectors or {"content": "main"}).items():
                    elements = page.query_selector_all(selector)
                    results[key] = [el.inner_text().strip() for el in elements]
                browser.close()
            return ToolResult(True, f"✓ Scraped login-protected page at '{url}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ scrape_login_protected failed: {e}")

    @staticmethod
    def extract_structured_data(url: str, schema_type: str = "all") -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
            soup = BeautifulSoup(resp.text, "lxml")
            results: dict = {}
            # JSON-LD
            if schema_type in ("all", "json-ld"):
                scripts = soup.find_all("script", type="application/ld+json")
                results["json_ld"] = []
                for s in scripts:
                    try: results["json_ld"].append(json.loads(s.string))
                    except Exception: pass
            # Meta tags
            if schema_type in ("all", "meta"):
                results["meta"] = {m.get("name", m.get("property", "")): m.get("content", "") for m in soup.find_all("meta") if m.get("content")}
            # Open Graph
            if schema_type in ("all", "opengraph"):
                results["opengraph"] = {m.get("property", ""): m.get("content", "") for m in soup.find_all("meta", property=re.compile("^og:"))}
            # Microdata
            if schema_type in ("all", "microdata"):
                items = soup.find_all(itemscope=True)
                results["microdata"] = [{"type": item.get("itemtype", ""), "props": {prop.get("itemprop"): prop.get("content", prop.get_text(strip=True)) for prop in item.find_all(itemprop=True)}} for item in items[:10]]
            return ToolResult(True, f"✓ Extracted structured data from '{url}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ extract_structured_data failed: {e}")

    @staticmethod
    def monitor_page_changes(
        url: str,
        interval: int = 60,
        alert_on_change: bool = True,
        cred_key: str = "alerts",
    ) -> ToolResult:
        try:
            import requests, hashlib
            def _get_hash() -> str:
                resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
                return hashlib.md5(resp.text.encode()).hexdigest()

            last_hash = _get_hash()

            def _monitor():
                nonlocal last_hash
                while True:
                    time.sleep(interval)
                    try:
                        current_hash = _get_hash()
                        if current_hash != last_hash:
                            last_hash = current_hash
                            if alert_on_change:
                                from tools_data_research import MonitoringTool as MT
                                MT = None  # avoid circular; just print
                                print(f"[NPM Agent] Page changed: {url}")
                    except Exception:
                        pass

            t = threading.Thread(target=_monitor, daemon=True)
            t.start()
            return ToolResult(True, f"✓ Monitoring '{url}' every {interval}s (background thread)")
        except Exception as e:
            return ToolResult(False, f"✗ monitor_page_changes failed: {e}")

    @staticmethod
    def bulk_scrape(
        urls: List[str],
        selectors: Dict[str, str],
        concurrent: int = 3,
        delay: float = 0.5,
        output: str = "bulk_scrape.json",
    ) -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            from concurrent.futures import ThreadPoolExecutor, as_completed

            def _scrape_one(url: str) -> dict:
                time.sleep(delay)
                try:
                    resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
                    soup = BeautifulSoup(resp.text, "lxml")
                    row  = {"url": url, "status": resp.status_code}
                    for key, sel in selectors.items():
                        els = soup.select(sel)
                        row[key] = [e.get_text(strip=True) for e in els[:20]]
                    return row
                except Exception as exc:
                    return {"url": url, "error": str(exc)}

            results = []
            with ThreadPoolExecutor(max_workers=concurrent) as ex:
                futures = {ex.submit(_scrape_one, u): u for u in urls}
                for fut in as_completed(futures):
                    results.append(fut.result())

            Path(output).write_text(json.dumps(results, indent=2))
            ok = sum(1 for r in results if "error" not in r)
            return ToolResult(True, f"✓ Scraped {ok}/{len(urls)} URLs, saved to '{output}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ bulk_scrape failed: {e}")

    @staticmethod
    def extract_emails_phones(url_or_text: str) -> ToolResult:
        try:
            if url_or_text.startswith("http"):
                import requests
                text = requests.get(url_or_text, headers={"User-Agent": "Mozilla/5.0"}, timeout=20).text
            else:
                text = url_or_text
            emails = list(set(re.findall(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", text)))
            phones = list(set(re.findall(r"[\+]?[(]?[0-9]{1,4}[)]?[-\s\.]?[(]?[0-9]{1,4}[)]?[-\s\.]?[0-9]{4,8}", text)))
            return ToolResult(True, f"✓ Found {len(emails)} emails, {len(phones)} phones", {"emails": emails, "phones": phones})
        except Exception as e:
            return ToolResult(False, f"✗ extract_emails_phones failed: {e}")

    @staticmethod
    def map_website_structure(
        url: str,
        depth: int = 2,
        output: str = "site_map.json",
    ) -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            from urllib.parse import urljoin, urlparse

            visited: set = set()
            site_map: dict = {}
            base_domain = urlparse(url).netloc

            def _crawl(current_url: str, current_depth: int):
                if current_depth > depth or current_url in visited:
                    return
                visited.add(current_url)
                try:
                    resp = requests.get(current_url, headers={"User-Agent": "Mozilla/5.0"}, timeout=10)
                    soup = BeautifulSoup(resp.text, "lxml")
                    links = []
                    for a in soup.find_all("a", href=True):
                        full = urljoin(current_url, a["href"])
                        if urlparse(full).netloc == base_domain and full not in visited:
                            links.append(full)
                    site_map[current_url] = links
                    for link in links[:10]:
                        _crawl(link, current_depth + 1)
                except Exception:
                    site_map[current_url] = []

            _crawl(url, 0)
            Path(output).write_text(json.dumps(site_map, indent=2))
            return ToolResult(True, f"✓ Mapped {len(site_map)} pages, saved to '{output}'", {"pages": len(site_map)})
        except Exception as e:
            return ToolResult(False, f"✗ map_website_structure failed: {e}")

    @staticmethod
    def take_full_screenshot(
        url: str,
        output: str = "screenshot.png",
        width: int = 1280,
        height: int = 900,
    ) -> ToolResult:
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True)
                page    = browser.new_page(viewport={"width": width, "height": height})
                page.goto(url, wait_until="networkidle", timeout=30000)
                page.screenshot(path=output, full_page=True)
                browser.close()
            size = Path(output).stat().st_size
            return ToolResult(True, f"✓ Screenshot saved to '{output}' ({size} bytes)")
        except Exception as e:
            return ToolResult(False, f"✗ take_full_screenshot failed: {e}")

    @staticmethod
    def extract_all_links(url: str, internal_only: bool = False) -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            from urllib.parse import urljoin, urlparse
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
            soup = BeautifulSoup(resp.text, "lxml")
            base_domain = urlparse(url).netloc
            links = []
            for a in soup.find_all("a", href=True):
                full = urljoin(url, a["href"])
                if internal_only and urlparse(full).netloc != base_domain:
                    continue
                links.append({"href": full, "text": a.get_text(strip=True)[:80]})
            unique_links = list({l["href"]: l for l in links}.values())
            return ToolResult(True, f"✓ Found {len(unique_links)} links", unique_links)
        except Exception as e:
            return ToolResult(False, f"✗ extract_all_links failed: {e}")

    @staticmethod
    def download_all_images(
        url: str,
        output_folder: str = "images",
        min_size: int = 5000,
    ) -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            from urllib.parse import urljoin
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
            soup = BeautifulSoup(resp.text, "lxml")
            Path(output_folder).mkdir(parents=True, exist_ok=True)
            downloaded = 0
            for img in soup.find_all("img", src=True):
                src = urljoin(url, img["src"])
                try:
                    r = requests.get(src, timeout=10)
                    if len(r.content) >= min_size:
                        ext  = Path(src.split("?")[0]).suffix or ".jpg"
                        name = hashlib.md5(src.encode()).hexdigest()[:12] + ext
                        (Path(output_folder) / name).write_bytes(r.content)
                        downloaded += 1
                except Exception:
                    pass
            return ToolResult(True, f"✓ Downloaded {downloaded} images to '{output_folder}'")
        except Exception as e:
            return ToolResult(False, f"✗ download_all_images failed: {e}")

    @staticmethod
    def fill_and_submit_form(
        url: str,
        fields: Dict[str, str],
        submit_selector: str = "button[type='submit']",
    ) -> ToolResult:
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True)
                page    = browser.new_page()
                page.goto(url, wait_until="networkidle", timeout=20000)
                for selector, value in fields.items():
                    try:
                        page.fill(selector, value)
                    except Exception:
                        page.select_option(selector, value)
                page.click(submit_selector)
                page.wait_for_load_state("networkidle", timeout=15000)
                final_url = page.url
                browser.close()
            return ToolResult(True, f"✓ Form submitted, landed on: {final_url}", {"final_url": final_url})
        except Exception as e:
            return ToolResult(False, f"✗ fill_and_submit_form failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 4. SearchResearchTool
# ─────────────────────────────────────────────────────────────────────────────

class SearchResearchTool:
    name = "search_research"
    description = (
        "Academic and web research: arXiv, PubMed, Semantic Scholar, Wikipedia, "
        "Google Scholar, news, trending topics, patents."
    )
    use = (
        """
Name of Tool:- SearchResearchTool,

Purpose of Tool:- 
The SearchResearchTool provides powerful academic, scientific, and general web research capabilities. 
It integrates with arXiv, PubMed, Semantic Scholar, Wikipedia, Google Scholar, news sources, trending topics, and patent databases. 
The tool supports searching, retrieving full details, downloading PDFs/abstracts, getting citations, and accessing structured data from authoritative sources. 
It is designed for literature reviews, patent research, news monitoring, fact-checking, trend analysis, and agentic research automation.

Methods:-
- search_arxiv: Searches arXiv for academic papers with optional category and date filtering.
- get_arxiv_paper: Retrieves detailed information about a specific arXiv paper (optionally downloads PDF).
- search_pubmed: Searches PubMed for biomedical literature and abstracts.
- search_semantic_scholar: Searches Semantic Scholar for scholarly papers with rich metadata.
- get_citations: Retrieves citations for a paper from Semantic Scholar.
- search_wikipedia: Searches or retrieves Wikipedia page summaries.
- get_wikipedia_page: Fetches full content of a Wikipedia page with cleaning options.
- search_google_scholar: Performs Google Scholar search with scraping for titles, snippets, and links.
- search_news: Searches for news articles (NewsAPI preferred, Google News RSS fallback).
- get_trending_topics: Retrieves current trending search topics by region.
- search_patents: Searches for patents with structured metadata.

How to use Tool Methods:-

1. search_arxiv:
   - Purpose: Searches the arXiv preprint repository for scientific papers.
   - Arguments:
     a) query: str - Search query (supports advanced syntax).
     b) max_results: int (default: 10).
     c) categories: List[str] (optional) - arXiv categories (e.g., ["cs.AI", "physics"]).
     d) date_from: str (optional) - Filter papers published after this date (YYYY-MM-DD).
     e) download_pdfs: bool (default: False) - Automatically download PDFs.
     f) output_folder: str (default: "arxiv_papers") - Folder for downloaded PDFs.
   - Returns: List of paper metadata including title, authors, abstract, PDF URL, etc.
   - How to call: 
     SearchResearchTool.search_arxiv(
         query="large language models",
         categories=["cs.AI"],
         max_results=15,
         download_pdfs=True
     )

2. get_arxiv_paper:
   - Purpose: Retrieves full details for a specific arXiv paper by ID.
   - Arguments:
     a) paper_id: str - arXiv ID (e.g., "2406.12345" or full URL).
     b) download: bool (default: False) - Download the PDF.
     c) output: str (default: ".") - Download directory.
   - Returns: Comprehensive paper metadata.
   - How to call: SearchResearchTool.get_arxiv_paper(paper_id="2406.12345", download=True)

3. search_pubmed:
   - Purpose: Searches PubMed for biomedical and life sciences literature.
   - Arguments:
     a) query: str
     b) max_results: int (default: 10)
     c) download_abstracts: bool (default: False) - Saves all abstracts to a text file.
   - Returns: List of PubMed IDs with links (abstracts optionally saved).
   - How to call: SearchResearchTool.search_pubmed(query="CRISPR gene editing", max_results=20)

4. search_semantic_scholar:
   - Purpose: Searches Semantic Scholar for academic papers with rich metadata (citations, authors, abstracts).
   - Arguments:
     a) query: str
     b) fields: List[str] (optional) - Fields to return.
     c) limit: int (default: 10)
   - Returns: List of paper objects.
   - How to call: SearchResearchTool.search_semantic_scholar(query="transformer architecture", limit=15)

5. get_citations:
   - Purpose: Retrieves papers that cite a given paper (Semantic Scholar).
   - Arguments:
     a) paper_id: str - Semantic Scholar paper ID.
     b) source: str (default: "semantic_scholar")
   - Returns: List of citing papers.
   - How to call: SearchResearchTool.get_citations(paper_id="paper_id_here")

6. search_wikipedia:
   - Purpose: Searches Wikipedia or retrieves a page summary.
   - Arguments:
     a) query: str
     b) language: str (default: "en")
     c) sentences: int (default: 5) - Summary length control.
   - Returns: Summary or search results.
   - How to call: SearchResearchTool.search_wikipedia(query="Quantum computing")

7. get_wikipedia_page:
   - Purpose: Fetches the full cleaned content of a Wikipedia article.
   - Arguments:
     a) title: str - Exact Wikipedia page title.
     b) language: str (default: "en")
     c) output_format: str (default: "text") - "text" or "html".
   - Returns: Page title and cleaned content.
   - How to call: SearchResearchTool.get_wikipedia_page(title="Artificial intelligence")

8. search_google_scholar:
   - Purpose: Scrapes Google Scholar for academic search results.
   - Arguments:
     a) query: str
     b) num_results: int (default: 10)
     c) year_from, year_to: int (optional) - Year range filter.
   - Returns: List of results with title, snippet, citations, etc.
   - How to call: SearchResearchTool.search_google_scholar(query="reinforcement learning from human feedback", num_results=15)

9. search_news:
   - Purpose: Searches for recent news articles (prefers NewsAPI, falls back to Google News RSS).
   - Arguments:
     a) query: str
     b) sources: List[str] (optional)
     c) language: str (default: "en")
     d) date_from, date_to: str (optional)
     e) max_results: int (default: 20)
   - Returns: List of news articles with title, source, URL, etc.
   - How to call: SearchResearchTool.search_news(query="AI breakthrough", date_from="2026-06-01")

10. get_trending_topics:
    - Purpose: Retrieves currently trending search topics by region.
    - Arguments:
      a) region: str (default: "US")
      b) category: str (default: "all")
    - Returns: List of trending topics with traffic estimates.
    - How to call: SearchResearchTool.get_trending_topics(region="IN")

11. search_patents:
    - Purpose: Searches for patents using PatentsView API.
    - Arguments:
      a) query: str
      b) country: str (default: "US")
      c) date_from, date_to: str (optional)
    - Returns: List of patent records with titles, abstracts, dates, etc.
    - How to call: SearchResearchTool.search_patents(query="neural network", country="US")
""")
    
    @staticmethod
    def search_arxiv(
        query: str,
        max_results: int = 10,
        categories: Optional[List[str]] = None,
        date_from: str = "",
        download_pdfs: bool = False,
        output_folder: str = "arxiv_papers",
    ) -> ToolResult:
        try:
            import arxiv
            search_query = query
            if categories:
                cat_filter = " OR ".join(f"cat:{c}" for c in categories)
                search_query = f"({query}) AND ({cat_filter})"
            client = arxiv.Client()
            search = arxiv.Search(query=search_query, max_results=max_results, sort_by=arxiv.SortCriterion.Relevance)
            papers = []
            for result in client.results(search):
                paper_data = {
                    "id":       result.get_short_id(),
                    "title":    result.title,
                    "authors":  [str(a) for a in result.authors[:5]],
                    "abstract": result.summary[:500],
                    "published":str(result.published.date()) if result.published else "",
                    "url":      result.entry_id,
                    "pdf_url":  result.pdf_url,
                    "categories": result.categories,
                }
                if date_from and paper_data["published"] < date_from:
                    continue
                papers.append(paper_data)
                if download_pdfs:
                    Path(output_folder).mkdir(parents=True, exist_ok=True)
                    try:
                        result.download_pdf(dirpath=output_folder)
                    except Exception:
                        pass
            return ToolResult(True, f"✓ Found {len(papers)} arXiv papers", papers)
        except Exception as e:
            return ToolResult(False, f"✗ search_arxiv failed: {e}")

    @staticmethod
    def get_arxiv_paper(
        paper_id: str,
        download: bool = False,
        output: str = ".",
    ) -> ToolResult:
        try:
            import arxiv
            client = arxiv.Client()
            search = arxiv.Search(id_list=[paper_id])
            results = list(client.results(search))
            if not results:
                return ToolResult(False, f"✗ Paper '{paper_id}' not found")
            r = results[0]
            data = {
                "id":         r.get_short_id(),
                "title":      r.title,
                "authors":    [str(a) for a in r.authors],
                "abstract":   r.summary,
                "published":  str(r.published.date()) if r.published else "",
                "pdf_url":    r.pdf_url,
                "categories": r.categories,
                "comment":    r.comment or "",
                "doi":        r.doi or "",
            }
            if download:
                r.download_pdf(dirpath=output)
                data["downloaded_to"] = output
            return ToolResult(True, f"✓ Got arXiv paper: {r.title[:60]}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_arxiv_paper failed: {e}")

    @staticmethod
    def search_pubmed(
        query: str,
        max_results: int = 10,
        download_abstracts: bool = False,
    ) -> ToolResult:
        try:
            from Bio import Entrez
            Entrez.email = "npmaiagent@example.com"
            handle = Entrez.esearch(db="pubmed", term=query, retmax=max_results)
            record = Entrez.read(handle)
            ids    = record.get("IdList", [])
            if not ids:
                return ToolResult(True, "✓ No PubMed results found", [])
            fetch_handle = Entrez.efetch(db="pubmed", id=ids, rettype="abstract", retmode="text")
            abstracts    = fetch_handle.read()
            papers = []
            for pmid in ids:
                papers.append({"pmid": pmid, "url": f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/"})
            if download_abstracts:
                Path("pubmed_abstracts.txt").write_text(abstracts)
                return ToolResult(True, f"✓ Found {len(papers)} PubMed papers, abstracts saved", {"papers": papers, "abstracts_file": "pubmed_abstracts.txt"})
            return ToolResult(True, f"✓ Found {len(papers)} PubMed papers", papers)
        except Exception as e:
            return ToolResult(False, f"✗ search_pubmed failed: {e}")

    @staticmethod
    def search_semantic_scholar(
        query: str,
        fields: Optional[List[str]] = None,
        limit: int = 10,
    ) -> ToolResult:
        try:
            import requests
            default_fields = "title,year,authors,abstract,citationCount,url"
            fields_str = ",".join(fields) if fields else default_fields
            resp = requests.get(
                "https://api.semanticscholar.org/graph/v1/paper/search",
                params={"query": query, "limit": limit, "fields": fields_str},
                timeout=20,
            )
            data  = resp.json()
            papers = data.get("data", [])
            return ToolResult(True, f"✓ Found {len(papers)} Semantic Scholar papers", papers)
        except Exception as e:
            return ToolResult(False, f"✗ search_semantic_scholar failed: {e}")

    @staticmethod
    def get_citations(paper_id: str, source: str = "semantic_scholar") -> ToolResult:
        try:
            import requests
            if source == "semantic_scholar":
                resp = requests.get(
                    f"https://api.semanticscholar.org/graph/v1/paper/{paper_id}/citations",
                    params={"fields": "title,year,authors", "limit": 50},
                    timeout=20,
                )
                citations = [c.get("citingPaper", {}) for c in resp.json().get("data", [])]
                return ToolResult(True, f"✓ Found {len(citations)} citations", citations)
            return ToolResult(False, f"✗ Source '{source}' not supported")
        except Exception as e:
            return ToolResult(False, f"✗ get_citations failed: {e}")

    @staticmethod
    def search_wikipedia(
        query: str,
        language: str = "en",
        sentences: int = 5,
    ) -> ToolResult:
        try:
            import requests
            resp = requests.get(
                f"https://{language}.wikipedia.org/api/rest_v1/page/summary/{query.replace(' ', '_')}",
                timeout=15,
            )
            if resp.status_code == 404:
                # try search
                search_resp = requests.get(
                    f"https://{language}.wikipedia.org/w/api.php",
                    params={"action": "query", "list": "search", "srsearch": query, "format": "json", "srlimit": 5},
                    timeout=15,
                )
                results = search_resp.json().get("query", {}).get("search", [])
                return ToolResult(True, f"✓ Wikipedia search: {len(results)} results", results)
            data = resp.json()
            return ToolResult(True, f"✓ Wikipedia: {data.get('title')}", {
                "title":   data.get("title"),
                "summary": data.get("extract", "")[:sentences * 200],
                "url":     data.get("content_urls", {}).get("desktop", {}).get("page", ""),
                "image":   data.get("originalimage", {}).get("source", ""),
            })
        except Exception as e:
            return ToolResult(False, f"✗ search_wikipedia failed: {e}")

    @staticmethod
    def get_wikipedia_page(
        title: str,
        language: str = "en",
        output_format: str = "text",
    ) -> ToolResult:
        try:
            import requests
            resp = requests.get(
                f"https://{language}.wikipedia.org/w/api.php",
                params={"action": "parse", "page": title, "prop": "wikitext|text", "format": "json"},
                timeout=20,
            )
            data = resp.json()
            if "error" in data:
                return ToolResult(False, f"✗ Wikipedia page not found: {title}")
            if output_format == "html":
                content = data.get("parse", {}).get("text", {}).get("*", "")
            else:
                raw = data.get("parse", {}).get("wikitext", {}).get("*", "")
                content = re.sub(r"\[\[([^\]|]+\|)?([^\]]+)\]\]", r"\2", raw)
                content = re.sub(r"\{\{.*?\}\}", "", content, flags=re.DOTALL)
                content = re.sub(r"<[^>]+>", "", content)
            return ToolResult(True, f"✓ Got Wikipedia page: {title}", {"title": title, "content": content[:10000]})
        except Exception as e:
            return ToolResult(False, f"✗ get_wikipedia_page failed: {e}")

    @staticmethod
    def search_google_scholar(
        query: str,
        num_results: int = 10,
        year_from: int = 0,
        year_to: int = 0,
    ) -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            params: dict = {"q": query, "hl": "en"}
            if year_from: params["as_ylo"] = year_from
            if year_to:   params["as_yhi"] = year_to
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            resp = requests.get("https://scholar.google.com/scholar", params=params, headers=headers, timeout=20)
            soup = BeautifulSoup(resp.text, "lxml")
            results = []
            for div in soup.select(".gs_r.gs_or.gs_scl")[:num_results]:
                title_tag   = div.select_one(".gs_rt a")
                snippet_tag = div.select_one(".gs_rs")
                info_tag    = div.select_one(".gs_a")
                cite_tag    = div.select_one(".gs_fl a")
                results.append({
                    "title":     title_tag.get_text(strip=True) if title_tag else "",
                    "url":       title_tag["href"] if title_tag and title_tag.get("href") else "",
                    "snippet":   snippet_tag.get_text(strip=True) if snippet_tag else "",
                    "info":      info_tag.get_text(strip=True) if info_tag else "",
                    "citations": cite_tag.get_text(strip=True) if cite_tag else "",
                })
            return ToolResult(True, f"✓ {len(results)} Google Scholar results", results)
        except Exception as e:
            return ToolResult(False, f"✗ search_google_scholar failed: {e}")

    @staticmethod
    def search_news(
        query: str,
        sources: Optional[List[str]] = None,
        language: str = "en",
        date_from: str = "",
        date_to: str = "",
        max_results: int = 20,
    ) -> ToolResult:
        try:
            import requests
            creds = CredStore.load("newsapi")
            api_key = creds.get("api_key", "")
            if api_key:
                params: dict = {"q": query, "language": language, "pageSize": max_results, "apiKey": api_key}
                if sources: params["sources"] = ",".join(sources)
                if date_from: params["from"] = date_from
                if date_to:   params["to"]   = date_to
                resp = requests.get("https://newsapi.org/v2/everything", params=params, timeout=20)
                data = resp.json()
                articles = [
                    {"title": a["title"], "source": a["source"]["name"], "url": a["url"], "published": a["publishedAt"], "description": a.get("description", "")}
                    for a in data.get("articles", [])
                ]
                return ToolResult(True, f"✓ {len(articles)} news articles", articles)
            # fallback: Google News RSS
            rss_url = f"https://news.google.com/rss/search?q={requests.utils.quote(query)}&hl={language}&gl=US"
            resp = requests.get(rss_url, timeout=20)
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(resp.text, "xml")
            items = [{"title": item.find("title").get_text(), "url": item.find("link").get_text() if item.find("link") else "", "published": item.find("pubDate").get_text() if item.find("pubDate") else ""} for item in soup.find_all("item")[:max_results]]
            return ToolResult(True, f"✓ {len(items)} news articles (RSS)", items)
        except Exception as e:
            return ToolResult(False, f"✗ search_news failed: {e}")

    @staticmethod
    def get_trending_topics(region: str = "US", category: str = "all") -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            url = f"https://trends.google.com/trends/trendingsearches/daily/rss?geo={region}"
            resp = requests.get(url, timeout=20)
            soup = BeautifulSoup(resp.text, "xml")
            trends = []
            for item in soup.find_all("item"):
                title_tag    = item.find("title")
                traffic_tag  = item.find("ht:approx_traffic") or item.find("approx_traffic")
                trends.append({
                    "topic":   title_tag.get_text() if title_tag else "",
                    "traffic": traffic_tag.get_text() if traffic_tag else "",
                })
            return ToolResult(True, f"✓ {len(trends)} trending topics in '{region}'", trends)
        except Exception as e:
            return ToolResult(False, f"✗ get_trending_topics failed: {e}")

    @staticmethod
    def search_patents(
        query: str,
        country: str = "US",
        date_from: str = "",
        date_to: str = "",
    ) -> ToolResult:
        try:
            import requests
            params: dict = {"q": query, "o": "json", "stats": "false"}
            if country:   params["country"] = country
            if date_from: params["dateRange"] = f"{date_from},{date_to or datetime.now().strftime('%Y-%m-%d')}"
            resp = requests.get(
                "https://patentsview.org/api/patents/query",
                params={"q": json.dumps({"_text_any": {"patent_title": query}}), "f": json.dumps(["patent_id", "patent_title", "patent_date", "patent_abstract", "assignee_organization"]), "o": json.dumps({"patent_date": "desc"}), "per_page": 25},
                timeout=30,
            )
            data    = resp.json()
            patents = data.get("patents", []) or []
            return ToolResult(True, f"✓ Found {len(patents)} patents", patents)
        except Exception as e:
            return ToolResult(False, f"✗ search_patents failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 5. FinancialDataTool
# ─────────────────────────────────────────────────────────────────────────────

class FinancialDataTool:
    name = "financial_data"
    description = (
        "Market and financial data: stocks, crypto, forex, commodities, "
        "technical indicators, earnings, economic indicators, options, portfolio analysis."
    )
    use = (
        """Name of Tool:- FinancialDataTool

Purpose of Tool:-
The FinancialDataTool is an all-in-one financial data ingestion, modeling, and analysis suite designed to pull and synthesize live market data across multiple asset classes. Leveraging popular programmatic wrappers and APIs (including `yfinance`, `ccxt`, `World Bank API`, and `CoinGecko`), it streamlines the retrieval of real-time quotes, technical overlays, fundamental statement data, and global macroeconomic metrics. Beyond simple data lookups, it includes processing methods for automated stock screening against valuation targets, standard financial derivative options-chain mapping, and dynamic interactive portfolio asset allocation visualization.

Methods:-
- get_stock_price: Retrieves standard historical equity charting and open-high-low-close pricing blocks for individual equities.
- get_multiple_stocks: Aggregates historical adjusted closing timelines for a list of equities concurrently.
- get_company_info: Compiles fundamental descriptive data profiles, capitalization structures, and operating margins for a public corporation.
- get_financial_statements: Parses and transforms standard accounting sheets, including the corporate income statement, balance sheet, or cash flow history.
- get_earnings_calendar: Extracts short-term projected corporate reporting timelines from live financial networks.
- get_economic_indicators: Interfaces with international databases to query historical sovereign economic health statistics.
- get_crypto_price: Collects daily open-high-low-close transaction series from primary token spot pairs across prominent global order books.
- get_crypto_info: Generates operational overviews, market caps, and circulating metrics for individual digital assets.
- get_forex_rate: Queries interbank foreign currency conversion quotes and generates multi-day performance profiles.
- get_commodity_prices: Pulls financial exchange pricing sequences for standard resource instruments.
- calculate_technical_indicators: Runs mathematical rolling transformations across pricing frames to build indicators like SMAs, RSI, MACD, and Bollinger Bands.
- screen_stocks: Evaluates filter bounds across standard index constituents to target equities fitting strict fundamental criteria.
- get_options_chain: Builds immediate expiration structures for call and put options contract listings.
- portfolio_analysis: Evaluates aggregated position weight allocations, calculates total unrealized adjustments, and generates HTML allocation pie charts.

How to use Tool Methods:-

1. get_stock_price:
   - Purpose: Obtains detailed historical data records for an isolated public ticker.
   - Arguments:
     a) symbol: str - Upper or lower-case company equity market ticker.
     b) period: str (default: "1mo") - Timing data scope window (e.g., "1d", "5d", "1mo", "1y").
     c) interval: str (default: "1d") - The frequency step dividing separate data candles (e.g., "1m", "5m", "1d").
   - Returns: ToolResult holding an array of historical dictionary records and a separate block for the latest market quote.
   - How to call: FinancialDataTool.get_stock_price(symbol="AAPL", period="6mo", interval="1wk")

2. get_multiple_stocks:
   - Purpose: Collects batch historical tracking records spanning multiple tickers inside an efficient structured package.
   - Arguments:
     a) symbols: List[str] - Target array containing clean company market indicators.
     b) period: str (default: "1mo") - Total backward duration time bucket.
   - Returns: ToolResult presenting combined column arrays mapped directly against specified dates.
   - How to call: FinancialDataTool.get_multiple_stocks(symbols=["MSFT", "GOOGL", "AMZN"], period="1y")

3. get_company_info:
   - Purpose: Audits corporate operational profiles, personnel counts, and market ratios.
   - Arguments:
     a) symbol: str - Targeted asset market ticker.
   - Returns: ToolResult filtering corporate descriptive metrics, 52-week ranges, forward P/E targets, and summary briefs.
   - How to call: FinancialDataTool.get_company_info(symbol="NVDA")

4. get_financial_statements:
   - Purpose: Downloads structured filings needed to calculate core corporate valuations.
   - Arguments:
     a) symbol: str - target financial market indicator.
     b) statement_type: str (default: "income") - Core structural accounting model type ("income", "balance", "cashflow").
   - Returns: ToolResult presenting timestamped accounting categories across multiple fiscal years.
   - How to call: FinancialDataTool.get_financial_statements(symbol="TSLA", statement_type="balance")

5. get_earnings_calendar:
   - Purpose: Identifies short-term public earnings announcements to handle sudden volatility risks.
   - Arguments:
     a) date_from: str (default: "") - Starting time window filter.
     b) date_to: str (default: "") - End point time window constraint.
   - Returns: ToolResult outputting upcoming company names, ticker references, and schedule timestamps.
   - How to call: FinancialDataTool.get_earnings_calendar()

6. get_economic_indicators:
   - Purpose: Imports sovereign structural indicators directly from global databases.
   - Arguments:
     a) indicator: str (default: "gdp") - Targeted metric category code or literal name ("gdp", "inflation", "unemployment", "interest_rate").
     b) country: str (default: "US") - standard ISO two or three-character country boundary code.
   - Returns: ToolResult mapping historical value data sequences explicitly against calendar years.
   - How to call: FinancialDataTool.get_economic_indicators(indicator="inflation", country="DE")

7. get_crypto_price:
   - Purpose: Tracks trading variations of major digital tokens against traditional fiat.
   - Arguments:
     a) symbol: str (default: "BTC") - The asset token network designation.
     b) currency: str (default: "USD") - Counter transaction pair quote medium.
     c) period: str (default: "30") - Number of retroactive historical days to sample.
   - Returns: ToolResult tracking price moves alongside the final logged candle.
   - How to call: FinancialDataTool.get_crypto_price(symbol="ETH", currency="EUR", period="15")

8. get_crypto_info:
   - Purpose: Collects capitalization weightings, market rank listings, and global descriptions for digital assets.
   - Arguments:
     a) symbol: str - Digital asset network tracking identifier.
   - Returns: ToolResult reporting overall rank, valuation metrics, and short informational paragraphs.
   - How to call: FinancialDataTool.get_crypto_info(symbol="SOL")

9. get_forex_rate:
   - Purpose: Ingests current conversion rates and changes affecting international currency markets.
   - Arguments:
     a) from_currency: str (default: "USD") - Base settlement denomination source.
     b) to_currency: str (default: "EUR") - Target exchange settlement conversion type.
     c) period: str (default: "1mo") - Total time length frame for baseline data mapping.
   - Returns: ToolResult with exact four-decimal exchange calculations alongside historical details.
   - How to call: FinancialDataTool.get_forex_rate(from_currency="GBP", to_currency="JPY")

10. get_commodity_prices:
    - Purpose: Gathers price information for physical resource products and raw trade items.
    - Arguments:
      a) commodity: str (default: "gold") - Chosen resource asset profile keyword ("gold", "silver", "oil", "natural_gas").
    - Returns: ToolResult showing active exchange market quotes and historical timelines.
    - How to call: FinancialDataTool.get_commodity_prices(commodity="oil")

11. calculate_technical_indicators:
    - Purpose: Generates algorithmic math overlays used to build quantitative trading indicators.
    - Arguments:
      a) symbol: str - target financial investment instrument ticker.
      b) indicators: Optional[List[str]] (default: None) - Requested metrics package array ("sma20", "rsi", "macd", "bollinger", "atr").
      c) period: str (default: "3mo") - Total structural tracking depth data window.
    - Returns: ToolResult capturing structural values for all calculated metrics.
    - How to call: FinancialDataTool.calculate_technical_indicators(symbol="AMD", indicators=["rsi", "macd"])

12. screen_stocks:
    - Purpose: Filters broad company indices down to high-value lists fitting strict financial metrics.
    - Arguments:
      a) market_cap_min: float (default: 1e9) - Minimum floor limitation required for total company market valuation.
      b) pe_max: float (default: 30.0) - Ceiling multiplier ceiling restriction parameter for Price-to-Earnings ratios.
      c) dividend_min: float (default: 0.0) - Minimum required percent yield parameter.
      d) sector: str (default: "") - Specific business industry structural designation string filter.
    - Returns: ToolResult delivering matching filtered stocks along with quick-glance data tables.
    - How to call: FinancialDataTool.screen_stocks(market_cap_min=50e9, pe_max=15.0, sector="Technology")

13. get_options_chain:
    - Purpose: Looks up derivative option lines to inspect strike ranges and implied volatility spreads.
    - Arguments:
      a) symbol: str - Asset instrument market option identifier.
      b) expiry: str (default: "") - Target string tracking date format (e.g., "YYYY-MM-DD"). Drops to nearest date if unspecified.
    - Returns: ToolResult separating Call and Put metrics including bids, asks, and open interest tallies.
    - How to call: FinancialDataTool.get_options_chain(symbol="META", expiry="2026-07-17")

14. portfolio_analysis:
    - Purpose: Analyzes investment balances, calculates gains/losses, and outputs visual allocation breakdowns.
    - Arguments:
      a) holdings: List[Dict[str, Any]] - Collection arrays grouping specific dictionaries with keys `symbol`, `shares`, and `avg_cost`.
      b) output: str (default: "portfolio.html") - File target storage destination path.
    - Returns: ToolResult declaring absolute valuation summaries, portfolio returns, and interactive chart exports.
    - How to call: FinancialDataTool.portfolio_analysis(holdings=[{"symbol": "AAPL", "shares": 50, "avg_cost": 175.0}, {"symbol": "MSFT", "shares": 20, "avg_cost": 350.0}])
    """)
    
    @staticmethod
    def get_stock_price(
        symbol: str,
        period: str = "1mo",
        interval: str = "1d",
    ) -> ToolResult:
        try:
            import yfinance as yf
            ticker = yf.Ticker(symbol.upper())
            hist   = ticker.history(period=period, interval=interval)
            if hist.empty:
                return ToolResult(False, f"✗ No data for '{symbol}'")
            data = hist.reset_index().to_dict(orient="records")
            latest = data[-1]
            return ToolResult(True, f"✓ {symbol}: latest close = {latest.get('Close', 'N/A'):.2f}", {"history": data, "latest": latest})
        except Exception as e:
            return ToolResult(False, f"✗ get_stock_price failed: {e}")

    @staticmethod
    def get_multiple_stocks(symbols: List[str], period: str = "1mo") -> ToolResult:
        try:
            import yfinance as yf
            import pandas as pd
            syms = [s.upper() for s in symbols]
            data = yf.download(syms, period=period, progress=False)["Close"]
            if isinstance(data, pd.Series):
                data = data.to_frame(name=syms[0])
            result = {"data": data.reset_index().to_dict(orient="records"), "symbols": syms}
            return ToolResult(True, f"✓ Downloaded {len(data)} rows for {len(syms)} symbols", result)
        except Exception as e:
            return ToolResult(False, f"✗ get_multiple_stocks failed: {e}")

    @staticmethod
    def get_company_info(symbol: str) -> ToolResult:
        try:
            import yfinance as yf
            info = yf.Ticker(symbol.upper()).info
            keys = ["longName", "sector", "industry", "country", "fullTimeEmployees",
                    "marketCap", "trailingPE", "forwardPE", "dividendYield", "beta",
                    "52WeekHigh", "52WeekLow", "longBusinessSummary", "website"]
            filtered = {k: info.get(k) for k in keys if k in info}
            return ToolResult(True, f"✓ Company info for {symbol}", filtered)
        except Exception as e:
            return ToolResult(False, f"✗ get_company_info failed: {e}")

    @staticmethod
    def get_financial_statements(
        symbol: str,
        statement_type: str = "income",
    ) -> ToolResult:
        try:
            import yfinance as yf
            ticker = yf.Ticker(symbol.upper())
            if statement_type == "income":
                stmt = ticker.financials
            elif statement_type == "balance":
                stmt = ticker.balance_sheet
            elif statement_type == "cashflow":
                stmt = ticker.cashflow
            else:
                return ToolResult(False, f"✗ Unknown statement type: {statement_type}")
            if stmt is None or stmt.empty:
                return ToolResult(False, f"✗ No {statement_type} statement data for {symbol}")
            data = stmt.to_dict()
            data = {str(k): {str(k2): v2 for k2, v2 in v.items()} for k, v in data.items()}
            return ToolResult(True, f"✓ {statement_type} statement for {symbol}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_financial_statements failed: {e}")

    @staticmethod
    def get_earnings_calendar(date_from: str = "", date_to: str = "") -> ToolResult:
        try:
            import requests
            from bs4 import BeautifulSoup
            url = "https://finance.yahoo.com/calendar/earnings"
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
            soup = BeautifulSoup(resp.text, "lxml")
            rows = soup.select("table tbody tr")
            earnings = []
            for row in rows[:30]:
                cells = row.find_all("td")
                if cells:
                    earnings.append({"symbol": cells[0].get_text(strip=True) if cells else "", "company": cells[1].get_text(strip=True) if len(cells) > 1 else "", "date": cells[2].get_text(strip=True) if len(cells) > 2 else ""})
            return ToolResult(True, f"✓ {len(earnings)} earnings events", earnings)
        except Exception as e:
            return ToolResult(False, f"✗ get_earnings_calendar failed: {e}")

    @staticmethod
    def get_economic_indicators(indicator: str = "gdp", country: str = "US") -> ToolResult:
        try:
            import requests
            creds   = CredStore.load("worldbank")
            iso_map = {"US": "USA", "UK": "GBR", "DE": "DEU", "CN": "CHN", "JP": "JPN"}
            iso3    = iso_map.get(country, country)
            code_map = {"gdp": "NY.GDP.MKTP.CD", "inflation": "FP.CPI.TOTL.ZG", "unemployment": "SL.UEM.TOTL.ZS", "interest_rate": "FR.INR.RINR"}
            indicator_code = code_map.get(indicator.lower(), indicator)
            resp = requests.get(
                f"https://api.worldbank.org/v2/country/{iso3}/indicator/{indicator_code}?format=json&per_page=20",
                timeout=20,
            )
            data = resp.json()
            if len(data) < 2:
                return ToolResult(False, f"✗ No data for {indicator} in {country}")
            records = [{"year": r.get("date"), "value": r.get("value")} for r in data[1] if r.get("value") is not None]
            return ToolResult(True, f"✓ {len(records)} data points for {indicator} ({country})", records)
        except Exception as e:
            return ToolResult(False, f"✗ get_economic_indicators failed: {e}")

    @staticmethod
    def get_crypto_price(
        symbol: str = "BTC",
        currency: str = "USD",
        period: str = "30",
    ) -> ToolResult:
        try:
            import ccxt
            exchange = ccxt.binance()
            pair = f"{symbol.upper()}/{currency.upper()}"
            ohlcv = exchange.fetch_ohlcv(pair, timeframe="1d", limit=int(period))
            data = [{"timestamp": o[0], "open": o[1], "high": o[2], "low": o[3], "close": o[4], "volume": o[5]} for o in ohlcv]
            latest = data[-1] if data else {}
            return ToolResult(True, f"✓ {pair}: latest close = {latest.get('close', 'N/A')}", {"history": data, "latest": latest})
        except Exception as e:
            # fallback: CoinGecko free API
            try:
                import requests
                id_map = {"BTC": "bitcoin", "ETH": "ethereum", "SOL": "solana", "ADA": "cardano", "BNB": "binancecoin"}
                cg_id  = id_map.get(symbol.upper(), symbol.lower())
                resp   = requests.get(f"https://api.coingecko.com/api/v3/simple/price?ids={cg_id}&vs_currencies={currency.lower()}", timeout=15)
                data   = resp.json()
                price  = data.get(cg_id, {}).get(currency.lower(), "N/A")
                return ToolResult(True, f"✓ {symbol} = {price} {currency}", {"symbol": symbol, "currency": currency, "price": price})
            except Exception as e2:
                return ToolResult(False, f"✗ get_crypto_price failed: {e} / {e2}")

    @staticmethod
    def get_crypto_info(symbol: str) -> ToolResult:
        try:
            import requests
            id_map = {"BTC": "bitcoin", "ETH": "ethereum", "SOL": "solana", "ADA": "cardano", "BNB": "binancecoin", "XRP": "ripple", "DOT": "polkadot", "DOGE": "dogecoin"}
            cg_id  = id_map.get(symbol.upper(), symbol.lower())
            resp   = requests.get(f"https://api.coingecko.com/api/v3/coins/{cg_id}", timeout=20)
            if resp.status_code != 200:
                return ToolResult(False, f"✗ Crypto '{symbol}' not found")
            raw = resp.json()
            data = {
                "name":        raw.get("name"),
                "symbol":      raw.get("symbol"),
                "rank":        raw.get("market_cap_rank"),
                "price_usd":   raw.get("market_data", {}).get("current_price", {}).get("usd"),
                "market_cap":  raw.get("market_data", {}).get("market_cap", {}).get("usd"),
                "change_24h":  raw.get("market_data", {}).get("price_change_percentage_24h"),
                "description": raw.get("description", {}).get("en", "")[:300],
            }
            return ToolResult(True, f"✓ Crypto info: {symbol}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_crypto_info failed: {e}")

    @staticmethod
    def get_forex_rate(
        from_currency: str = "USD",
        to_currency: str = "EUR",
        period: str = "1mo",
    ) -> ToolResult:
        try:
            import yfinance as yf
            pair    = f"{from_currency.upper()}{to_currency.upper()}=X"
            ticker  = yf.Ticker(pair)
            hist    = ticker.history(period=period)
            if hist.empty:
                return ToolResult(False, f"✗ No forex data for {from_currency}/{to_currency}")
            latest_rate = float(hist["Close"].iloc[-1])
            data = {"pair": pair, "latest_rate": latest_rate, "history": hist["Close"].reset_index().to_dict(orient="records")}
            return ToolResult(True, f"✓ {from_currency}/{to_currency} = {latest_rate:.4f}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_forex_rate failed: {e}")

    @staticmethod
    def get_commodity_prices(commodity: str = "gold") -> ToolResult:
        try:
            import yfinance as yf
            sym_map = {"gold": "GC=F", "silver": "SI=F", "oil": "CL=F", "natural_gas": "NG=F", "corn": "ZC=F", "wheat": "ZW=F", "copper": "HG=F"}
            sym = sym_map.get(commodity.lower(), commodity)
            ticker = yf.Ticker(sym)
            hist   = ticker.history(period="1mo")
            if hist.empty:
                return ToolResult(False, f"✗ No data for commodity '{commodity}'")
            latest = float(hist["Close"].iloc[-1])
            return ToolResult(True, f"✓ {commodity}: {latest:.2f}", {"symbol": sym, "latest_price": latest, "history": hist["Close"].reset_index().to_dict(orient="records")})
        except Exception as e:
            return ToolResult(False, f"✗ get_commodity_prices failed: {e}")

    @staticmethod
    def calculate_technical_indicators(
        symbol: str,
        indicators: Optional[List[str]] = None,
        period: str = "3mo",
    ) -> ToolResult:
        try:
            import yfinance as yf
            import numpy as np
            ticker = yf.Ticker(symbol.upper())
            hist   = ticker.history(period=period)
            if hist.empty:
                return ToolResult(False, f"✗ No data for {symbol}")
            close  = hist["Close"]
            volume = hist["Volume"]
            inds   = indicators or ["sma20", "sma50", "rsi", "macd", "bollinger"]
            result: dict = {}
            if "sma20"  in inds: result["sma20"]  = float(close.rolling(20).mean().iloc[-1])
            if "sma50"  in inds: result["sma50"]  = float(close.rolling(50).mean().iloc[-1])
            if "ema12"  in inds: result["ema12"]  = float(close.ewm(span=12).mean().iloc[-1])
            if "ema26"  in inds: result["ema26"]  = float(close.ewm(span=26).mean().iloc[-1])
            if "rsi" in inds:
                delta  = close.diff()
                gain   = delta.where(delta > 0, 0).rolling(14).mean()
                loss   = (-delta.where(delta < 0, 0)).rolling(14).mean()
                rs     = gain / loss
                result["rsi"] = float((100 - 100 / (1 + rs)).iloc[-1])
            if "macd" in inds:
                ema12  = close.ewm(span=12).mean()
                ema26  = close.ewm(span=26).mean()
                macd   = ema12 - ema26
                signal = macd.ewm(span=9).mean()
                result["macd"]        = float(macd.iloc[-1])
                result["macd_signal"] = float(signal.iloc[-1])
                result["macd_hist"]   = float((macd - signal).iloc[-1])
            if "bollinger" in inds:
                sma = close.rolling(20).mean()
                std = close.rolling(20).std()
                result["bollinger_upper"] = float((sma + 2 * std).iloc[-1])
                result["bollinger_lower"] = float((sma - 2 * std).iloc[-1])
                result["bollinger_mid"]   = float(sma.iloc[-1])
            if "atr" in inds:
                high = hist["High"]; low = hist["Low"]
                tr   = (high - low).rolling(14).mean()
                result["atr"] = float(tr.iloc[-1])
            return ToolResult(True, f"✓ Technical indicators for {symbol}", result)
        except Exception as e:
            return ToolResult(False, f"✗ calculate_technical_indicators failed: {e}")

    @staticmethod
    def screen_stocks(
        market_cap_min: float = 1e9,
        pe_max: float = 30.0,
        dividend_min: float = 0.0,
        sector: str = "",
    ) -> ToolResult:
        try:
            import yfinance as yf
            # Use a predefined list of S&P 500 tickers for screening
            sp500_tickers = ["AAPL", "MSFT", "AMZN", "GOOGL", "META", "NVDA", "TSLA", "BRK-B",
                             "JNJ", "V", "PG", "JPM", "UNH", "HD", "MA", "DIS", "BAC", "INTC",
                             "VZ", "NFLX", "ADBE", "CRM", "PYPL", "PFE", "KO", "PEP", "XOM",
                             "CVX", "WMT", "MCD", "ABT", "T", "ABBV", "WFC", "CSCO", "AVGO"]
            matched = []
            for sym in sp500_tickers:
                try:
                    info = yf.Ticker(sym).info
                    mc   = info.get("marketCap", 0) or 0
                    pe   = info.get("trailingPE", 9999) or 9999
                    div  = info.get("dividendYield", 0) or 0
                    sec  = info.get("sector", "")
                    if mc >= market_cap_min and pe <= pe_max and div >= dividend_min:
                        if not sector or sector.lower() in sec.lower():
                            matched.append({"symbol": sym, "name": info.get("shortName"), "sector": sec, "market_cap": mc, "pe": pe, "dividend_yield": round(div * 100, 2)})
                except Exception:
                    pass
            return ToolResult(True, f"✓ {len(matched)} stocks match criteria", matched)
        except Exception as e:
            return ToolResult(False, f"✗ screen_stocks failed: {e}")

    @staticmethod
    def get_options_chain(symbol: str, expiry: str = "") -> ToolResult:
        try:
            import yfinance as yf
            ticker   = yf.Ticker(symbol.upper())
            expiries = ticker.options
            if not expiries:
                return ToolResult(False, f"✗ No options data for {symbol}")
            target = expiry if expiry in expiries else expiries[0]
            chain  = ticker.option_chain(target)
            data   = {
                "expiry":  target,
                "calls":   chain.calls[["strike", "lastPrice", "bid", "ask", "impliedVolatility", "openInterest"]].head(20).to_dict(orient="records"),
                "puts":    chain.puts[["strike",  "lastPrice", "bid", "ask", "impliedVolatility", "openInterest"]].head(20).to_dict(orient="records"),
            }
            return ToolResult(True, f"✓ Options chain for {symbol} (expiry: {target})", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_options_chain failed: {e}")

    @staticmethod
    def portfolio_analysis(
        holdings: List[Dict[str, Any]],
        output: str = "portfolio.html",
    ) -> ToolResult:
        """holdings: [{"symbol": "AAPL", "shares": 10, "avg_cost": 150.0}, ...]"""
        try:
            import yfinance as yf
            import pandas as pd
            import plotly.express as px
            portfolio = []
            total_value = 0.0
            for h in holdings:
                sym  = h["symbol"].upper()
                info = yf.Ticker(sym).fast_info
                current_price = getattr(info, "last_price", None) or 0
                shares        = h.get("shares", 0)
                avg_cost      = h.get("avg_cost", current_price)
                market_value  = shares * current_price
                gain_loss     = (current_price - avg_cost) * shares
                pct_return    = ((current_price - avg_cost) / avg_cost * 100) if avg_cost else 0
                total_value  += market_value
                portfolio.append({
                    "symbol":        sym,
                    "shares":        shares,
                    "avg_cost":      avg_cost,
                    "current_price": current_price,
                    "market_value":  round(market_value, 2),
                    "gain_loss":     round(gain_loss, 2),
                    "pct_return":    round(pct_return, 2),
                })
            df = pd.DataFrame(portfolio)
            fig = px.pie(df, values="market_value", names="symbol", title="Portfolio Allocation")
            fig.write_html(output)
            summary = {"total_value": round(total_value, 2), "total_gain_loss": round(sum(p["gain_loss"] for p in portfolio), 2), "positions": portfolio}
            return ToolResult(True, f"✓ Portfolio: ${total_value:,.2f} total value, chart saved to '{output}'", summary)
        except Exception as e:
            return ToolResult(False, f"✗ portfolio_analysis failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 6. SocialMediaDataTool
# ─────────────────────────────────────────────────────────────────────────────

class SocialMediaDataTool:
    name = "social_media_data"
    description = (
        "Social media intelligence: Twitter/X, Reddit, YouTube, Instagram profiles, "
        "HackerNews — fetch posts, comments, search, and user info."
    )
    use = (
        """Name of Tool:- SocialMediaDataTool

Purpose of Tool:-
The SocialMediaDataTool provides unified access to open-source and authenticated intelligence across primary social media networks, including Twitter/X, Reddit, YouTube, Instagram, and HackerNews. It abstracts the underlying API clients (`tweepy`, `praw`, Google's YouTube Data API, etc.) to fetch user account states, post timelines, recent search results, comments, and real-time geographic trends. It is designed to facilitate sentiment modeling, social trend monitoring, or brand reputation tracking without requiring separate boilerplate configurations for each platform.

Methods:-
- get_twitter_user: Fetches a target Twitter/X profile's basic descriptive field metadata, verified status, and core public engagement metrics.
- get_twitter_timeline: Retrieves a recent list of published tweets and corresponding interaction counters directly from an user's public feed.
- search_twitter: Queries the recent Twitter/X global index using explicit strings, language settings, and optional chronological boundaries.
- get_twitter_trends: Compiles geographic trending topics and approximate tweet transaction volumes based on platform-specific WOEID targets.
- get_reddit_posts: Structures specific subreddit post queues across distinct sorting patterns like hot, top, or new.
- search_reddit: Executes broad relevance scans across Reddit or restricted sub-forums matching input keyword strings.
- get_reddit_comments: Generates flat structural threads parsing the conversation history of an isolated Reddit text or URL thread.
- get_subreddit_info: Inspects aggregate community subscriber scales, active participant tracking, and content flags for a specific community.
- get_youtube_video_info: Resolves descriptive briefs, publication timestamps, duration details, and global engagement counts for a video ID.
- get_youtube_channel_info: Aggregates total view history, global subscriber counts, and video counts matching channel handles or IDs.
- get_youtube_comments: Extracts top-tier algorithmic or chronological feedback lists attached to a video.
- search_youtube: Executes algorithmic video search queries across the global platform index to isolate video assets matching search strings.
- get_instagram_profile: Leverages public endpoints and structured pattern fallbacks to extract user biographic summaries and core follower metrics.
- get_hackernews_top: Hits official database streams to collect active top-ranked stories, scores, and discussion sizes on HackerNews.

How to use Tool Methods:-

1. get_twitter_user:
   - Purpose: Collects an operational metadata profile snapshot for a specified Twitter/X username.
   - Arguments:
     a) username: str - Plain handle pointer for the destination account.
     b) cred_key: str (default: "twitter") - Credential dictionary alias key inside the runtime vault.
   - Returns: ToolResult holding structured profile creation logs and total public follower metric trees.
   - How to call: SocialMediaDataTool.get_twitter_user(username="nasa")

2. get_twitter_timeline:
   - Purpose: Reviews recent explicit user broadcasts and associated interaction metrics.
   - Arguments:
     a) username: str - Target profile account name string.
     b) max_results: int (default: 20) - Cap constraining maximum output array items returned.
     c) cred_key: str (default: "twitter") - Active client profile vault indicator.
   - Returns: ToolResult storing a list of distinct timeline elements.
   - How to call: SocialMediaDataTool.get_twitter_timeline(username="techcrunch", max_results=10)

3. search_twitter:
   - Purpose: Scans recent platform entries for specific topic matches or keywords.
   - Arguments:
     a) query: str - Expression string target for lookups.
     b) max_results: int (default: 20) - Structural query length restriction boundary.
     c) lang: str (default: "en") - ISO platform language constraint filter code.
     b) date_from: str (default: "") - UTC start boundary timestamp filter.
     e) cred_key: str (default: "twitter") - Vault storage credential link.
   - Returns: ToolResult passing matching historical content dictionaries.
   - How to call: SocialMediaDataTool.search_twitter(query="artificial intelligence", lang="en")

4. get_twitter_trends:
   - Purpose: Evaluates geographic regions to highlight hot tracking words and viral conversations.
   - Arguments:
     a) location: str (default: "worldwide") - Geographic regional label mapping to target codes (e.g., "worldwide", "us", "uk", "india").
     b) cred_key: str (default: "twitter") - Associated configuration token identification string.
   - Returns: ToolResult presenting ranked regional trends and calculated topic engagement volumes.
   - How to call: SocialMediaDataTool.get_twitter_trends(location="us")

5. get_reddit_posts:
   - Purpose: Collects core tracking details for standard forum post timelines.
   - Arguments:
     a) subreddit: str - target forum context target name.
     b) sort: str (default: "hot") - Platform queue arrangement option strings ("hot", "top", "new").
     c) limit: int (default: 25) - Quantitative index size ceiling parameter.
     d) time_filter: str (default: "week") - Chronological evaluation depth scope used for "top" sort queries ("day", "week", "month", "all").
     e) cred_key: str (default: "reddit") - Active identity keys location mapping reference.
   - Returns: ToolResult outputting titles, tracking links, upvote ratios, and score objects.
   - How to call: SocialMediaDataTool.get_reddit_posts(subreddit="Python", sort="top", time_filter="month")

6. search_reddit:
   - Purpose: Finds relevant text entries or URLs matching user strings.
   - Arguments:
     a) query: str - Explicit text snippet mapping parameters.
     b) subreddit: str (default: "all") - Boundary context parameter limiting matching logic.
     c) limit: int (default: 25) - Frame results limit ceiling value.
     d) cred_key: str (default: "reddit") - Access profile identification sequence.
   - Returns: ToolResult packing a list of structured entry documents.
   - How to call: SocialMediaDataTool.search_reddit(query="machine learning", subreddit="all")

7. get_reddit_comments:
   - Purpose: Flattens forum discussions to extract user feedback and replies.
   - Arguments:
     a) post_id: str - Explicit target forum submission reference ID string.
     b) limit: int (default: 50) - Numeric limit tracking depth scale for entries.
     c) sort: str (default: "top") - Direction structural sequence filter ("top", "new", "controversial").
     d) cred_key: str (default: "reddit") - Target environment key designation reference.
   - Returns: ToolResult storing separate comment blocks containing scores and message fragments.
   - How to call: SocialMediaDataTool.get_reddit_comments(post_id="123xyz", limit=20)

8. get_subreddit_info:
   - Purpose: Validates forum parameters, active user limits, and overall size.
   - Arguments:
     a) subreddit: str - Plain sub-forum identity indicator name.
     b) cred_key: str (default: "reddit") - Target configuration file mapping reference.
   - Returns: ToolResult reporting total subscriber weightings, overview descriptions, and explicit age flags.
   - How to call: SocialMediaDataTool.get_subreddit_info(subreddit="dataisbeautiful")

9. get_youtube_video_info:
   - Purpose: Looks up specific public video metrics and asset tracking attributes.
   - Arguments:
     a) video_id: str - target platform video parameter link character sequence.
     b) cred_key: str (default: "youtube") - Identity store target token lookup key.
   - Returns: ToolResult structuring channel names, view stats, descriptions, and duration trackers.
   - How to call: SocialMediaDataTool.get_youtube_video_info(video_id="dQw4w9WgXcQ")

10. get_youtube_channel_info:
    - Purpose: Gathers overview parameters, total asset counters, and scale profiles for creators.
    - Arguments:
      a) channel_id_or_name: str - target alphanumeric tracking code or handle token identifier.
      b) cred_key: str (default: "youtube") - Storage access key profile variable.
    - Returns: ToolResult showing user creation sizes, view scales, and profile paragraphs.
    - How to call: SocialMediaDataTool.get_youtube_channel_info(channel_id_or_name="UCBR8-60-B28hp2BmDPdntcQ")

11. get_youtube_comments:
    - Purpose: Collects public audience responses and discussion chains from videos.
    - Arguments:
      a) video_id: str - alphanumeric video marker.
      b) max_results: int (default: 50) - Collection list size ceiling restriction.
      c) cred_key: str (default: "youtube") - Client access validation block tag.
    - Returns: ToolResult listing comments with author names, likes, and timestamps.
    - How to call: SocialMediaDataTool.get_youtube_comments(video_id="dQw4w9WgXcQ", max_results=10)

12. search_youtube:
    - Purpose: Indexes online video assets to match query strings.
    - Arguments:
      a) query: str - Target scanning topic phrase context string.
      b) max_results: int (default: 10) - Allocation count limit parameter.
      c) cred_key: str (default: "youtube") - Key file marker string reference.
    - Returns: ToolResult supplying matching links, channel tags, and content snippets.
    - How to call: SocialMediaDataTool.search_youtube(query="python programming tutorials", max_results=5)

13. get_instagram_profile:
    - Purpose: Extracts biographic configurations and audience count snapshots from public profiles.
    - Arguments:
      a) username: str - Target profile screen handle.
    - Returns: ToolResult mapping total following statistics, follower tallies, and text bios.
    - How to call: SocialMediaDataTool.get_instagram_profile(username="natgeo")

14. get_hackernews_top:
    - Purpose: Monitors trending tech stories from the tech ecosystem.
    - Arguments:
      a) limit: int (default: 30) - Item indexing depth limit metric.
    - Returns: ToolResult packaging active technical stories, discussion tracking links, and submission scores.
    - How to call: SocialMediaDataTool.get_hackernews_top(limit=15)
    """)
    
    @staticmethod
    def _twitter_client(cred_key: str = "twitter"):
        import tweepy
        c = CredStore.load(cred_key)
        return tweepy.Client(
            bearer_token=c.get("bearer_token", ""),
            consumer_key=c.get("api_key", ""),
            consumer_secret=c.get("api_secret", ""),
            access_token=c.get("access_token", ""),
            access_token_secret=c.get("access_token_secret", ""),
            wait_on_rate_limit=True,
        )

    @staticmethod
    def _reddit_client(cred_key: str = "reddit"):
        import praw
        c = CredStore.load(cred_key)
        return praw.Reddit(
            client_id     = c.get("client_id", ""),
            client_secret = c.get("client_secret", ""),
            user_agent    = "NPMAgent/1.0",
        )

    @staticmethod
    def get_twitter_user(username: str, cred_key: str = "twitter") -> ToolResult:
        try:
            client = SocialMediaDataTool._twitter_client(cred_key)
            user   = client.get_user(username=username, user_fields=["public_metrics", "description", "location", "created_at", "verified"])
            if not user.data:
                return ToolResult(False, f"✗ Twitter user '{username}' not found")
            u = user.data
            data = {"id": str(u.id), "name": u.name, "username": u.username, "description": u.description, "metrics": u.public_metrics, "created_at": str(u.created_at)}
            return ToolResult(True, f"✓ Twitter user @{username}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_twitter_user failed: {e}")

    @staticmethod
    def get_twitter_timeline(username: str, max_results: int = 20, cred_key: str = "twitter") -> ToolResult:
        try:
            client = SocialMediaDataTool._twitter_client(cred_key)
            user   = client.get_user(username=username)
            if not user.data:
                return ToolResult(False, f"✗ User '{username}' not found")
            tweets = client.get_users_tweets(user.data.id, max_results=min(max_results, 100), tweet_fields=["created_at", "public_metrics", "text"])
            data   = [{"id": str(t.id), "text": t.text, "created_at": str(t.created_at), "metrics": t.public_metrics} for t in (tweets.data or [])]
            return ToolResult(True, f"✓ {len(data)} tweets from @{username}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_twitter_timeline failed: {e}")

    @staticmethod
    def search_twitter(
        query: str,
        max_results: int = 20,
        lang: str = "en",
        date_from: str = "",
        cred_key: str = "twitter",
    ) -> ToolResult:
        try:
            client = SocialMediaDataTool._twitter_client(cred_key)
            q = f"{query} lang:{lang}"
            kwargs: dict = {"query": q, "max_results": min(max_results, 100), "tweet_fields": ["created_at", "public_metrics", "author_id"]}
            if date_from: kwargs["start_time"] = date_from
            resp   = client.search_recent_tweets(**kwargs)
            tweets = [{"id": str(t.id), "text": t.text, "created_at": str(t.created_at), "metrics": t.public_metrics} for t in (resp.data or [])]
            return ToolResult(True, f"✓ {len(tweets)} tweets for '{query}'", tweets)
        except Exception as e:
            return ToolResult(False, f"✗ search_twitter failed: {e}")

    @staticmethod
    def get_twitter_trends(location: str = "worldwide", cred_key: str = "twitter") -> ToolResult:
        try:
            import tweepy
            c  = CredStore.load(cred_key)
            v1 = tweepy.API(tweepy.OAuth1UserHandler(c.get("api_key"), c.get("api_secret"), c.get("access_token"), c.get("access_token_secret")))
            woeid_map = {"worldwide": 1, "us": 23424977, "uk": 23424975, "india": 23424848, "canada": 23424775}
            woeid  = woeid_map.get(location.lower(), 1)
            trends = v1.get_place_trends(woeid)[0]["trends"]
            data   = [{"name": t["name"], "tweet_volume": t.get("tweet_volume"), "url": t.get("url")} for t in trends[:30]]
            return ToolResult(True, f"✓ {len(data)} trending topics in '{location}'", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_twitter_trends failed: {e}")

    @staticmethod
    def get_reddit_posts(
        subreddit: str,
        sort: str = "hot",
        limit: int = 25,
        time_filter: str = "week",
        cred_key: str = "reddit",
    ) -> ToolResult:
        try:
            r   = SocialMediaDataTool._reddit_client(cred_key)
            sub = r.subreddit(subreddit)
            fn  = getattr(sub, sort, sub.hot)
            posts = []
            for post in fn(limit=limit, time_filter=time_filter if sort == "top" else "all"):
                posts.append({"id": post.id, "title": post.title, "score": post.score, "upvote_ratio": post.upvote_ratio, "num_comments": post.num_comments, "url": post.url, "created_utc": post.created_utc, "author": str(post.author)})
            return ToolResult(True, f"✓ {len(posts)} posts from r/{subreddit}", posts)
        except Exception as e:
            return ToolResult(False, f"✗ get_reddit_posts failed: {e}")

    @staticmethod
    def search_reddit(query: str, subreddit: str = "all", limit: int = 25, cred_key: str = "reddit") -> ToolResult:
        try:
            r       = SocialMediaDataTool._reddit_client(cred_key)
            sub     = r.subreddit(subreddit)
            results = []
            for post in sub.search(query, limit=limit, sort="relevance"):
                results.append({"id": post.id, "title": post.title, "score": post.score, "subreddit": str(post.subreddit), "url": post.url, "num_comments": post.num_comments})
            return ToolResult(True, f"✓ {len(results)} Reddit results for '{query}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ search_reddit failed: {e}")

    @staticmethod
    def get_reddit_comments(post_id: str, limit: int = 50, sort: str = "top", cred_key: str = "reddit") -> ToolResult:
        try:
            r       = SocialMediaDataTool._reddit_client(cred_key)
            post    = r.submission(id=post_id)
            post.comment_sort = sort
            post.comments.replace_more(limit=0)
            comments = [{"id": c.id, "body": c.body[:500], "score": c.score, "author": str(c.author), "created_utc": c.created_utc} for c in post.comments.list()[:limit]]
            return ToolResult(True, f"✓ {len(comments)} comments for post '{post_id}'", comments)
        except Exception as e:
            return ToolResult(False, f"✗ get_reddit_comments failed: {e}")

    @staticmethod
    def get_subreddit_info(subreddit: str, cred_key: str = "reddit") -> ToolResult:
        try:
            r   = SocialMediaDataTool._reddit_client(cred_key)
            sub = r.subreddit(subreddit)
            data = {"name": sub.display_name, "title": sub.title, "description": sub.public_description[:500], "subscribers": sub.subscribers, "active": sub.accounts_active, "created_utc": sub.created_utc, "over_18": sub.over18, "url": f"https://www.reddit.com/r/{subreddit}"}
            return ToolResult(True, f"✓ Subreddit r/{subreddit}: {sub.subscribers:,} subscribers", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_subreddit_info failed: {e}")

    @staticmethod
    def get_youtube_video_info(video_id: str, cred_key: str = "youtube") -> ToolResult:
        try:
            import requests
            api_key = CredStore.load(cred_key).get("api_key", "")
            if not api_key:
                return ToolResult(False, "✗ No YouTube API key in credentials (key: 'youtube')")
            resp = requests.get(
                "https://www.googleapis.com/youtube/v3/videos",
                params={"part": "snippet,statistics,contentDetails", "id": video_id, "key": api_key},
                timeout=15,
            )
            items = resp.json().get("items", [])
            if not items:
                return ToolResult(False, f"✗ Video '{video_id}' not found")
            item = items[0]
            data = {"id": video_id, "title": item["snippet"]["title"], "channel": item["snippet"]["channelTitle"], "published": item["snippet"]["publishedAt"], "description": item["snippet"]["description"][:500], "stats": item.get("statistics", {}), "duration": item["contentDetails"].get("duration", "")}
            return ToolResult(True, f"✓ YouTube video: {data['title'][:60]}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_youtube_video_info failed: {e}")

    @staticmethod
    def get_youtube_channel_info(channel_id_or_name: str, cred_key: str = "youtube") -> ToolResult:
        try:
            import requests
            api_key = CredStore.load(cred_key).get("api_key", "")
            param_key = "id" if channel_id_or_name.startswith("UC") else "forUsername"
            resp = requests.get(
                "https://www.googleapis.com/youtube/v3/channels",
                params={"part": "snippet,statistics", param_key: channel_id_or_name, "key": api_key},
                timeout=15,
            )
            items = resp.json().get("items", [])
            if not items:
                return ToolResult(False, f"✗ Channel '{channel_id_or_name}' not found")
            item = items[0]
            data = {"id": item["id"], "title": item["snippet"]["title"], "description": item["snippet"]["description"][:300], "subscribers": item["statistics"].get("subscriberCount"), "views": item["statistics"].get("viewCount"), "videos": item["statistics"].get("videoCount")}
            return ToolResult(True, f"✓ YouTube channel: {data['title']}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_youtube_channel_info failed: {e}")

    @staticmethod
    def get_youtube_comments(video_id: str, max_results: int = 50, cred_key: str = "youtube") -> ToolResult:
        try:
            import requests
            api_key = CredStore.load(cred_key).get("api_key", "")
            resp = requests.get(
                "https://www.googleapis.com/youtube/v3/commentThreads",
                params={"part": "snippet", "videoId": video_id, "maxResults": min(max_results, 100), "order": "relevance", "key": api_key},
                timeout=20,
            )
            items    = resp.json().get("items", [])
            comments = [{"author": i["snippet"]["topLevelComment"]["snippet"]["authorDisplayName"], "text": i["snippet"]["topLevelComment"]["snippet"]["textDisplay"][:500], "likes": i["snippet"]["topLevelComment"]["snippet"]["likeCount"], "published": i["snippet"]["topLevelComment"]["snippet"]["publishedAt"]} for i in items]
            return ToolResult(True, f"✓ {len(comments)} YouTube comments", comments)
        except Exception as e:
            return ToolResult(False, f"✗ get_youtube_comments failed: {e}")

    @staticmethod
    def search_youtube(query: str, max_results: int = 10, cred_key: str = "youtube") -> ToolResult:
        try:
            import requests
            api_key = CredStore.load(cred_key).get("api_key", "")
            resp = requests.get(
                "https://www.googleapis.com/youtube/v3/search",
                params={"part": "snippet", "q": query, "maxResults": max_results, "type": "video", "key": api_key},
                timeout=20,
            )
            items   = resp.json().get("items", [])
            results = [{"video_id": i["id"]["videoId"], "title": i["snippet"]["title"], "channel": i["snippet"]["channelTitle"], "published": i["snippet"]["publishedAt"], "description": i["snippet"]["description"][:200]} for i in items]
            return ToolResult(True, f"✓ {len(results)} YouTube results for '{query}'", results)
        except Exception as e:
            return ToolResult(False, f"✗ search_youtube failed: {e}")

    @staticmethod
    def get_instagram_profile(username: str) -> ToolResult:
        try:
            import requests
            resp = requests.get(
                f"https://www.instagram.com/{username}/?__a=1&__d=dis",
                headers={"User-Agent": "Mozilla/5.0", "X-Requested-With": "XMLHttpRequest"},
                timeout=15,
            )
            if resp.status_code != 200:
                # fallback: scrape public page
                resp2 = requests.get(f"https://www.instagram.com/{username}/", headers={"User-Agent": "Mozilla/5.0"}, timeout=15)
                follower_match = re.search(r'"edge_followed_by":\{"count":(\d+)', resp2.text)
                following_match = re.search(r'"edge_follow":\{"count":(\d+)', resp2.text)
                posts_match = re.search(r'"edge_owner_to_timeline_media":\{"count":(\d+)', resp2.text)
                bio_match   = re.search(r'"biography":"([^"]*)"', resp2.text)
                return ToolResult(True, f"✓ Instagram @{username} (public)", {
                    "username":  username,
                    "followers": follower_match.group(1) if follower_match else "N/A",
                    "following": following_match.group(1) if following_match else "N/A",
                    "posts":     posts_match.group(1) if posts_match else "N/A",
                    "bio":       bio_match.group(1) if bio_match else "",
                })
            data = resp.json()
            user = data.get("graphql", {}).get("user", data.get("data", {}).get("user", {}))
            return ToolResult(True, f"✓ Instagram @{username}", {"username": username, "followers": user.get("edge_followed_by", {}).get("count"), "following": user.get("edge_follow", {}).get("count"), "posts": user.get("edge_owner_to_timeline_media", {}).get("count"), "bio": user.get("biography", "")})
        except Exception as e:
            return ToolResult(False, f"✗ get_instagram_profile failed: {e}")

    @staticmethod
    def get_hackernews_top(limit: int = 30) -> ToolResult:
        try:
            import requests
            top_ids = requests.get("https://hacker-news.firebaseio.com/v0/topstories.json", timeout=15).json()
            stories = []
            for story_id in top_ids[:limit]:
                item = requests.get(f"https://hacker-news.firebaseio.com/v0/item/{story_id}.json", timeout=10).json()
                if item and item.get("type") == "story":
                    stories.append({"id": story_id, "title": item.get("title"), "score": item.get("score"), "by": item.get("by"), "comments": item.get("descendants", 0), "url": item.get("url", f"https://news.ycombinator.com/item?id={story_id}")})
            return ToolResult(True, f"✓ Top {len(stories)} HackerNews stories", stories)
        except Exception as e:
            return ToolResult(False, f"✗ get_hackernews_top failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 7. WeatherGeoTool
# ─────────────────────────────────────────────────────────────────────────────

class WeatherGeoTool:
    name = "weather_geo"
    description = (
        "Weather and geolocation: current weather, forecasts, historical data, "
        "alerts, geocoding, reverse geocode, timezone, elevation, nearby places, air quality."
    )
    use = (
        """Name of Tool:- WeatherGeoTool

Purpose of Tool:-
The WeatherGeoTool delivers a comprehensive global environment intelligence interface. It unifies public geospatial endpoints with major mapping services (including OpenWeatherMap, Google Maps, and Open-Elevation) into a clean programmatic structure. The suite handles structural geocoding conversions, parses atmospheric and metric variables, tracks chronological weather patterns, calculates physical road matrices, maps global elevations, and updates dynamic regional health variables like real-time UV indices or particulate air pollution components.

Methods:-
- get_current_weather: Pulls live metric or imperial meteorological readings, cloud categories, and wind speed parameters for an exact location string.
- get_forecast: Extracts multi-day segmented atmospheric projection sequences grouped into detailed 3-hour structural candle increments.
- get_historical_weather: References specialized atmospheric lookups to isolate complete climate snapshots for specific past dates.
- get_weather_alerts: Polls international disaster prevention networks for active meteorological hazard warnings or life-safety notices.
- geocode: Maps literal textual addresses directly to standard decimal coordinate pairs (latitude and longitude).
- reverse_geocode: Translates raw terrestrial decimal coordinate matrices into formatted addresses.
- get_timezone: Determines regional identity tags and systemic UTC temporal offset data for exact geometric points.
- calculate_distance: Resolves multi-destination matrices calculating exact road miles and estimated driving transit times from a primary address.
- get_elevation: Intersects target coordinates against topographic terrain maps to return vertical height meters relative to sea level.
- get_nearby_places: Scans precise structural radial perimeters around user locations to index local businesses or restaurants by category type.
- get_air_quality: Samples live particulate records to map an Air Quality Index (AQI) alongside distinct carbon, ozone, and sulfur micro-metrics.
- get_uv_index: Gathers immediate solar radiation values to grade relative UV exposure hazards across standardized warning scales.

How to use Tool Methods:-

1. get_current_weather:
   - Purpose: Collects active ground-level atmospheric updates for a target region.
   - Arguments:
     a) location: str - Town, state, or global city destination text string.
     b) units: str (default: "metric") - Metric measurement standard choice ("metric" for Celsius, "imperial" for Fahrenheit).
     c) cred_key: str (default: "openweather") - Vault indicator configuration mapping key.
   - Returns: ToolResult passing exact temperature readings, relative humidity, wind metrics, and current weather text descriptions.
   - How to call: WeatherGeoTool.get_current_weather(location="London, UK", units="metric")

2. get_forecast:
   - Purpose: Monitors intermediate atmospheric transformations to determine approaching weather trends.
   - Arguments:
     a) location: str - target alphanumeric location designation.
     b) days: int (default: 5) - Forward chronological day scale ceiling length constraint.
     c) units: str (default: "metric") - Internal conversion scalar formatting framework.
     d) cred_key: str (default: "openweather") - Active credentials link location.
   - Returns: ToolResult packaging consecutive segmented weather prediction frames.
   - How to call: WeatherGeoTool.get_forecast(location="Tokyo", days=3, units="metric")

3. get_historical_weather:
   - Purpose: Reconstructs exact past weather criteria to assist tracking calculations or climate studies.
   - Arguments:
     a) location: str - Target environment description.
     b) date: str - Historical calendar checkpoint formatted explicitly as "YYYY-MM-DD".
     c) units: str (default: "metric") - Output standard variable format.
     d) cred_key: str (default: "openweather") - Configuration secret validation reference pointer.
   - Returns: ToolResult delivering the chosen date's climate properties.
   - How to call: WeatherGeoTool.get_historical_weather(location="New York", date="2025-07-04")

4. get_weather_alerts:
   - Purpose: Captures high-severity local safety warnings or weather events.
   - Arguments:
     a) location: str - Alphanumeric community identity label string.
     b) cred_key: str (default: "openweather") - Internal API authentication identifier map string.
   - Returns: ToolResult supplying critical emergency broadcast alert data blocks.
   - How to call: WeatherGeoTool.get_weather_alerts(location="Miami, FL")

5. geocode:
   - Purpose: Translates descriptive text strings into specific mathematical terrestrial parameters.
   - Arguments:
     a) address: str - Street layout coordinates, city names, or region labels.
     b) cred_key: str (default: "googlemaps") - Associated routing data license locator string.
   - Returns: ToolResult yielding decimal values for coordinate parameters along with formatted labels.
   - How to call: WeatherGeoTool.geocode(address="1600 Amphitheatre Pkwy, Mountain View, CA")

6. reverse_geocode:
   - Purpose: Decodes complex geographical points back into traditional mailing addresses.
   - Arguments:
     a) lat: float - Terrestrial horizontal meridian position point identifier.
     b) lng: float - Terrestrial vertical meridian position point identifier.
     c) cred_key: str (default: "googlemaps") - Key profile variable marker.
   - Returns: ToolResult with formatted geographic attributes and structural address details.
   - How to call: WeatherGeoTool.reverse_geocode(lat=40.7128, lng=-74.0060)

7. get_timezone:
   - Purpose: Resolves local timezone rules and daylight saving offsets for any coordinate pair.
   - Arguments:
     a) lat: float - target horizontal location parameter.
     b) lng: float - target vertical location parameter.
     c) cred_key: str (default: "googlemaps") - API validation authorization pointer variable.
   - Returns: ToolResult storing standard IANA timezone strings and current offset structures.
   - How to call: WeatherGeoTool.get_timezone(lat=35.6762, lng=139.6503)

8. calculate_distance:
   - Purpose: Computes dynamic travel metrics spanning various routes and destinations.
   - Arguments:
     a) from_location: str - True origins departure baseline address context string.
     b) to_locations: List[str] - Destination coordinate array index.
     c) mode: str (default: "driving") - Transit behavior model ("driving", "walking", "bicycling", "transit").
     d) cred_key: str (default: "googlemaps") - Google engine verification tag mapping reference.
   - Returns: ToolResult outputting absolute distances alongside logical travel durations.
   - How to call: WeatherGeoTool.calculate_distance(from_location="Paris", to_locations=["Lyon", "Marseille"], mode="driving")

9. get_elevation:
   - Purpose: Inspects raw physical vertical heights relative to international base sea measurements.
   - Arguments:
     a) lat: float - Latitude parameter element value.
     b) lng: float - Longitude parameter element value.
     c) cred_key: str (default: "openweather") - Internal API vault storage key mapping text.
   - Returns: ToolResult mapping altitude properties in meters.
   - How to call: WeatherGeoTool.get_elevation(lat=46.5588, lng=10.4368)

10. get_nearby_places:
    - Purpose: Locates local points of interest inside specified spatial ranges.
    - Arguments:
      a) lat: float - Baseline anchor horizontal tracking axis coordinate value.
      b) lng: float - Baseline anchor vertical tracking axis coordinate value.
      c) type: str (default: "restaurant") - Specific structural classification marker category (e.g., "hotel", "atm", "gas_station").
      d) radius: int (default: 1000) - Circular constraint scanning distance parameter in meters.
      e) cred_key: str (default: "googlemaps") - Platform engine tracking keys mapping path string.
    - Returns: ToolResult summarizing local destination names, addresses, and user rating metrics.
    - How to call: WeatherGeoTool.get_nearby_places(lat=48.8566, lng=2.3522, type="cafe", radius=500)

11. get_air_quality:
    - Purpose: Evaluates environmental purity variables and structural pollution components.
    - Arguments:
      a) location: str - target structural community address query context.
      b) cred_key: str (default: "openweather") - Active verification token path selector string.
    - Returns: ToolResult evaluating AQI threat ranges alongside exact structural concentration maps.
    - How to call: WeatherGeoTool.get_air_quality(location="Beijing")

12. get_uv_index:
    - Purpose: Gauges potential ultraviolet burn hazards to manage outdoor risks.
    - Arguments:
      a) location: str - Targeted regional tracking location indicator name.
      b) cred_key: str (default: "openweather") - Validation credential configuration profile tag.
    - Returns: ToolResult grading absolute radiation levels across qualitative risk brackets.
    - How to call: WeatherGeoTool.get_uv_index(location="Sydney, Australia")
    """)
    
    @staticmethod
    def _owm_key(cred_key: str = "openweather") -> str:
        return CredStore.load(cred_key).get("api_key", "")

    @staticmethod
    def _geocode_coords(location: str, cred_key: str = "openweather") -> Tuple[float, float]:
        import requests
        api_key = WeatherGeoTool._owm_key(cred_key)
        resp    = requests.get(
            "https://api.openweathermap.org/geo/1.0/direct",
            params={"q": location, "limit": 1, "appid": api_key},
            timeout=10,
        )
        data = resp.json()
        if not data:
            raise ValueError(f"Location not found: {location}")
        return data[0]["lat"], data[0]["lon"]

    @staticmethod
    def get_current_weather(location: str, units: str = "metric", cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key = WeatherGeoTool._owm_key(cred_key)
            if not api_key:
                return ToolResult(False, "✗ No OpenWeatherMap API key (cred_key: 'openweather')")
            resp = requests.get(
                "https://api.openweathermap.org/data/2.5/weather",
                params={"q": location, "units": units, "appid": api_key},
                timeout=15,
            )
            w = resp.json()
            if resp.status_code != 200:
                return ToolResult(False, f"✗ OWM error: {w.get('message')}")
            data = {
                "location":    w["name"],
                "country":     w["sys"]["country"],
                "temp":        w["main"]["temp"],
                "feels_like":  w["main"]["feels_like"],
                "humidity":    w["main"]["humidity"],
                "description": w["weather"][0]["description"],
                "wind_speed":  w["wind"]["speed"],
                "visibility":  w.get("visibility", 0),
                "units":       units,
            }
            unit_sym = "°C" if units == "metric" else "°F"
            return ToolResult(True, f"✓ {data['location']}: {data['temp']}{unit_sym}, {data['description']}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_current_weather failed: {e}")

    @staticmethod
    def get_forecast(location: str, days: int = 5, units: str = "metric", cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key = WeatherGeoTool._owm_key(cred_key)
            resp    = requests.get(
                "https://api.openweathermap.org/data/2.5/forecast",
                params={"q": location, "units": units, "cnt": days * 8, "appid": api_key},
                timeout=15,
            )
            data = resp.json()
            if resp.status_code != 200:
                return ToolResult(False, f"✗ OWM error: {data.get('message')}")
            forecasts = [
                {"datetime": f["dt_txt"], "temp": f["main"]["temp"], "feels_like": f["main"]["feels_like"], "description": f["weather"][0]["description"], "humidity": f["main"]["humidity"], "wind": f["wind"]["speed"]}
                for f in data.get("list", [])
            ]
            return ToolResult(True, f"✓ {days}-day forecast for '{location}': {len(forecasts)} data points", forecasts)
        except Exception as e:
            return ToolResult(False, f"✗ get_forecast failed: {e}")

    @staticmethod
    def get_historical_weather(location: str, date: str, units: str = "metric", cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key  = WeatherGeoTool._owm_key(cred_key)
            lat, lon = WeatherGeoTool._geocode_coords(location, cred_key)
            dt       = int(datetime.strptime(date, "%Y-%m-%d").timestamp())
            resp     = requests.get(
                "https://api.openweathermap.org/data/3.0/onecall/timemachine",
                params={"lat": lat, "lon": lon, "dt": dt, "units": units, "appid": api_key},
                timeout=15,
            )
            data = resp.json()
            if resp.status_code != 200:
                return ToolResult(False, f"✗ OWM error: {data.get('message')}")
            return ToolResult(True, f"✓ Historical weather for '{location}' on {date}", data.get("data", [{}])[0])
        except Exception as e:
            return ToolResult(False, f"✗ get_historical_weather failed: {e}")

    @staticmethod
    def get_weather_alerts(location: str, cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key  = WeatherGeoTool._owm_key(cred_key)
            lat, lon = WeatherGeoTool._geocode_coords(location, cred_key)
            resp     = requests.get(
                "https://api.openweathermap.org/data/3.0/onecall",
                params={"lat": lat, "lon": lon, "exclude": "current,minutely,hourly,daily", "appid": api_key},
                timeout=15,
            )
            alerts = resp.json().get("alerts", [])
            return ToolResult(True, f"✓ {len(alerts)} weather alerts for '{location}'", alerts)
        except Exception as e:
            return ToolResult(False, f"✗ get_weather_alerts failed: {e}")

    @staticmethod
    def geocode(address: str, cred_key: str = "googlemaps") -> ToolResult:
        try:
            api_key = CredStore.load(cred_key).get("api_key", "")
            if api_key:
                import googlemaps
                gmaps  = googlemaps.Client(key=api_key)
                result = gmaps.geocode(address)
                if not result:
                    return ToolResult(False, f"✗ Address not found: {address}")
                loc = result[0]["geometry"]["location"]
                return ToolResult(True, f"✓ Geocoded: {result[0]['formatted_address']}", {"lat": loc["lat"], "lng": loc["lng"], "formatted_address": result[0]["formatted_address"]})
            # fallback: Nominatim
            import requests
            resp = requests.get("https://nominatim.openstreetmap.org/search", params={"q": address, "format": "json", "limit": 1}, headers={"User-Agent": "NPMAgent/1.0"}, timeout=15)
            data = resp.json()
            if not data:
                return ToolResult(False, f"✗ Address not found: {address}")
            return ToolResult(True, f"✓ Geocoded: {data[0]['display_name']}", {"lat": float(data[0]["lat"]), "lng": float(data[0]["lon"]), "formatted_address": data[0]["display_name"]})
        except Exception as e:
            return ToolResult(False, f"✗ geocode failed: {e}")

    @staticmethod
    def reverse_geocode(lat: float, lng: float, cred_key: str = "googlemaps") -> ToolResult:
        try:
            api_key = CredStore.load(cred_key).get("api_key", "")
            if api_key:
                import googlemaps
                gmaps  = googlemaps.Client(key=api_key)
                result = gmaps.reverse_geocode((lat, lng))
                if not result:
                    return ToolResult(False, f"✗ No address found for ({lat}, {lng})")
                return ToolResult(True, f"✓ {result[0]['formatted_address']}", result[0])
            import requests
            resp = requests.get("https://nominatim.openstreetmap.org/reverse", params={"lat": lat, "lon": lng, "format": "json"}, headers={"User-Agent": "NPMAgent/1.0"}, timeout=15)
            data = resp.json()
            return ToolResult(True, f"✓ {data.get('display_name', '')}", data)
        except Exception as e:
            return ToolResult(False, f"✗ reverse_geocode failed: {e}")

    @staticmethod
    def get_timezone(lat: float, lng: float, cred_key: str = "googlemaps") -> ToolResult:
        try:
            import requests
            api_key = CredStore.load(cred_key).get("api_key", "")
            if api_key:
                resp = requests.get(
                    "https://maps.googleapis.com/maps/api/timezone/json",
                    params={"location": f"{lat},{lng}", "timestamp": int(time.time()), "key": api_key},
                    timeout=15,
                )
                data = resp.json()
                return ToolResult(data.get("status") == "OK", f"✓ Timezone: {data.get('timeZoneId')}", data)
            # fallback: timeapi.io
            resp = requests.get(f"https://timeapi.io/api/TimeZone/coordinate?latitude={lat}&longitude={lng}", timeout=15)
            data = resp.json()
            return ToolResult(True, f"✓ Timezone: {data.get('timeZone')}", data)
        except Exception as e:
            return ToolResult(False, f"✗ get_timezone failed: {e}")

    @staticmethod
    def calculate_distance(
        from_location: str,
        to_locations: List[str],
        mode: str = "driving",
        cred_key: str = "googlemaps",
    ) -> ToolResult:
        try:
            import googlemaps
            api_key = CredStore.load(cred_key).get("api_key", "")
            if not api_key:
                return ToolResult(False, "✗ Google Maps API key required for distance calculations")
            gmaps   = googlemaps.Client(key=api_key)
            matrix  = gmaps.distance_matrix(from_location, to_locations, mode=mode)
            results = []
            for i, dest in enumerate(to_locations):
                elem = matrix["rows"][0]["elements"][i]
                results.append({"destination": dest, "distance": elem.get("distance", {}).get("text", ""), "duration": elem.get("duration", {}).get("text", ""), "status": elem.get("status")})
            return ToolResult(True, f"✓ Distance matrix for {len(to_locations)} destinations", results)
        except Exception as e:
            return ToolResult(False, f"✗ calculate_distance failed: {e}")

    @staticmethod
    def get_elevation(lat: float, lng: float, cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            resp = requests.get(
                f"https://api.open-elevation.com/api/v1/lookup",
                params={"locations": f"{lat},{lng}"},
                timeout=15,
            )
            data      = resp.json()
            elevation = data.get("results", [{}])[0].get("elevation", 0)
            return ToolResult(True, f"✓ Elevation at ({lat}, {lng}): {elevation}m", {"lat": lat, "lng": lng, "elevation_m": elevation})
        except Exception as e:
            return ToolResult(False, f"✗ get_elevation failed: {e}")

    @staticmethod
    def get_nearby_places(
        lat: float,
        lng: float,
        type: str = "restaurant",
        radius: int = 1000,
        cred_key: str = "googlemaps",
    ) -> ToolResult:
        try:
            import googlemaps
            api_key = CredStore.load(cred_key).get("api_key", "")
            if not api_key:
                return ToolResult(False, "✗ Google Maps API key required")
            gmaps  = googlemaps.Client(key=api_key)
            result = gmaps.places_nearby(location=(lat, lng), radius=radius, type=type)
            places = [{"name": p["name"], "rating": p.get("rating"), "address": p.get("vicinity"), "place_id": p["place_id"]} for p in result.get("results", [])[:20]]
            return ToolResult(True, f"✓ {len(places)} {type}s within {radius}m", places)
        except Exception as e:
            return ToolResult(False, f"✗ get_nearby_places failed: {e}")

    @staticmethod
    def get_air_quality(location: str, cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key  = WeatherGeoTool._owm_key(cred_key)
            lat, lon = WeatherGeoTool._geocode_coords(location, cred_key)
            resp     = requests.get(
                "https://api.openweathermap.org/data/2.5/air_pollution",
                params={"lat": lat, "lon": lon, "appid": api_key},
                timeout=15,
            )
            data = resp.json()
            aqi_map = {1: "Good", 2: "Fair", 3: "Moderate", 4: "Poor", 5: "Very Poor"}
            if "list" not in data:
                return ToolResult(False, f"✗ No air quality data: {data}")
            current = data["list"][0]
            aqi     = current["main"]["aqi"]
            components = current["components"]
            return ToolResult(True, f"✓ Air quality in '{location}': {aqi_map.get(aqi, aqi)}", {"aqi": aqi, "status": aqi_map.get(aqi), "components": components})
        except Exception as e:
            return ToolResult(False, f"✗ get_air_quality failed: {e}")

    @staticmethod
    def get_uv_index(location: str, cred_key: str = "openweather") -> ToolResult:
        try:
            import requests
            api_key  = WeatherGeoTool._owm_key(cred_key)
            lat, lon = WeatherGeoTool._geocode_coords(location, cred_key)
            resp     = requests.get(
                "https://api.openweathermap.org/data/3.0/onecall",
                params={"lat": lat, "lon": lon, "exclude": "minutely,hourly,daily,alerts", "appid": api_key},
                timeout=15,
            )
            data = resp.json()
            uv   = data.get("current", {}).get("uvi", 0)
            risk_map = [(3, "Low"), (6, "Moderate"), (8, "High"), (11, "Very High")]
            risk = "Extreme"
            for threshold, label in risk_map:
                if uv < threshold:
                    risk = label
                    break
            return ToolResult(True, f"✓ UV Index in '{location}': {uv} ({risk})", {"uv_index": uv, "risk": risk, "location": location})
        except Exception as e:
            return ToolResult(False, f"✗ get_uv_index failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 8. TextAnalyticsTool
# ─────────────────────────────────────────────────────────────────────────────

class TextAnalyticsTool:
    name = "text_analytics"
    description = (
        "NLP and text intelligence: sentiment, classification, entity extraction, "
        "keywords, summarization, translation, language detection, grammar, readability, "
        "topic modeling, similarity, embeddings, semantic search, plagiarism detection."
    )
    use = (
        """Name of Tool:- TextAnalyticsTool

Purpose of Tool:-
The TextAnalyticsTool acts as an all-in-one Natural Language Processing (NLP) framework. It bundles lightweight heuristic models, statistical processors, and machine learning pipelines (such as `TextBlob`, `scikit-learn`, `Helsinki-NLP`, and local LLMs via `Ollama`). The tool handles textual feature discovery, semantic indexing, cross-lingual translation, structural grammar profiling, readability math verification, and multidimensional semantic vector alignment.

Methods:-
- sentiment_analysis: Scores explicit text strings on polar scales to classify the underlying emotional sentiment context.
- classify_text: Assigns zero-shot classifications across user-provided target labels.
- extract_entities: Groups structural strings into classified token types like proper names, organizations, dates, or emails.
- extract_keywords: Scores distinctive text tokens using TF-IDF matrices, noun phrase parsing, or literal structural word frequencies.
- summarize: Condenses long documents into abstractive concepts or extractive sentences.
- translate: Converts source strings into matching destination target language arrays using transformers or generative frameworks.
- detect_language: Analyzes input strings to generate a prioritized list of predicted structural language matches.
- check_grammar: Runs language syntax checkers to detect errors and export fully corrected strings.
- calculate_readability: Solves traditional linguistic mathematical formulas to return Flesch Reading Ease and equivalent school grade rankings.
- topic_modeling: Deconstructs text arrays into clear semantic thematic word groups using unsupervised LDA or NMF algorithms.
- text_similarity: Calculates literal structural or deeper semantic document similarities via Edit-Distance or Vector Space matrices.
- generate_embeddings: Converts plain text into dense dimensional vector embeddings for custom vector tasks.
- semantic_search: Scores query vector positions against document corpora to pull matching contextual fragments.
- detect_plagiarism: Audits test texts against an array of reference documents using token overlap thresholds.

How to use Tool Methods:-

1. sentiment_analysis:
   - Purpose: Examines contextual tone across one or more texts to evaluate emotional leanings.
   - Arguments:
     a) texts: Union[str, List[str]] - A single string block or a list of target text sentences.
     b) model: str (default: "textblob") - Operational pipeline engine logic selector ("textblob" or "ollama").
   - Returns: ToolResult passing polarity scores, subjective weightings, and qualitative labels.
   - How to call: TextAnalyticsTool.sentiment_analysis(texts="I love using this new software tool!", model="textblob")

2. classify_text:
   - Purpose: Flags unstructured strings using predefined descriptive evaluation labels.
   - Arguments:
     a) text: str - Source text content snippet.
     b) labels: List[str] - Target classification labels.
     c) model: str (default: "transformers") - Core algorithm runtime vehicle ("transformers" or "ollama").
   - Returns: ToolResult outputting the most relevant labels alongside relative model confidence metrics.
   - How to call: TextAnalyticsTool.classify_text(text="The stock market crashed today.", labels=["politics", "finance", "sports"])

3. extract_entities:
   - Purpose: Traverses phrases to detect and categorize real-world entities.
   - Arguments:
     a) text: str - Document body being mapped.
     b) entity_types: Optional[List[str]] (default: None) - Constraints limiting output to specific classes (e.g., ["PER", "LOC"]).
   - Returns: ToolResult storing structured arrays of tokens, structural classifications, and pipeline scores.
   - How to call: TextAnalyticsTool.extract_entities(text="Elon Musk visited the Tesla factory in Berlin.")

4. extract_keywords:
   - Purpose: Extracts the most informative and distinct words or phrases from a text block.
   - Arguments:
     a) text: str - Main raw document target text string.
     b) n: int (default: 10) - Total allocation count limit for final keyword list.
     c) method: str (default: "tfidf") - Statistical model selector strategy ("tfidf", "textrank", "frequency").
   - Returns: ToolResult mapping relevant keyword elements along with their calculated relevance indicators.
   - How to call: TextAnalyticsTool.extract_keywords(text="Data science relies on statistics and machine learning.", n=5)

5. summarize:
   - Purpose: Shrinks documents down to their key sentences or ideas.
   - Arguments:
     a) text: str - Target body of text.
     b) length: str (default: "medium") - Intended footprint restriction scale ("short", "medium", "long").
     c) style: str (default: "abstractive") - Compression structural strategy selection ("abstractive" or "extractive").
     d) model: str (default: "ollama") - Engine execution driver identity string ("ollama" or "transformers").
   - Returns: ToolResult featuring a condensed summary string.
   - How to call: TextAnalyticsTool.summarize(text="Very long article text here...", length="short", style="abstractive")

6. translate:
   - Purpose: Rewrites text from a source language into a different target language.
   - Arguments:
     a) text: str - Source text block meant for translation.
     b) target_language: str (default: "es") - Destination language code string abbreviation (e.g., "es", "fr", "hi").
     c) model: str (default: "ollama") - Runtime pipeline system platform reference ("ollama" or "transformers").
   - Returns: ToolResult yielding the translated text block.
   - How to call: TextAnalyticsTool.translate(text="Hello, how are you today?", target_language="fr")

7. detect_language:
   - Purpose: Evaluates phonetic and syntax structural indicators to identify the source language.
   - Arguments:
     a) text: str - Plain language string query block.
   - Returns: ToolResult logging top predicted ISO codes coupled with clear statistical probability scales.
   - How to call: TextAnalyticsTool.detect_language(text="Bonjour tout le monde")

8. check_grammar:
   - Purpose: Scans documents to identify spelling issues, typing slipups, and grammatical errors.
   - Arguments:
     a) text: str - Input copy awaiting processing evaluation.
   - Returns: ToolResult returning list maps highlighting exact context offsets, error details, and updated replacement text string variants.
   - How to call: TextAnalyticsTool.check_grammar(text="He do not know where to go.")

9. calculate_readability:
   - Purpose: Grades text difficulty to measure how easy it is for humans to comprehend.
   - Arguments:
     a) text: str - Target document copy.
   - Returns: ToolResult exporting structural word variables alongside Flesch and SMOG grade metric scores.
   - How to call: TextAnalyticsTool.calculate_readability(text="The quick brown fox jumps over the lazy dog.")

10. topic_modeling:
    - Purpose: Discovers hidden thematic structures within a collection of text documents.
    - Arguments:
      a) texts: List[str] - Collection index holding separate text documents.
      b) n_topics: int (default: 5) - Target allocation count for cluster extractions.
      b) method: str (default: "lda") - Matrix processing algorithm strategy ("lda" or "nmf").
      d) output: str (default: "topics.json") - Path destination file name.
    - Returns: ToolResult presenting key keyword matrix bundles sorted into clean topic groups.
    - How to call: TextAnalyticsTool.topic_modeling(texts=["doc one content", "doc two content"], n_topics=3)

11. text_similarity:
    - Purpose: Measures the closeness between two texts using either raw characters or conceptual meaning.
    - Arguments:
      a) text1: str - Baseline comparison source string.
      b) text2: str - Destination tracking target comparison text.
      c) method: str (default: "cosine") - Algorithmic tracking engine selection ("cosine", "embedding", "levenshtein").
    - Returns: ToolResult supplying explicit mathematical similarity values scaled from 0 to 1.
    - How to call: TextAnalyticsTool.text_similarity(text1="Machine Learning", text2="Deep Learning", method="embedding")

12. generate_embeddings:
    - Purpose: Encodes texts into high-dimensional vector representations for mathematical distance modeling.
    - Arguments:
      a) texts: Union[str, List[str]] - Singular string sequence or list array containing target strings.
      b) model: str (default: "all-MiniLM-L6-v2") - SentenceTransformer vector mapping identity reference.
    - Returns: ToolResult containing multi-dimensional vector arrays alongside dimension structural shapes.
    - How to call: TextAnalyticsTool.generate_embeddings(texts=["Deep learning is a subset of AI."])

13. semantic_search:
    - Purpose: Finds relevant matches within a text database based on the core meaning of a search query rather than exact keywords.
    - Arguments:
      a) query: str - Plain language search statement context query.
      b) corpus: List[str] - Data baseline array indexing available document references.
      c) top_k: int (default: 5) - Maximum length limit defining returned results array.
    - Returns: ToolResult delivering ranked documents paired with their structural cosine value proximity records.
    - How to call: TextAnalyticsTool.semantic_search(query="fruits", corpus=["An apple a day", "Cars drive fast", "Bananas are yellow"])

14. detect_plagiarism:
    - Purpose: Compares a test string against known resource text arrays to isolate copycat text issues.
    - Arguments:
      a) text: str - Document undergoing audit verification checks.
      b) reference_texts: List[str] - Known source materials checklist array.
      c) threshold: float (default: 0.7) - Evaluation flag ceiling constraint index (0.0 to 1.0).
    - Returns: ToolResult reporting individual match scores and active plagiarism alert flags.
    - How to call: TextAnalyticsTool.detect_plagiarism(text="Sample check essay text", reference_texts=["Original reference source one"])
    """)
        
    @staticmethod
    def sentiment_analysis(texts: Union[str, List[str]], model: str = "textblob") -> ToolResult:
        try:
            if isinstance(texts, str): texts = [texts]
            results = []
            if model == "textblob":
                from textblob import TextBlob
                for t in texts:
                    blob  = TextBlob(t)
                    pol   = blob.sentiment.polarity
                    label = "positive" if pol > 0.05 else "negative" if pol < -0.05 else "neutral"
                    results.append({"text": t[:100], "polarity": round(pol, 3), "subjectivity": round(blob.sentiment.subjectivity, 3), "label": label})
            elif model == "ollama":
                try:
                    from npmai import Ollama
                    llm = Ollama(model="llama3.2:3b", temperature=0.1, change=True, Models=["mistral:7b"])
                    for t in texts:
                        resp  = llm.invoke(f"Classify sentiment as positive/negative/neutral and give a score -1 to 1. Reply JSON only: {{\"label\":\"...\",\"score\":0.0}}. Text: {t[:500]}")
                        clean = re.search(r'\{.*?\}', resp, re.DOTALL)
                        data  = json.loads(clean.group()) if clean else {"label": "neutral", "score": 0.0}
                        results.append({"text": t[:100], **data})
                except ImportError:
                    return ToolResult(False, "✗ npmai not available for ollama sentiment")
            return ToolResult(True, f"✓ Sentiment analyzed for {len(results)} texts", results)
        except Exception as e:
            return ToolResult(False, f"✗ sentiment_analysis failed: {e}")

    @staticmethod
    def classify_text(text: str, labels: List[str], model: str = "transformers") -> ToolResult:
        try:
            if model == "transformers":
                from transformers import pipeline
                classifier = pipeline("zero-shot-classification", model="facebook/bart-large-mnli")
                result     = classifier(text[:1000], candidate_labels=labels)
                return ToolResult(True, f"✓ Top label: {result['labels'][0]} ({result['scores'][0]:.2f})", {"labels": result["labels"], "scores": result["scores"]})
            # Ollama fallback
            from npmai import Ollama
            llm   = Ollama(model="llama3.2:3b", temperature=0.1, change=True, Models=["mistral:7b"])
            resp  = llm.invoke(f"Classify the following text into ONE of these categories: {labels}. Reply with JSON only: {{\"label\":\"...\",\"confidence\":0.0}}.\nText: {text[:500]}")
            clean = re.search(r'\{.*?\}', resp, re.DOTALL)
            data  = json.loads(clean.group()) if clean else {"label": labels[0], "confidence": 0.5}
            return ToolResult(True, f"✓ Classified as: {data['label']}", data)
        except Exception as e:
            return ToolResult(False, f"✗ classify_text failed: {e}")

    @staticmethod
    def extract_entities(text: str, entity_types: Optional[List[str]] = None) -> ToolResult:
        try:
            from transformers import pipeline
            ner = pipeline("ner", model="dbmdz/bert-large-cased-finetuned-conll03-english", aggregation_strategy="simple")
            raw = ner(text[:512])
            entities = [{"entity": e["entity_group"], "word": e["word"], "score": round(e["score"], 3)} for e in raw]
            if entity_types:
                entities = [e for e in entities if e["entity"] in entity_types]
            return ToolResult(True, f"✓ Extracted {len(entities)} entities", entities)
        except Exception as e:
            # fallback: simple regex-based NER
            try:
                entities = []
                # emails
                for m in re.finditer(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", text):
                    entities.append({"entity": "EMAIL", "word": m.group()})
                # URLs
                for m in re.finditer(r"https?://[^\s]+", text):
                    entities.append({"entity": "URL", "word": m.group()})
                # Dates
                for m in re.finditer(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b|\b\d{4}-\d{2}-\d{2}\b", text):
                    entities.append({"entity": "DATE", "word": m.group()})
                # Capitalized sequences (potential names/orgs)
                for m in re.finditer(r"(?<!\. )\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)\b", text):
                    entities.append({"entity": "PERSON_OR_ORG", "word": m.group()})
                return ToolResult(True, f"✓ Extracted {len(entities)} entities (regex fallback)", entities)
            except Exception as e2:
                return ToolResult(False, f"✗ extract_entities failed: {e} / {e2}")

    @staticmethod
    def extract_keywords(text: str, n: int = 10, method: str = "tfidf") -> ToolResult:
        try:
            if method == "tfidf":
                from sklearn.feature_extraction.text import TfidfVectorizer
                import numpy as np
                vectorizer = TfidfVectorizer(stop_words="english", max_features=200, ngram_range=(1, 2))
                tfidf      = vectorizer.fit_transform([text])
                scores     = zip(vectorizer.get_feature_names_out(), tfidf.toarray()[0])
                keywords   = sorted(scores, key=lambda x: x[1], reverse=True)[:n]
                return ToolResult(True, f"✓ Top {n} keywords (TF-IDF)", [{"keyword": k, "score": round(s, 4)} for k, s in keywords])
            elif method == "textrank":
                from textblob import TextBlob
                blob = TextBlob(text)
                # noun phrases as keywords
                nps = list(set(blob.noun_phrases))[:n]
                return ToolResult(True, f"✓ {len(nps)} noun phrase keywords", nps)
            else:
                # frequency-based
                words  = re.findall(r"\b[a-zA-Z]{4,}\b", text.lower())
                stops  = {"this", "that", "with", "from", "have", "been", "will", "they", "their", "there", "were", "what", "when", "which", "your", "into", "over", "such", "also", "than", "then", "them", "some", "more"}
                counts = {}
                for w in words:
                    if w not in stops:
                        counts[w] = counts.get(w, 0) + 1
                top = sorted(counts.items(), key=lambda x: x[1], reverse=True)[:n]
                return ToolResult(True, f"✓ Top {n} keywords (frequency)", [{"keyword": k, "count": c} for k, c in top])
        except Exception as e:
            return ToolResult(False, f"✗ extract_keywords failed: {e}")

    @staticmethod
    def summarize(
        text: str,
        length: str = "medium",
        style: str = "abstractive",
        model: str = "ollama",
    ) -> ToolResult:
        try:
            max_words = {"short": 50, "medium": 150, "long": 300}.get(length, 150)
            if style == "extractive":
                from textblob import TextBlob
                sentences = TextBlob(text).sentences
                top_n     = {"short": 2, "medium": 4, "long": 6}.get(length, 4)
                summary   = " ".join(str(s) for s in sentences[:top_n])
                return ToolResult(True, f"✓ Extractive summary ({top_n} sentences)", summary)
            # abstractive via Ollama or transformers
            if model == "ollama":
                from npmai import Ollama
                llm  = Ollama(model="llama3.2:3b", temperature=0.2, change=True, Models=["mistral:7b"])
                resp = llm.invoke(f"Summarize the following text in approximately {max_words} words:\n\n{text[:4000]}")
                return ToolResult(True, f"✓ Summary ({len(resp.split())} words)", resp)
            else:
                from transformers import pipeline
                summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
                max_len    = min(max_words * 2, 150)
                result     = summarizer(text[:1024], max_length=max_len, min_length=20, do_sample=False)
                return ToolResult(True, f"✓ Summary", result[0]["summary_text"])
        except Exception as e:
            return ToolResult(False, f"✗ summarize failed: {e}")

    @staticmethod
    def translate(text: str, target_language: str = "es", model: str = "ollama") -> ToolResult:
        try:
            if model == "ollama":
                from npmai import Ollama
                lang_names = {"es": "Spanish", "fr": "French", "de": "German", "ja": "Japanese", "zh": "Chinese", "ar": "Arabic", "pt": "Portuguese", "hi": "Hindi", "ru": "Russian", "ko": "Korean"}
                lang_name  = lang_names.get(target_language, target_language)
                llm   = Ollama(model="llama3.2:3b", temperature=0.1, change=True, Models=["mistral:7b"])
                resp  = llm.invoke(f"Translate the following text to {lang_name}. Output ONLY the translation, nothing else:\n\n{text}")
                return ToolResult(True, f"✓ Translated to {lang_name}", resp.strip())
            from transformers import pipeline
            translator = pipeline("translation", model=f"Helsinki-NLP/opus-mt-en-{target_language}")
            result     = translator(text[:512])
            return ToolResult(True, f"✓ Translated", result[0]["translation_text"])
        except Exception as e:
            return ToolResult(False, f"✗ translate failed: {e}")

    @staticmethod
    def detect_language(text: str) -> ToolResult:
        try:
            from langdetect import detect, detect_langs
            lang  = detect(text)
            probs = detect_langs(text)
            langs = [{"lang": str(p).split(":")[0], "prob": float(str(p).split(":")[1])} for p in probs]
            return ToolResult(True, f"✓ Detected language: {lang}", {"language": lang, "probabilities": langs})
        except Exception as e:
            return ToolResult(False, f"✗ detect_language failed: {e}")

    @staticmethod
    def check_grammar(text: str) -> ToolResult:
        try:
            import language_tool_python
            tool    = language_tool_python.LanguageTool("en-US")
            matches = tool.check(text)
            errors  = [{"message": m.message, "context": m.context, "replacements": m.replacements[:3], "offset": m.offset, "length": m.errorLength} for m in matches[:20]]
            corrected = language_tool_python.utils.correct(text, matches)
            tool.close()
            return ToolResult(True, f"✓ {len(errors)} grammar issues found", {"errors": errors, "corrected": corrected})
        except ImportError:
            return ToolResult(False, "✗ language-tool-python not installed. Run: pip install language-tool-python")
        except Exception as e:
            return ToolResult(False, f"✗ check_grammar failed: {e}")

    @staticmethod
    def calculate_readability(text: str) -> ToolResult:
        try:
            words     = text.split()
            sentences = re.split(r"[.!?]+", text)
            sentences = [s.strip() for s in sentences if s.strip()]
            syllables = sum(max(1, len(re.findall(r"[aeiouAEIOU]", w))) for w in words)
            num_words = len(words) or 1
            num_sents = len(sentences) or 1
            # Flesch Reading Ease
            fre = 206.835 - 1.015 * (num_words / num_sents) - 84.6 * (syllables / num_words)
            # Flesch-Kincaid Grade Level
            fkgl = 0.39 * (num_words / num_sents) + 11.8 * (syllables / num_words) - 15.59
            # SMOG
            poly_count = sum(1 for w in words if len(re.findall(r"[aeiouAEIOU]", w)) >= 3)
            smog = 3.1291 + 1.0430 * (poly_count * (30 / num_sents)) ** 0.5
            data = {
                "flesch_reading_ease":   round(fre, 1),
                "flesch_kincaid_grade":  round(fkgl, 1),
                "smog_grade":            round(smog, 1),
                "words":                 num_words,
                "sentences":             num_sents,
                "avg_words_per_sentence":round(num_words / num_sents, 1),
                "syllables_per_word":    round(syllables / num_words, 2),
            }
            ease = "Very Easy" if fre > 80 else "Easy" if fre > 60 else "Average" if fre > 40 else "Difficult" if fre > 20 else "Very Difficult"
            return ToolResult(True, f"✓ Readability: {ease} (FRE={fre:.1f}, Grade={fkgl:.1f})", data)
        except Exception as e:
            return ToolResult(False, f"✗ calculate_readability failed: {e}")

    @staticmethod
    def topic_modeling(
        texts: List[str],
        n_topics: int = 5,
        method: str = "lda",
        output: str = "topics.json",
    ) -> ToolResult:
        try:
            from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer
            from sklearn.decomposition import LatentDirichletAllocation, NMF
            vectorizer = CountVectorizer(stop_words="english", max_features=1000, min_df=1) if method == "lda" else TfidfVectorizer(stop_words="english", max_features=1000)
            X = vectorizer.fit_transform(texts)
            model = LatentDirichletAllocation(n_components=n_topics, random_state=42) if method == "lda" else NMF(n_components=n_topics, random_state=42)
            model.fit(X)
            feature_names = vectorizer.get_feature_names_out()
            topics = []
            for i, comp in enumerate(model.components_):
                top_words = [feature_names[j] for j in comp.argsort()[-10:][::-1]]
                topics.append({"topic": i + 1, "top_words": top_words})
            Path(output).write_text(json.dumps(topics, indent=2))
            return ToolResult(True, f"✓ {n_topics} topics extracted ({method.upper()}), saved to '{output}'", topics)
        except Exception as e:
            return ToolResult(False, f"✗ topic_modeling failed: {e}")

    @staticmethod
    def text_similarity(text1: str, text2: str, method: str = "cosine") -> ToolResult:
        try:
            if method in ("cosine", "tfidf"):
                from sklearn.feature_extraction.text import TfidfVectorizer
                from sklearn.metrics.pairwise import cosine_similarity
                tfidf  = TfidfVectorizer().fit_transform([text1, text2])
                sim    = float(cosine_similarity(tfidf[0], tfidf[1])[0][0])
                return ToolResult(True, f"✓ Cosine similarity: {sim:.4f}", {"similarity": sim, "method": "cosine_tfidf"})
            elif method == "embedding":
                from sentence_transformers import SentenceTransformer
                from sklearn.metrics.pairwise import cosine_similarity
                import numpy as np
                model = SentenceTransformer("all-MiniLM-L6-v2")
                embs  = model.encode([text1, text2])
                sim   = float(cosine_similarity([embs[0]], [embs[1]])[0][0])
                return ToolResult(True, f"✓ Semantic similarity: {sim:.4f}", {"similarity": sim, "method": "sentence_transformer"})
            elif method == "levenshtein":
                import numpy as np
                m, n    = len(text1), len(text2)
                dp      = [[0] * (n + 1) for _ in range(m + 1)]
                for i in range(m + 1): dp[i][0] = i
                for j in range(n + 1): dp[0][j] = j
                for i in range(1, m + 1):
                    for j in range(1, n + 1):
                        if text1[i-1] == text2[j-1]: dp[i][j] = dp[i-1][j-1]
                        else: dp[i][j] = 1 + min(dp[i-1][j], dp[i][j-1], dp[i-1][j-1])
                dist = dp[m][n]
                sim  = 1 - dist / max(m, n, 1)
                return ToolResult(True, f"✓ Levenshtein similarity: {sim:.4f}", {"similarity": sim, "edit_distance": dist})
            return ToolResult(False, f"✗ Unknown method: {method}")
        except Exception as e:
            return ToolResult(False, f"✗ text_similarity failed: {e}")

    @staticmethod
    def generate_embeddings(texts: Union[str, List[str]], model: str = "all-MiniLM-L6-v2") -> ToolResult:
        try:
            from sentence_transformers import SentenceTransformer
            if isinstance(texts, str): texts = [texts]
            encoder    = SentenceTransformer(model)
            embeddings = encoder.encode(texts, convert_to_numpy=True)
            return ToolResult(True, f"✓ Generated {len(embeddings)} embeddings (dim={embeddings.shape[1]})", {"embeddings": embeddings.tolist(), "shape": list(embeddings.shape)})
        except Exception as e:
            return ToolResult(False, f"✗ generate_embeddings failed: {e}")

    @staticmethod
    def semantic_search(query: str, corpus: List[str], top_k: int = 5) -> ToolResult:
        try:
            from sentence_transformers import SentenceTransformer
            from sklearn.metrics.pairwise import cosine_similarity
            import numpy as np
            model  = SentenceTransformer("all-MiniLM-L6-v2")
            q_emb  = model.encode([query])
            c_embs = model.encode(corpus)
            sims   = cosine_similarity(q_emb, c_embs)[0]
            top_indices = np.argsort(sims)[::-1][:top_k]
            results = [{"index": int(i), "text": corpus[i][:200], "score": float(sims[i])} for i in top_indices]
            return ToolResult(True, f"✓ Top {top_k} semantic matches", results)
        except Exception as e:
            return ToolResult(False, f"✗ semantic_search failed: {e}")

    @staticmethod
    def detect_plagiarism(
        text: str,
        reference_texts: List[str],
        threshold: float = 0.7,
    ) -> ToolResult:
        try:
            from sklearn.feature_extraction.text import TfidfVectorizer
            from sklearn.metrics.pairwise import cosine_similarity
            all_texts  = [text] + reference_texts
            vectorizer = TfidfVectorizer(stop_words="english")
            tfidf      = vectorizer.fit_transform(all_texts)
            sims       = cosine_similarity(tfidf[0], tfidf[1:])[0]
            matches = [{"reference_index": i, "similarity": float(sims[i]), "flag": sims[i] >= threshold, "excerpt": reference_texts[i][:200]} for i in range(len(reference_texts))]
            flagged  = [m for m in matches if m["flag"]]
            return ToolResult(True, f"✓ {len(flagged)}/{len(reference_texts)} references flagged (threshold={threshold})", {"matches": matches, "flagged_count": len(flagged)})
        except Exception as e:
            return ToolResult(False, f"✗ detect_plagiarism failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 9. DatabaseTool
# ─────────────────────────────────────────────────────────────────────────────

class DatabaseTool:
    name = "database"
    description = (
        "Multi-database operations: PostgreSQL, MySQL, MongoDB, Redis, SQLite. "
        "Query execution, transactions, backup/restore, schema inspection, CSV import/export."
    )
    use = (
        """Name of Tool:- DatabaseTool

Purpose of Tool:-
The DatabaseTool provides a unified interface for executing multi-engine database operations across relational (PostgreSQL, MySQL, SQLite), document-oriented (MongoDB), and key-value/in-memory systems (Redis). It abstracts routine administrative workflows including connection state management, single/transactional query dispatch, structural schema indexing, backup/restore streams, and data serialization between database tables and local flat CSV sheets.

Methods:-
- connect_postgres: Initializes a connection instance to a PostgreSQL server cluster.
- execute_query: Runs structured queries against an open PostgreSQL or MySQL relational connection, converting rows into dynamic dictionaries.
- execute_transaction: Wraps a sequence of modification commands within an explicit ACID transaction block.
- backup_postgres: Triggers a shell subprocess executing `pg_dump` to capture full snapshots of a given database.
- restore_postgres: Feeds a saved snapshot file directly back into a targeted database connection utilizing native command utilities.
- get_schema: Queries the relational database information schema to map public column specifications and type requirements.
- connect_mysql: Opens a connection stream targeting a MySQL database instance.
- connect_mongodb: Resolves a standard MongoDB connection string connection into an operational database mapping client.
- mongo_find: Queries a Document store cluster collection, translating object primitives into standard JSON serializable objects.
- mongo_insert: Packs single or multi-document dictionary definitions straight into a designated NoSQL collection.
- mongo_update: Finds and replaces targeted attributes across one or many structured documents.
- mongo_delete: Triggers structural document deletions for objects meeting explicit filter parameters.
- connect_redis: Verifies connectivity to a live key-value Redis storage server instance.
- redis_set: Maps an individual string variable to a unique key string with optional Time-To-Live parameters.
- redis_get: Extracts cached string payloads mapped to specific key locations.
- redis_hset: Groups complex multi-attribute dictionary maps inside structured Redis hash definitions.
- redis_hget: Retreives individual string value fragments nested inside specified Redis hashes.
- redis_lpush: Prepends sequence vectors directly onto a tracking Redis array index location.
- export_to_csv: Pipelines data records straight into isolated tabular `.csv` sheets using data science frameworks.
- import_from_csv: Parses local flat `.csv` records to populate a live target database table framework using automatic fallbacks.
- create_sqlite_db: Instantiates a standalone relational SQLite database file embedded on local disk media.
- query_sqlite: Executes commands targeting a local embedded SQLite file.

How to use Tool Methods:-

1. connect_postgres:
   - Purpose: Instantiates an open connection handle to a live PostgreSQL relational instance.
   - Arguments:
     a) host: str (default: "localhost") - Destination server cluster IP network endpoint text.
     b) port: int (default: 5432) - Active networking portal path constraints.
     c) database: str - Target schema namespace text identifier.
     b) user: str - Access profile identity authorization text string.
     e) password: str - Authentication entry passphrase string.
     f) cred_key: str (default: "postgres") - Internal credential vault profile index configuration string.
   - Returns: Connection object allowing communication down the pipeline.
   - How to call: DatabaseTool.connect_postgres(database="production_db", user="admin")

2. execute_query:
   - Purpose: Dispatches custom commands against an active engine connection handle.
   - Arguments:
     a) conn: object - Active, instantiated connection session reference hook.
     b) query: str - Structured query language string statement text.
     c) params: Optional[tuple] (default: None) - Variable mapping parameters to prevent injections.
   - Returns: ToolResult passing structured dict rows or record count change tallies.
   - How to call: DatabaseTool.execute_query(conn, "SELECT * FROM users WHERE status = %s", ("active",))

3. execute_transaction:
   - Purpose: Ensures uniform processing of multi-query changes via structural rollback handling.
   - Arguments:
     a) conn: object - Destination target engine connection session wrapper.
     b) queries: List[str] - Sequence array containing raw manipulation queries.
   - Returns: ToolResult validating transaction execution status properties.
   - How to call: DatabaseTool.execute_transaction(conn, ["UPDATE accounts SET balance = 100 WHERE id = 1", "INSERT INTO logs (msg) VALUES ('updated')"])

4. backup_postgres:
   - Purpose: Dumps full schema structures and inline records down into isolated recovery files.
   - Arguments:
     a) conn_string: str - Server connectivity profile description array.
     b) output_file: str - Target workspace dump output path tracking string.
   - Returns: ToolResult tracking processing performance metrics along with byte sizes.
   - How to call: DatabaseTool.backup_postgres(conn_string="postgresql://user:pass@localhost/db", output_file="/backups/daily.sql")

5. restore_postgres:
   - Purpose: Imports complete schema layouts back into active database configurations.
   - Arguments:
     a) conn_string: str - Server initialization credentials statement string.
     b) backup_file: str - Workspace path indexing the target resource backup source file.
   - Returns: ToolResult capturing standard error stream outputs and success statements.
   - How to call: DatabaseTool.restore_postgres(conn_string="postgresql://user:pass@localhost/db", backup_file="/backups/daily.sql")

6. get_schema:
   - Purpose: Audits structure types to help write clean queries.
   - Arguments:
     a) conn: object - Live active connection context platform session wrapper.
     b) table_or_all: str (default: "all") - Targeted structural context table selection label.
   - Returns: ToolResult breaking down schemas, defaults, and data properties.
   - How to call: DatabaseTool.get_schema(conn, table_or_all="orders")

7. connect_mysql:
   - Purpose: Establishes a connection interface link targeting a live MySQL relational endpoint.
   - Arguments:
     a) host: str (default: "localhost") - Network host server interface string.
     b) port: int (default: 3306) - System communication interface ports channel index.
     c) database: str - Targeted database name context.
     d) user: str - Connection authorization credentials identity text string.
     e) password: str - Authentication code access passphrase string.
     f) cred_key: str (default: "mysql") - Credentials database tracking index mapping path text.
   - Returns: Connection resource interface allowing programmatic transactional pipeline routines.
   - How to call: DatabaseTool.connect_mysql(database="ecommerce", user="root")

8. connect_mongodb:
   - Purpose: Initiates a document client engine targeting distributed NoSQL databases.
   - Arguments:
     a) uri: str - Multi-node connection initialization endpoint string text.
     b) database: str - Target collection namespace definition label string.
     c) cred_key: str (default: "mongodb") - Internal configuration management credential vault index pointer.
   - Returns: Pymongo workspace Database layer reference hook or raw Client index engine.
   - How to call: DatabaseTool.connect_mongodb(uri="mongodb://localhost:27017/", database="app_store")

9. mongo_find:
   - Purpose: Searches document repositories to filter and return serializable objects.
   - Arguments:
     a) db: object - Verified document client instance profile.
     b) collection: str - target data collection grouping label.
     c) filter: Optional[Dict] (default: None) - Filter conditions dict.
     d) projection: Optional[Dict] (default: None) - Target field inclusion or exclusion rules.
     e) limit: int (default: 100) - Maximum return threshold capacity limits integer constraints.
     f) sort: Optional[List[Tuple[str, int]]] (default: None) - Target arrangement order sorting tuple indices array.
   - Returns: ToolResult passing document arrays without native Object primitives.
   - How to call: DatabaseTool.mongo_find(db, collection="logs", filter={"level": "error"}, limit=10)

10. mongo_insert:
    - Purpose: Saves single document records or batch arrays into a given collection.
    - Arguments:
      a) db: object - target database wrapper anchor context session object.
      b) collection: str - target data storage domain identifier group.
      b) documents: Union[Dict, List[Dict]] - Individual dictionary record map or sequence list array containing document dictionaries.
    - Returns: ToolResult reflecting successful item counts along with standard document string references.
    - How to call: DatabaseTool.mongo_insert(db, collection="users", documents={"name": "Alice", "role": "dev"})

11. mongo_update:
    - Purpose: Performs field manipulations inside a document layer collection framework.
    - Arguments:
      a) db: object - target database driver environment handle.
      b) collection: str - Target destination storage array namespace label.
      c) filter: Dict - Target locating search query criteria rule statement dictionary.
      d) update: Dict - Modification instructions structured dictionary.
      e) many: bool (default: False) - System structural wide scaling directive choices option flag.
    - Returns: ToolResult tracking matching records vs modified asset counts.
    - How to call: DatabaseTool.mongo_update(db, collection="users", filter={"role": "dev"}, update={"$set": {"status": "verified"}}, many=True)

12. mongo_delete:
    - Purpose: Clears database document objects based on matching filter queries.
    - Arguments:
      a) db: object - Core platform collection cluster runtime layer workspace connector.
      b) collection: str - Targeted system tracking storage space category identifier.
      c) filter: Dict - Structural parameters dictionary isolating elements targeted for removal.
      d) many: bool (default: False) - Multi-element operational sweep selection tracking flag.
    - Returns: ToolResult showing total records deleted from the collection.
    - How to call: DatabaseTool.mongo_delete(db, collection="sessions", filter={"expired": True}, many=True)

13. connect_redis:
    - Purpose: Opens up quick channels to interact with high-speed volatile storage grids.
    - Arguments:
      a) host: str (default: "localhost") - Network host server interface string.
      b) port: int (default: 6379) - Port parameter configuration mapping location identifier.
      c) password: str - Target security string profile access code passphrase.
      d) db: int (default: 0) - Active virtual memory tracking frame slot address allocation.
      o) cred_key: str (default: "redis") - Storage secrets locator lookup profile identifier text.
    - Returns: Active Redis pipeline client utility engine entity context model.
    - How to call: DatabaseTool.connect_redis(host="cache.local", db=1)

14. redis_set:
    - Purpose: Assigns scalar entries straight onto global fast-access tracking namespaces.
    - Arguments:
      a) r: object - Active initialized Redis engine workspace handle.
      b) key: str - Reference identity variable identifier tracking string name.
      c) value: str - Target values literal data element payload.
      d) ttl: Optional[int] (default: None) - Life window constraints duration limit metrics in seconds.
    - Returns: ToolResult confirming key allocation properties.
    - How to call: DatabaseTool.redis_set(r, key="session_101", value="active_user", ttl=3600)

15. redis_get:
    - Purpose: Pulls key values from the memory database.
    - Arguments:
      a) r: object - Verified workspace interface context runner handle.
      b) key: str - Destination tracking reference identity key term.
    - Returns: ToolResult housing string values or target omission errors.
    - How to call: DatabaseTool.redis_get(r, key="session_101")

16. redis_hset:
    - Purpose: Serializes nested dictionary tables directly inside localized high-speed memory maps.
    - Arguments:
      a) r: object - Live environment tracking memory platform driver context layer.
      b) name: str - Target group reference identity tag term string text.
      c) mapping: Dict[str, str] - Multi-attribute parameter layout element configurations data table.
    - Returns: ToolResult displaying field initialization stats.
    - How to call: DatabaseTool.redis_hset(r, name="user:1001", mapping={"name": "Bob", "auth": "true"})

17. redis_hget:
    - Purpose: Extracts isolated attributes from multi-field memory maps.
    - Arguments:
      a) r: object - Active volatile engine system connection reference session hook.
      b) name: str - Target cluster group identification key tracking text label.
      c) key: str - target attribute tracking item property parameter identifier string name.
    - Returns: ToolResult delivering the internal attribute value.
    - How to call: DatabaseTool.redis_hget(r, name="user:1001", key="auth")

18. redis_lpush:
    - Purpose: Pushes entry parameters directly onto volatile queuing pipelines.
    - Arguments:
      a) r: object - Active programmatic memory grid connector handle session.
      b) key: str - Target tracker pipeline registry identity tag name string text.
      c) values: List[str] - Data components array index sequence collection.
    - Returns: ToolResult showing updated length calculations.
    - How to call: DatabaseTool.redis_lpush(r, key="task_queue", values=["task_1", "task_2"])

19. export_to_csv:
    - Purpose: Formats tabular record data tables into file blocks.
    - Arguments:
      a) conn: object - Valid open query engine driver configuration handler connection.
      b) query: str - Target source collection extraction query framework directive.
      c) output: str - Target workspace output destination file configuration tracking framework.
    - Returns: ToolResult tracking row extraction limits exported down onto local storage media.
    - How to call: DatabaseTool.export_to_csv(conn, "SELECT * FROM products", "products_export.csv")

20. import_from_csv:
    - Purpose: Parses flat file tables to ingest records into existing tables.
    - Arguments:
      a) conn: object - Operational destination cluster runtime layer environment pipeline handle.
      b) table: str - Destination framework target schema matrix block identifier text.
      c) csv_path: str - Local spreadsheet resource data origin file tracking location path string text.
      d) if_exists: str (default: "append") - Action rules if structural conflicts manifest ("fail", "replace", "append").
    - Returns: ToolResult validating absolute load metrics indicators.
    - How to call: DatabaseTool.import_from_csv(conn, table="subscribers", csv_path="subscribers.csv")

21. create_sqlite_db:
    - Purpose: Provisions zero-infrastructure database environments straight onto local storage spaces.
    - Arguments:
      a) path: str - Target destination disk execution directory path text location indicator.
      b) schema: str - Multi-table SQL layout schema build scripts declaration block text.
    - Returns: ToolResult indicating environment deployment parameters.
    - How to call: DatabaseTool.create_sqlite_db(path="local_cache.db", schema="CREATE TABLE cache (id INT, data TEXT);")

22. query_sqlite:
    - Purpose: Interacts with embedded database single-file instances without network overhead.
    - Arguments:
      a) path: str - Workspace database deployment system file path tracking string text.
      b) query: str - Command script syntax target direction directive text label.
      c) params: Optional[tuple] (default: None) - Safe insertion query values data parameter tuple components.
    - Returns: ToolResult returning target dynamic list map rows or database commit tracking flags.
    - How to call: DatabaseTool.query_sqlite("local_cache.db", "SELECT * FROM cache WHERE id = ?", (101,))
    """)
    
    # ── PostgreSQL ────────────────────────────────────────────────────────────

    @staticmethod
    def connect_postgres(
        host: str = "localhost",
        port: int = 5432,
        database: str = "",
        user: str = "",
        password: str = "",
        cred_key: str = "postgres",
    ):
        try:
            import psycopg2
            c = CredStore.load(cred_key) if not database else {}
            conn = psycopg2.connect(
                host     = host     or c.get("host", "localhost"),
                port     = port     or int(c.get("port", 5432)),
                database = database or c.get("database", ""),
                user     = user     or c.get("user", ""),
                password = password or c.get("password", ""),
            )
            conn.autocommit = True
            return conn
        except Exception as e:
            raise ConnectionError(f"PostgreSQL connect failed: {e}")

    @staticmethod
    def execute_query(conn, query: str, params: Optional[tuple] = None) -> ToolResult:
        try:
            import psycopg2.extras
            with conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor) as cur:
                cur.execute(query, params)
                try:
                    rows = [dict(r) for r in cur.fetchall()]
                    return ToolResult(True, f"✓ {len(rows)} rows returned", rows)
                except Exception:
                    return ToolResult(True, f"✓ Query executed, {cur.rowcount} rows affected")
        except Exception as e:
            return ToolResult(False, f"✗ execute_query failed: {e}")

    @staticmethod
    def execute_transaction(conn, queries: List[str]) -> ToolResult:
        try:
            old_autocommit   = conn.autocommit
            conn.autocommit  = False
            with conn.cursor() as cur:
                for q in queries:
                    cur.execute(q)
            conn.commit()
            conn.autocommit = old_autocommit
            return ToolResult(True, f"✓ Transaction committed ({len(queries)} queries)")
        except Exception as e:
            try: conn.rollback()
            except Exception: pass
            return ToolResult(False, f"✗ execute_transaction failed (rolled back): {e}")

    @staticmethod
    def backup_postgres(conn_string: str, output_file: str) -> ToolResult:
        try:
            result = subprocess.run(
                ["pg_dump", conn_string, "-f", output_file],
                capture_output=True, text=True, timeout=300,
            )
            if result.returncode != 0:
                return ToolResult(False, f"✗ pg_dump failed: {result.stderr}")
            size = Path(output_file).stat().st_size
            return ToolResult(True, f"✓ PostgreSQL backup saved to '{output_file}' ({size} bytes)")
        except Exception as e:
            return ToolResult(False, f"✗ backup_postgres failed: {e}")

    @staticmethod
    def restore_postgres(conn_string: str, backup_file: str) -> ToolResult:
        try:
            result = subprocess.run(
                ["psql", conn_string, "-f", backup_file],
                capture_output=True, text=True, timeout=600,
            )
            if result.returncode != 0:
                return ToolResult(False, f"✗ psql restore failed: {result.stderr}")
            return ToolResult(True, f"✓ PostgreSQL restored from '{backup_file}'")
        except Exception as e:
            return ToolResult(False, f"✗ restore_postgres failed: {e}")

    @staticmethod
    def get_schema(conn, table_or_all: str = "all") -> ToolResult:
        try:
            import psycopg2.extras
            with conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor) as cur:
                if table_or_all == "all":
                    cur.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'public'")
                    tables = [r["table_name"] for r in cur.fetchall()]
                    schema = {}
                    for t in tables:
                        cur.execute("SELECT column_name, data_type, is_nullable FROM information_schema.columns WHERE table_name = %s ORDER BY ordinal_position", (t,))
                        schema[t] = [dict(r) for r in cur.fetchall()]
                    return ToolResult(True, f"✓ Schema for {len(tables)} tables", schema)
                else:
                    cur.execute("SELECT column_name, data_type, is_nullable, column_default FROM information_schema.columns WHERE table_name = %s ORDER BY ordinal_position", (table_or_all,))
                    cols = [dict(r) for r in cur.fetchall()]
                    return ToolResult(True, f"✓ Schema for table '{table_or_all}'", cols)
        except Exception as e:
            return ToolResult(False, f"✗ get_schema failed: {e}")

    # ── MySQL ─────────────────────────────────────────────────────────────────

    @staticmethod
    def connect_mysql(
        host: str = "localhost",
        port: int = 3306,
        database: str = "",
        user: str = "",
        password: str = "",
        cred_key: str = "mysql",
    ):
        try:
            import mysql.connector
            c    = CredStore.load(cred_key) if not database else {}
            conn = mysql.connector.connect(
                host     = host     or c.get("host", "localhost"),
                port     = port     or int(c.get("port", 3306)),
                database = database or c.get("database", ""),
                user     = user     or c.get("user", ""),
                password = password or c.get("password", ""),
            )
            return conn
        except Exception as e:
            raise ConnectionError(f"MySQL connect failed: {e}")

    # ── MongoDB ───────────────────────────────────────────────────────────────

    @staticmethod
    def connect_mongodb(uri: str = "", database: str = "", cred_key: str = "mongodb"):
        try:
            import pymongo
            c    = CredStore.load(cred_key)
            uri  = uri  or c.get("uri",  "mongodb://localhost:27017/")
            db   = database or c.get("database", "")
            client = pymongo.MongoClient(uri)
            return client[db] if db else client
        except Exception as e:
            raise ConnectionError(f"MongoDB connect failed: {e}")

    @staticmethod
    def mongo_find(
        db,
        collection: str,
        filter: Optional[Dict] = None,
        projection: Optional[Dict] = None,
        limit: int = 100,
        sort: Optional[List[Tuple[str, int]]] = None,
    ) -> ToolResult:
        try:
            coll   = db[collection]
            cursor = coll.find(filter or {}, projection)
            if sort: cursor = cursor.sort(sort)
            cursor = cursor.limit(limit)
            docs   = [{k: (str(v) if str(type(v)) == "<class 'bson.objectid.ObjectId'>" else v) for k, v in d.items()} for d in cursor]
            return ToolResult(True, f"✓ {len(docs)} documents from '{collection}'", docs)
        except Exception as e:
            return ToolResult(False, f"✗ mongo_find failed: {e}")

    @staticmethod
    def mongo_insert(db, collection: str, documents: Union[Dict, List[Dict]]) -> ToolResult:
        try:
            coll = db[collection]
            if isinstance(documents, dict):
                result = coll.insert_one(documents)
                return ToolResult(True, f"✓ Inserted 1 document, id={result.inserted_id}")
            result = coll.insert_many(documents)
            return ToolResult(True, f"✓ Inserted {len(result.inserted_ids)} documents", {"inserted_ids": [str(i) for i in result.inserted_ids]})
        except Exception as e:
            return ToolResult(False, f"✗ mongo_insert failed: {e}")

    @staticmethod
    def mongo_update(
        db,
        collection: str,
        filter: Dict,
        update: Dict,
        many: bool = False,
    ) -> ToolResult:
        try:
            coll = db[collection]
            if many:
                result = coll.update_many(filter, update)
            else:
                result = coll.update_one(filter, update)
            return ToolResult(True, f"✓ Updated {result.modified_count} document(s)", {"matched": result.matched_count, "modified": result.modified_count})
        except Exception as e:
            return ToolResult(False, f"✗ mongo_update failed: {e}")

    @staticmethod
    def mongo_delete(db, collection: str, filter: Dict, many: bool = False) -> ToolResult:
        try:
            coll = db[collection]
            if many:
                result = coll.delete_many(filter)
            else:
                result = coll.delete_one(filter)
            return ToolResult(True, f"✓ Deleted {result.deleted_count} document(s)")
        except Exception as e:
            return ToolResult(False, f"✗ mongo_delete failed: {e}")

    # ── Redis ─────────────────────────────────────────────────────────────────

    @staticmethod
    def connect_redis(
        host: str = "localhost",
        port: int = 6379,
        password: str = "",
        db: int = 0,
        cred_key: str = "redis",
    ):
        try:
            import redis
            c = CredStore.load(cred_key)
            r = redis.Redis(
                host     = host     or c.get("host", "localhost"),
                port     = port     or int(c.get("port", 6379)),
                password = password or c.get("password", None) or None,
                db       = db,
                decode_responses=True,
            )
            r.ping()
            return r
        except Exception as e:
            raise ConnectionError(f"Redis connect failed: {e}")

    @staticmethod
    def redis_set(r, key: str, value: str, ttl: Optional[int] = None) -> ToolResult:
        try:
            if ttl: r.setex(key, ttl, value)
            else:   r.set(key, value)
            return ToolResult(True, f"✓ Redis SET '{key}'")
        except Exception as e:
            return ToolResult(False, f"✗ redis_set failed: {e}")

    @staticmethod
    def redis_get(r, key: str) -> ToolResult:
        try:
            val = r.get(key)
            if val is None:
                return ToolResult(False, f"✗ Key '{key}' not found in Redis")
            return ToolResult(True, f"✓ Redis GET '{key}'", val)
        except Exception as e:
            return ToolResult(False, f"✗ redis_get failed: {e}")

    @staticmethod
    def redis_hset(r, name: str, mapping: Dict[str, str]) -> ToolResult:
        try:
            r.hset(name, mapping=mapping)
            return ToolResult(True, f"✓ Redis HSET '{name}' ({len(mapping)} fields)")
        except Exception as e:
            return ToolResult(False, f"✗ redis_hset failed: {e}")

    @staticmethod
    def redis_hget(r, name: str, key: str) -> ToolResult:
        try:
            val = r.hget(name, key)
            if val is None:
                return ToolResult(False, f"✗ Hash key '{name}.{key}' not found")
            return ToolResult(True, f"✓ Redis HGET '{name}.{key}'", val)
        except Exception as e:
            return ToolResult(False, f"✗ redis_hget failed: {e}")

    @staticmethod
    def redis_lpush(r, key: str, values: List[str]) -> ToolResult:
        try:
            length = r.lpush(key, *values)
            return ToolResult(True, f"✓ Redis LPUSH '{key}', list length={length}")
        except Exception as e:
            return ToolResult(False, f"✗ redis_lpush failed: {e}")

    # ── CSV import/export ─────────────────────────────────────────────────────

    @staticmethod
    def export_to_csv(conn, query: str, output: str) -> ToolResult:
        try:
            import pandas as pd
            df = pd.read_sql_query(query, conn)
            df.to_csv(output, index=False)
            return ToolResult(True, f"✓ Exported {len(df)} rows to '{output}'")
        except Exception as e:
            return ToolResult(False, f"✗ export_to_csv failed: {e}")

    @staticmethod
    def import_from_csv(conn, table: str, csv_path: str, if_exists: str = "append") -> ToolResult:
        try:
            import pandas as pd
            from sqlalchemy import create_engine
            df = pd.read_csv(csv_path)
            # use pandas to_sql if sqlalchemy available
            try:
                engine = create_engine(conn.dsn if hasattr(conn, "dsn") else str(conn))
                df.to_sql(table, engine, if_exists=if_exists, index=False)
            except Exception:
                # fallback: manual INSERT
                cols    = ", ".join(df.columns)
                placeholders = ", ".join(["%s"] * len(df.columns))
                with conn.cursor() as cur:
                    for _, row in df.iterrows():
                        cur.execute(f"INSERT INTO {table} ({cols}) VALUES ({placeholders})", tuple(row))
                conn.commit()
            return ToolResult(True, f"✓ Imported {len(df)} rows into '{table}'")
        except Exception as e:
            return ToolResult(False, f"✗ import_from_csv failed: {e}")

    # ── SQLite ────────────────────────────────────────────────────────────────

    @staticmethod
    def create_sqlite_db(path: str, schema: str) -> ToolResult:
        try:
            conn = sqlite3.connect(path)
            conn.executescript(schema)
            conn.commit()
            conn.close()
            return ToolResult(True, f"✓ SQLite database created at '{path}'")
        except Exception as e:
            return ToolResult(False, f"✗ create_sqlite_db failed: {e}")

    @staticmethod
    def query_sqlite(
        path: str,
        query: str,
        params: Optional[tuple] = None,
    ) -> ToolResult:
        try:
            conn = sqlite3.connect(path)
            conn.row_factory = sqlite3.Row
            cur  = conn.execute(query, params or ())
            try:
                rows = [dict(r) for r in cur.fetchall()]
                conn.close()
                return ToolResult(True, f"✓ {len(rows)} rows from SQLite", rows)
            except Exception:
                conn.commit()
                conn.close()
                return ToolResult(True, f"✓ Query executed")
        except Exception as e:
            return ToolResult(False, f"✗ query_sqlite failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# 10. ReportGeneratorTool
# ─────────────────────────────────────────────────────────────────────────────

class ReportGeneratorTool:
    name = "report_generator"
    description = (
        "Automated report generation: PDF, Word, Excel, PowerPoint, research reports, "
        "scheduled reports, dashboard HTML reports, and template-based generation."
    )
    use = (
        """Name of Tool: ReportGeneratorTool

Purpose of Tool:
The ReportGeneratorTool provides an automated, unified interface to generate multi-format corporate, technical, and research reports (including PDF, Word, Excel, PowerPoint, and HTML). It handles end-to-end publishing workflows, incorporating structural design elements like tables of contents, multi-sheet analytical books with embedded charts, dynamic dashboard layouts, background scheduling runners, template string substitutions, and automated AI research drafting.

Methods:
- create_pdf_report: Generates a publication-grade PDF file including custom covers, dynamic pagination markers, explicit section styles, and automatic dataset tables.
- create_word_report: Assembles an editable Microsoft Word document (.docx) tracking document metadata, structured section hierarchies, and native tabular layouts.
- create_excel_report: Compiles structural data frames into styled Excel workbooks (.xlsx) featuring zebra-striping formatting, custom widths, and auto-generated bar charts.
- create_presentation: Provisions high-impact, programmatic presentation slides (.pptx) adhering to modern widescreen parameters and layout themes.
- generate_research_report: Drives local or remote Large Language Models to architect multi-section academic outlines, compose long-form prose, and compile relevant citation strings into a polished PDF artifact.
- schedule_report: Spins up concurrent background scheduler threads to execute periodic reporting blocks and dispatch completions via secure SMTP channels.
- create_dashboard_report: Packages dynamic metric maps alongside external data visualization components into an inline styled HTML dashboard application interface.
- generate_from_template: Evaluates flat text template structures to complete macro variable substitutions and emit completed documents to target disk targets using multiple format formats.

How to use Tool Methods:

1. create_pdf_report:
   - Purpose: Generates a publication-ready PDF document complete with cover page options, section rules, and embedded table definitions.
   - Arguments:
     a) title: str - Primary header text mapping for the main cover page.
     b) sections: List[Dict[str, str]] - Data structures containing 'title', 'content', and optional matrix list arrays mapped to 'table'.
     c) output: str (default: "report.pdf") - Target workspace file storage location string.
     d) logo: str - Optional filesystem string pointing to a corporate header graphic.
     e) author: str (default: "NPM Agent") - Identification name tag tracking responsibility metadata.
     f) date: str - Specific formatting placeholder; falls back to runtime capture if omitted.
     g) toc: bool (default: True) - Flag controlling table of contents index structure generation.
   - Returns: ToolResult mapping successful compilation flags along with complete file metrics.
   - How to call: ReportGeneratorTool.create_pdf_report(title="Q2 Audit", sections=[{"title": "Summary", "content": "All clear.", "table": [["ID", "Status"], [1, "OK"]]}] )

2. create_word_report:
   - Purpose: Assembles standard editable document hierarchies targeting word processing software suites.
   - Arguments:
     a) title: str - Document root header tracking text string.
     b) sections: List[Dict[str, str]] - Sequence collections parsing out text paragraphs and structural arrays.
     c) output: str (default: "report.docx") - Local system target path storage tracking variable.
     d) logo: str - Local file location reference locating corporate image assets.
     e) author: str (default: "NPM Agent") - Standard workspace system identification text indicator.
     f) styles: Optional[Dict[str, Any]] (default: None) - Structural formatting directives profile dictionary map.
   - Returns: ToolResult tracking processing performance traits along with total target bytes metrics.
   - How to call: ReportGeneratorTool.create_word_report(title="Project Scope", sections=[{"title": "Overview", "content": "Text here."}])

3. create_excel_report:
   - Purpose: Compiles data frame grids across independent workbook tabs applying specific visibility parameters.
   - Arguments:
     a) data_dict: Dict[str, Any] - Core data tracking dictionary containing data structures or Pandas DataFrames indexed by sheet labels.
     b) output: str (default: "report.xlsx") - Target output spreadsheet location string locator.
     c) charts: bool (default: True) - Structural directive flag forcing parsing of numerical tracks to build companion graphics.
     d) formatting: bool (default: True) - Determines whether grid design overrides are activated.
   - Returns: ToolResult identifying processed sheet counts and underlying binary output sizes.
   - How to call: ReportGeneratorTool.create_excel_report(data_dict={"Sales": [{"Item": "A", "Qty": 20}, {"Item": "B", "Qty": 35}]})

4. create_presentation:
   - Purpose: Programmatically renders complete slide graphics frameworks based on quick-reference bullet specifications.
   - Arguments:
     a) slides_data: List[Dict[str, Any]] - Collection arrays grouping title parameters, body statements, and lists of components.
     b) output: str (default: "presentation.pptx") - Document delivery deployment path selector.
     c) theme: str (default: "dark") - Canvas color rules layout template switcher switch ("light" or "dark").
   - Returns: ToolResult providing validation of absolute page count metrics completed.
   - How to call: ReportGeneratorTool.create_presentation(slides_data=[{"title": "Intro", "bullets": ["Point 1", "Point 2"]}])

5. generate_research_report:
   - Purpose: Coordinates semantic processing loops to draft authoritative materials based on flat subject strings.
   - Arguments:
     a) topic: str - Descriptive subject string directing the AI drafting scope.
     b) model: str (default: "llama3.2:3b") - Core language model selection index tracking framework identifier.
     c) output: str (default: "research_report.pdf") - Target workspace file storage location string.
     d) include_citations: bool (default: True) - Flag adding a references section with plausible source listings.
   - Returns: ToolResult showing completion states and sub-component generation tallies.
   - How to call: ReportGeneratorTool.generate_research_report(topic="Quantum Cryptography Evolution", model="mistral:7b")

6. schedule_report:
   - Purpose: Spawns independent execution engines tasked with compiling updates based on explicit chronological intervals.
   - Arguments:
     a) report_func: Callable - Programmatic target block parameter defining the reporting action.
     b) schedule: str - Timing configuration expression syntax (e.g., "every 10 minutes", "monday at 08:00").
     c) email_to: str - Optional routing parameters assigning transmission goals to a mail agent.
     d) output_folder: str (default: "scheduled_reports") - Filesystem directory root handling asset sorting.
   - Returns: ToolResult registering task acceptance inside the global background runner system context.
   - How to call: ReportGeneratorTool.schedule_report(report_func=my_pdf_builder, schedule="every 5 minutes", email_to="lead@firm.com")

7. create_dashboard_report:
   - Purpose: Merges standalone numerical summaries and external charting sources into an isolated HTML display.
   - Arguments:
     a) metrics: Dict[str, Any] - Core flat parameters index capturing tracking KPIs.
     b) charts: List[str] - Local target document paths capturing vector or layout visual blocks.
     c) output_html: str (default: "dashboard_report.html") - Web deployment string location file index parameter.
     d) title: str (default: "Dashboard Report") - Core application structural view header label text.
   - Returns: ToolResult documenting dashboard component tracking distributions.
   - How to call: ReportGeneratorTool.create_dashboard_report(metrics={"ROI": "14%", "Users": 1200}, charts=["clicks.html"])

8. generate_from_template:
   - Purpose: Scans raw layout fields to resolve tag expressions using specified data maps.
   - Arguments:
     a) template_path: str - Local filesystem index highlighting source format structure files.
     b) data: Dict[str, Any] - Complete value mapping dict matching target document keys.
     c) output_format: str (default: "html") - Output file type identifier choice ("html", "pdf", "docx").
     d) output_path: str (default: "output_report") - Runtime structural name pointer identifying target deployment targets.
   - Returns: ToolResult noting validation profiles, fallback scenarios, or byte limits completed.
   - How to call: ReportGeneratorTool.generate_from_template(template_path="invoice.tmpl", data={"user": "John Doe", "total": "$50"}, output_format="html")
   """)
    
    @staticmethod
    def create_pdf_report(
        title: str,
        sections: List[Dict[str, str]],
        output: str = "report.pdf",
        logo: str = "",
        author: str = "NPM Agent",
        date: str = "",
        toc: bool = True,
    ) -> ToolResult:
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, HRFlowable
            from reportlab.platypus import TableOfContents, Table, TableStyle
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

            doc    = SimpleDocTemplate(output, pagesize=A4, rightMargin=inch, leftMargin=inch, topMargin=inch, bottomMargin=inch)
            styles = getSampleStyleSheet()
            story  = []

            # Title page
            story.append(Spacer(1, 2 * inch))
            story.append(Paragraph(title, ParagraphStyle("Title", fontSize=28, leading=36, alignment=TA_CENTER, textColor=colors.HexColor("#2C3E50"))))
            story.append(Spacer(1, 0.3 * inch))
            story.append(Paragraph(f"Author: {author}", ParagraphStyle("Author", fontSize=12, alignment=TA_CENTER, textColor=colors.grey)))
            story.append(Paragraph(f"Date: {date or datetime.now().strftime('%B %d, %Y')}", ParagraphStyle("Date", fontSize=12, alignment=TA_CENTER, textColor=colors.grey)))
            story.append(PageBreak())

            # TOC
            if toc:
                story.append(Paragraph("Table of Contents", styles["Heading1"]))
                story.append(Spacer(1, 0.2 * inch))
                for i, sec in enumerate(sections, 1):
                    story.append(Paragraph(f"{i}. {sec.get('title', '')}", styles["Normal"]))
                story.append(PageBreak())

            # Sections
            for sec in sections:
                story.append(Paragraph(sec.get("title", ""), styles["Heading1"]))
                story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#3498DB")))
                story.append(Spacer(1, 0.1 * inch))
                body = sec.get("content", "")
                for para in body.split("\n\n"):
                    if para.strip():
                        story.append(Paragraph(para.strip(), ParagraphStyle("Body", fontSize=10, leading=14, alignment=TA_JUSTIFY)))
                        story.append(Spacer(1, 0.1 * inch))
                # table data if provided
                if "table" in sec:
                    table_data = sec["table"]
                    if table_data:
                        t = Table(table_data)
                        t.setStyle(TableStyle([
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2C3E50")),
                            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
                            ("ALIGN",      (0, 0), (-1, -1), "CENTER"),
                            ("FONTSIZE",   (0, 0), (-1, -1), 9),
                            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#ECF0F1")]),
                            ("GRID",       (0, 0), (-1, -1), 0.5, colors.grey),
                        ]))
                        story.append(t)
                        story.append(Spacer(1, 0.2 * inch))
                story.append(PageBreak())

            doc.build(story)
            size = Path(output).stat().st_size
            return ToolResult(True, f"✓ PDF report '{title}' saved to '{output}' ({size} bytes)", {"output": output, "sections": len(sections)})
        except Exception as e:
            return ToolResult(False, f"✗ create_pdf_report failed: {e}")

    @staticmethod
    def create_word_report(
        title: str,
        sections: List[Dict[str, str]],
        output: str = "report.docx",
        logo: str = "",
        author: str = "NPM Agent",
        styles: Optional[Dict[str, Any]] = None,
    ) -> ToolResult:
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc  = Document()
            core = doc.core_properties
            core.author = author
            core.title  = title

            # Title
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            if logo and Path(logo).exists():
                doc.add_picture(logo, width=Inches(2))

            doc.add_paragraph(f"Generated by {author} on {datetime.now().strftime('%B %d, %Y')}")
            doc.add_page_break()

            # Sections
            for sec in sections:
                doc.add_heading(sec.get("title", "Untitled"), level=1)
                content = sec.get("content", "")
                for para in content.split("\n\n"):
                    if para.strip():
                        doc.add_paragraph(para.strip())
                if "table" in sec:
                    table_data = sec["table"]
                    if table_data:
                        t = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                        t.style = "Table Grid"
                        for i, row in enumerate(table_data):
                            for j, cell_val in enumerate(row):
                                t.cell(i, j).text = str(cell_val)

            doc.save(output)
            size = Path(output).stat().st_size
            return ToolResult(True, f"✓ Word report saved to '{output}' ({size} bytes)", {"output": output})
        except Exception as e:
            return ToolResult(False, f"✗ create_word_report failed: {e}")

    @staticmethod
    def create_excel_report(
        data_dict: Dict[str, Any],
        output: str = "report.xlsx",
        charts: bool = True,
        formatting: bool = True,
    ) -> ToolResult:
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.chart import BarChart, Reference
            from openpyxl.utils import get_column_letter
            import pandas as pd

            wb = openpyxl.Workbook()
            wb.remove(wb.active)

            header_fill   = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")
            header_font   = Font(color="FFFFFF", bold=True, size=11)
            alt_fill      = PatternFill(start_color="ECF0F1", end_color="ECF0F1", fill_type="solid")
            thin_border   = Border(
                left=Side(style="thin", color="BDBDBD"), right=Side(style="thin", color="BDBDBD"),
                top=Side(style="thin",  color="BDBDBD"), bottom=Side(style="thin", color="BDBDBD"),
            )

            for sheet_name, data in data_dict.items():
                ws   = wb.create_sheet(title=str(sheet_name)[:31])
                df   = pd.DataFrame(data) if not isinstance(data, pd.DataFrame) else data
                cols = list(df.columns)

                # Header row
                for j, col in enumerate(cols, 1):
                    cell = ws.cell(row=1, column=j, value=str(col))
                    if formatting:
                        cell.fill   = header_fill
                        cell.font   = header_font
                        cell.border = thin_border
                        cell.alignment = Alignment(horizontal="center")

                # Data rows
                for i, row in enumerate(df.itertuples(index=False), 2):
                    for j, val in enumerate(row, 1):
                        cell = ws.cell(row=i, column=j, value=val)
                        if formatting:
                            if i % 2 == 0: cell.fill = alt_fill
                            cell.border = thin_border

                # Auto-width
                for col_idx, col_name in enumerate(cols, 1):
                    max_len = max(len(str(col_name)), df[col_name].astype(str).str.len().max() if len(df) > 0 else 0)
                    ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 3, 40)

                # Add bar chart for numeric columns
                if charts and len(df) > 1:
                    num_cols = df.select_dtypes(include="number").columns.tolist()
                    if num_cols:
                        try:
                            chart    = BarChart()
                            chart.title = f"{sheet_name} Chart"
                            num_col_idx = cols.index(num_cols[0]) + 1
                            data_ref = Reference(ws, min_col=num_col_idx, min_row=1, max_row=min(len(df)+1, 20))
                            chart.add_data(data_ref, titles_from_data=True)
                            ws.add_chart(chart, f"A{len(df)+5}")
                        except Exception:
                            pass

            wb.save(output)
            size = Path(output).stat().st_size
            return ToolResult(True, f"✓ Excel report saved to '{output}' ({len(data_dict)} sheets, {size} bytes)", {"output": output, "sheets": list(data_dict.keys())})
        except Exception as e:
            return ToolResult(False, f"✗ create_excel_report failed: {e}")

    @staticmethod
    def create_presentation(
        slides_data: List[Dict[str, Any]],
        output: str = "presentation.pptx",
        theme: str = "dark",
    ) -> ToolResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            bg_color   = RGBColor(0x2C, 0x3E, 0x50) if theme == "dark" else RGBColor(0xFF, 0xFF, 0xFF)
            text_color = RGBColor(0xFF, 0xFF, 0xFF) if theme == "dark" else RGBColor(0x2C, 0x3E, 0x50)
            accent     = RGBColor(0x3A, 0x9B, 0xD5)

            for slide_data in slides_data:
                layout = prs.slide_layouts[6]  # blank
                slide  = prs.slides.add_slide(layout)

                # Background
                bg  = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

                # Title
                title_text = slide_data.get("title", "")
                if title_text:
                    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
                    tf     = tx_box.text_frame
                    tf.word_wrap = True
                    p = tf.add_paragraph()
                    p.text = title_text
                    p.font.size  = Pt(32)
                    p.font.bold  = True
                    p.font.color.rgb = accent
                    p.alignment  = PP_ALIGN.LEFT

                # Content
                content = slide_data.get("content", "")
                if content:
                    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5.5))
                    tf     = tx_box.text_frame
                    tf.word_wrap = True
                    for line in content.split("\n"):
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size  = Pt(16)
                        p.font.color.rgb = text_color
                        if line.startswith("•") or line.startswith("-"):
                            p.level = 1

                # Bullet points
                bullets = slide_data.get("bullets", [])
                if bullets:
                    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5.5))
                    tf     = tx_box.text_frame
                    tf.word_wrap = True
                    for bullet in bullets:
                        p = tf.add_paragraph()
                        p.text  = f"• {bullet}"
                        p.font.size  = Pt(18)
                        p.font.color.rgb = text_color
                        p.space_after = Pt(8)

            prs.save(output)
            size = Path(output).stat().st_size
            return ToolResult(True, f"✓ Presentation ({len(slides_data)} slides) saved to '{output}'", {"output": output, "slides": len(slides_data)})
        except Exception as e:
            return ToolResult(False, f"✗ create_presentation failed: {e}")

    @staticmethod
    def generate_research_report(
        topic: str,
        model: str = "llama3.2:3b",
        output: str = "research_report.pdf",
        include_citations: bool = True,
    ) -> ToolResult:
        try:
            from npmai import Ollama
            llm = Ollama(model=model, temperature=0.3, change=True, Models=["mistral:7b"])

            # Generate structured content
            outline_prompt = f"Create a detailed research report outline for the topic: '{topic}'. Return JSON: {{\"sections\": [{{\"title\": \"...\", \"key_points\": [\"...\"]}}]}} — 5-7 sections, no explanation."
            outline_raw    = llm.invoke(outline_prompt)
            clean = re.search(r'\{.*\}', outline_raw, re.DOTALL)
            outline = json.loads(clean.group()) if clean else {"sections": [{"title": "Overview", "key_points": [topic]}]}

            sections = []
            for sec in outline.get("sections", []):
                content_prompt = f"Write a detailed, professional section for a research report.\nSection: {sec['title']}\nKey points to cover: {sec.get('key_points', [])}\nTopic context: {topic}\nWrite 3-5 detailed paragraphs. Be informative and precise."
                content = llm.invoke(content_prompt)
                sections.append({"title": sec["title"], "content": content})

            if include_citations:
                cit_prompt = f"List 5 credible academic or authoritative sources (real or plausible) for the topic '{topic}'. Format as a numbered list."
                citations  = llm.invoke(cit_prompt)
                sections.append({"title": "References", "content": citations})

            result = ReportGeneratorTool.create_pdf_report(
                title    = f"Research Report: {topic}",
                sections = sections,
                output   = output,
                author   = "NPM Agent AI",
            )
            return ToolResult(result.success, f"✓ Research report on '{topic}' saved to '{output}'" if result.success else result.output, {"output": output, "sections": len(sections)})
        except ImportError:
            return ToolResult(False, "✗ npmai not available. Install npmai package.")
        except Exception as e:
            return ToolResult(False, f"✗ generate_research_report failed: {e}")

    @staticmethod
    def schedule_report(
        report_func: Callable,
        schedule: str,
        email_to: str = "",
        output_folder: str = "scheduled_reports",
    ) -> ToolResult:
        try:
            import schedule as sched
            Path(output_folder).mkdir(parents=True, exist_ok=True)

            def _job():
                ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
                outf = str(Path(output_folder) / f"report_{ts}.pdf")
                result = report_func(output=outf)
                if result.success and email_to:
                    # attempt email notification
                    try:
                        import smtplib
                        from email.mime.text import MIMEText
                        creds = CredStore.load("gmail")
                        msg   = MIMEText(f"Scheduled report generated: {outf}")
                        msg["Subject"] = f"NPM Agent: Scheduled Report {ts}"
                        msg["From"]    = creds.get("email", "")
                        msg["To"]      = email_to
                        with smtplib.SMTP("smtp.gmail.com", 587) as s:
                            s.starttls(); s.login(creds.get("email", ""), creds.get("password", ""))
                            s.sendmail(creds.get("email", ""), email_to, msg.as_string())
                    except Exception:
                        pass

            parts = schedule.lower().split()
            if "minutes" in parts:
                sched.every(int(parts[1])).minutes.do(_job)
            elif "hours" in parts:
                sched.every(int(parts[1])).hours.do(_job)
            elif "day" in parts and "at" in parts:
                t = parts[parts.index("at") + 1]
                sched.every().day.at(t).do(_job)
            elif "monday" in parts:
                sched.every().monday.at(parts[-1]).do(_job)

            def _run():
                while True:
                    sched.run_pending()
                    time.sleep(30)
            threading.Thread(target=_run, daemon=True).start()
            return ToolResult(True, f"✓ Report scheduled: '{schedule}'" + (f", emailing to {email_to}" if email_to else ""))
        except Exception as e:
            return ToolResult(False, f"✗ schedule_report failed: {e}")

    @staticmethod
    def create_dashboard_report(
        metrics: Dict[str, Any],
        charts: List[str],
        output_html: str = "dashboard_report.html",
        title: str = "Dashboard Report",
    ) -> ToolResult:
        try:
            # Build metric cards
            metric_cards = ""
            for k, v in metrics.items():
                color = "#3498DB"
                if isinstance(v, (int, float)) and v < 0: color = "#E74C3C"
                elif isinstance(v, (int, float)) and v > 0: color = "#2ECC71"
                metric_cards += f"""
                <div class="metric-card">
                    <div class="metric-label">{k}</div>
                    <div class="metric-value" style="color:{color}">{v}</div>
                </div>"""

            # Embed charts
            chart_divs = ""
            for chart_path in charts:
                if Path(chart_path).exists():
                    content = Path(chart_path).read_text()
                    body_m  = re.search(r"<body[^>]*>(.*?)</body>", content, re.DOTALL)
                    inner   = body_m.group(1) if body_m else content
                    chart_divs += f'<div class="chart-container">{inner}</div>'

            html = f"""<!DOCTYPE html>
<html>
<head>
  <title>{title}</title>
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0 }}
    body {{ font-family: 'Segoe UI', sans-serif; background: #1A1A2E; color: #E0E0E0; padding: 20px }}
    h1   {{ text-align: center; color: #3498DB; margin: 20px 0; font-size: 28px }}
    .subtitle {{ text-align:center; color:#888; margin-bottom:30px }}
    .metrics   {{ display: flex; flex-wrap: wrap; gap: 16px; justify-content: center; margin-bottom: 30px }}
    .metric-card {{ background: #16213E; border-radius: 12px; padding: 20px 30px; text-align: center; min-width: 160px; box-shadow: 0 4px 20px rgba(0,0,0,0.4) }}
    .metric-label {{ font-size: 12px; color: #888; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 8px }}
    .metric-value {{ font-size: 28px; font-weight: bold }}
    .charts    {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(560px, 1fr)); gap: 20px }}
    .chart-container {{ background: #16213E; border-radius: 12px; padding: 16px; overflow: hidden }}
  </style>
</head>
<body>
  <h1>📊 {title}</h1>
  <p class="subtitle">Generated {datetime.now().strftime('%B %d, %Y at %H:%M')}</p>
  <div class="metrics">{metric_cards}</div>
  <div class="charts">{chart_divs}</div>
</body>
</html>"""
            Path(output_html).write_text(html)
            size = Path(output_html).stat().st_size
            return ToolResult(True, f"✓ Dashboard report saved to '{output_html}' ({size} bytes)", {"output": output_html, "metrics": len(metrics), "charts": len(charts)})
        except Exception as e:
            return ToolResult(False, f"✗ create_dashboard_report failed: {e}")

    @staticmethod
    def generate_from_template(
        template_path: str,
        data: Dict[str, Any],
        output_format: str = "html",
        output_path: str = "output_report",
    ) -> ToolResult:
        try:
            template_str = Path(template_path).read_text()
            # Simple {{variable}} substitution
            rendered = template_str
            for k, v in data.items():
                rendered = rendered.replace(f"{{{{{k}}}}}", str(v))
                rendered = rendered.replace(f"{{{{ {k} }}}}", str(v))

            fmt = output_format.lower()
            if not output_path.endswith(f".{fmt}"):
                output_path = f"{output_path}.{fmt}"

            if fmt == "html":
                Path(output_path).write_text(rendered)
            elif fmt == "pdf":
                try:
                    import weasyprint
                    weasyprint.HTML(string=rendered).write_pdf(output_path)
                except ImportError:
                    # fallback: save HTML and note PDF conversion unavailable
                    html_path = output_path.replace(".pdf", ".html")
                    Path(html_path).write_text(rendered)
                    return ToolResult(True, f"✓ Template rendered as HTML (PDF requires weasyprint): '{html_path}'")
            elif fmt in ("docx", "word"):
                from docx import Document
                doc = Document()
                for line in rendered.split("\n"):
                    if line.strip():
                        doc.add_paragraph(line)
                doc.save(output_path)
            else:
                Path(output_path).write_text(rendered)

            size = Path(output_path).stat().st_size
            return ToolResult(True, f"✓ Report from template saved to '{output_path}' ({size} bytes)", {"output": output_path})
        except Exception as e:
            return ToolResult(False, f"✗ generate_from_template failed: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# Tool registry
# ─────────────────────────────────────────────────────────────────────────────

DATA_RESEARCH_TOOLS = {
    DataAnalysisTool.name:        DataAnalysisTool,
    VisualizationTool.name:       VisualizationTool,
    WebScrapingAdvancedTool.name: WebScrapingAdvancedTool,
    SearchResearchTool.name:      SearchResearchTool,
    FinancialDataTool.name:       FinancialDataTool,
    SocialMediaDataTool.name:     SocialMediaDataTool,
    WeatherGeoTool.name:          WeatherGeoTool,
    TextAnalyticsTool.name:       TextAnalyticsTool,
    DatabaseTool.name:            DatabaseTool,
    ReportGeneratorTool.name:     ReportGeneratorTool,
}

DATA_RESEARCH_TOOLS_SUMMARY = "\n".join(
    f"- {k}: {v.description}" for k, v in DATA_RESEARCH_TOOLS.items()
)

__all__ = [
    "DataAnalysisTool",
    "VisualizationTool",
    "WebScrapingAdvancedTool",
    "SearchResearchTool",
    "FinancialDataTool",
    "SocialMediaDataTool",
    "WeatherGeoTool",
    "TextAnalyticsTool",
    "DatabaseTool",
    "ReportGeneratorTool",
    "DATA_RESEARCH_TOOLS",
    "DATA_RESEARCH_TOOLS_SUMMARY",
]
